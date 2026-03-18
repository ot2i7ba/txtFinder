#!/usr/bin/env python3

# txtfinder.py
# Copyright (c) 2026 ot2i7ba
# https://github.com/ot2i7ba/
# This code is licensed under the MIT License (see LICENSE for details).

"""
txtFinder — Search filenames across document collections and highlight matches.

Interactive CLI tool with six menu options:

  1  Generate a filenames list (full, hash-only, or stem-only)
  2  Generate a hash list file (MD5 or SHA-256)
  3  Search within filenames.json (substring, keywords, hashes, filelist comparison)
  4  Search in files — PDF, JPG, DOCX/XLSX/PPTX/ODT, TXT
  5  Browse search history
  6  Edit runtime settings

Supported search modes: exact substring, fuzzy (Levenshtein-based), stem-only, regex.
Matches are highlighted directly in output copies of each searched file.
Per-file PDF reports and a CSV summary are generated after every search run.
"""

import os
# Suppress NumExpr thread-count warning before any library imports trigger it
os.environ.setdefault("NUMEXPR_MAX_THREADS", str(min(os.cpu_count() or 4, 8)))

import csv
import hashlib
import io
import json
import logging

# ---------------------------------------------------------------------------
# Optional fast JSON backend — transparent drop-in via module-level patching.
# orjson is 5–10× faster than stdlib json; falls back silently if not installed.
# orjson.JSONDecodeError is a subclass of json.JSONDecodeError, so all existing
# except-clauses remain fully compatible without modification.
# ---------------------------------------------------------------------------
try:
    import orjson as _orjson

    def _orjson_dumps(obj, *, indent=None, ensure_ascii=True, **_kw) -> str:
        """Wrap orjson.dumps to match stdlib json.dumps signature.

        Args:
            obj: Python object to serialise.
            indent: When truthy, apply 2-space indentation (orjson only supports
                indent=2 natively; other values are treated identically).
            ensure_ascii: Accepted for API compatibility; orjson always emits
                valid UTF-8 and ignores this parameter.
            **_kw: Any remaining stdlib kwargs are silently absorbed.

        Returns:
            JSON string with optional indentation.
        """
        option = _orjson.OPT_INDENT_2 if indent else None
        return _orjson.dumps(obj, option=option).decode()

    # Patch stdlib json module in-place — all downstream json.* calls
    # automatically use orjson without any further code changes.
    json.loads = _orjson.loads                                          # type: ignore[method-assign]
    json.dumps = _orjson_dumps                                          # type: ignore[method-assign]
    json.load  = lambda f, **_kw: _orjson.loads(f.read())              # type: ignore[method-assign]
    json.dump  = lambda obj, f, **kw: f.write(_orjson_dumps(obj, **kw))# type: ignore[method-assign]

except ImportError:
    pass  # stdlib json used transparently — no action needed
# ---------------------------------------------------------------------------
import atexit
import logging.handlers
import platform
import re
import shutil
import sys
import threading
import time
import unicodedata
from concurrent.futures import ThreadPoolExecutor, as_completed
from dataclasses import dataclass
from datetime import datetime
from functools import lru_cache, partial
from pathlib import Path
from typing import Any

# ---------------------------------------------------------------------------
# readline — persistent input history (up/down arrow) across sessions.
# stdlib on Linux/macOS; on Windows pyreadline3 is tried as a drop-in
# (pip install pyreadline3).  Silently disabled if neither is available.
# History file stored in the user's home directory so it persists across
# sessions regardless of the working directory txtFinder is launched from.
# Note: isatty() is intentionally NOT used as a guard — PowerShell and many
# other Windows terminal hosts report isatty()=False even in interactive use.
# ---------------------------------------------------------------------------
_READLINE_HISTORY = str(Path.home() / ".txtfinder_readline_history")
_READLINE_AVAILABLE = False
try:
    import readline as _readline
except ImportError:
    try:
        import pyreadline3 as _readline  # type: ignore[no-redef]  # Windows fallback
    except ImportError:
        _readline = None  # type: ignore[assignment]
if _readline is not None:
    try:
        if hasattr(_readline, "set_history_length"):
            _readline.set_history_length(200)
    except Exception:
        pass
    try:
        if hasattr(_readline, "read_history_file"):
            _readline.read_history_file(_READLINE_HISTORY)
    except Exception:
        pass  # file missing or unreadable — start with empty history
    try:
        if hasattr(_readline, "write_history_file"):
            atexit.register(_readline.write_history_file, _READLINE_HISTORY)
        _READLINE_AVAILABLE = True
    except Exception:
        pass

import fitz  # PyMuPDF
from rapidfuzz import fuzz as _rfuzz  # fast fuzzy matching
import typer

# Configure logging with rotation (max 1 MB per file, 3 backups kept)
_log_handler = logging.handlers.RotatingFileHandler(
    "txtfinder.log", maxBytes=1_048_576, backupCount=3, encoding="utf-8"
)
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[_log_handler, logging.StreamHandler()],
)
logger = logging.getLogger(__name__)

# Constants
APP_VERSION = "1.0"
FILENAMES_JSON = "filenames.json"
FILENAMES_META_JSON = "filenames_meta.json"  # Companion metadata (scan_root, hash_mode, …)
EXCLUDED_FILES = ['filenames.json', 'txtfinder.log', 'txtfinder.py', 'txtfinder.exe']
HASH_CHUNK_SIZE = 32768  # 32KB chunks for hash calculation

# ============================================================================
# MARKING CONFIGURATION - User Configurable
# ============================================================================
# Highlight color for PDF annotations (RGB: 0-1 float)
HIGHLIGHT_COLOR = (1, 0, 0)

# Text file markers
TXT_MARKER = "[[FOUND]]"
TXT_MARKER_POSITION = "before"  # "before" or "after"

# Output file suffixes
OUTPUT_SUFFIX = "_checked"
REPORT_SUFFIX = "_report"
HASHES_SUFFIX = "_hashes"
IMAGE_FOUND_SUFFIX = "_found"          # e.g. original.jpg → original_found.jpg
HASHLIST_REPORT_SUFFIX = "_hashlists"
FILELIST_REPORT_SUFFIX = "_filelist"
HITLIST_SUFFIX = "_hitlist"
HASHLISTS_DIR = "hashes"               # Directory scanned recursively for external hash list files
EXPORT_DIR = "export"                  # Root directory for hash match file exports

# Output directory (empty string = same directory as the source file)
OUTPUT_DIR = ""

# Report generation behaviour
REPORT_ON_MATCH_ONLY = True  # When True, skip report generation for files with 0 matches

# OCR Configuration
OCR_MIN_CONFIDENCE = 30  # Minimum confidence score for OCR word detection
OCR_LANGUAGES = "eng+deu"  # Tesseract language codes
OCR_PSM = 6    # Tesseract Page Segmentation Mode (0-13); 3=auto, 6=uniform block, 11=sparse text
OCR_DPI = 200  # Resolution for OCR page rendering

# Performance Configuration
MAX_SEARCH_PATTERNS = 500  # Maximum number of unique search patterns
MIN_STEM_LENGTH = 4  # Minimum stem length for pattern generation (avoids false positives)
STEMS_MIN_LENGTH = 5  # Minimum stem length for stems-only search (0 = no limit)
MAX_WORKERS = int(os.environ.get("TXTFINDER_WORKERS", min(os.cpu_count() or 4, 8)))

# Fuzzy Matching Configuration
FUZZY_THRESHOLD = 0.80  # Minimum similarity ratio for fuzzy matching (0.0-1.0)

# Search Profiles
PROFILES_DIR = "search_profiles"

# Directory Scanning
INPUT_DIR = "input"  # Default input directory name
EXCLUDED_DIRS = {"search_profiles", "__pycache__", ".git"}
HISTORY_FILE = "search_history.json"   # Search history log
HISTORY_MAX_ENTRIES = 100              # Maximum history entries to keep

# File type filter sets (used by "Generate filenames list" media/documents modes)
# Stored as frozensets for O(1) per-file extension lookup.
# Configurable at runtime via Settings menu (comma-separated string input).
MEDIA_EXTENSIONS: frozenset = frozenset({
    ".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".webm", ".m4v",
    ".mpg", ".mpeg", ".3gp", ".ts", ".mts", ".m2ts",
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".tif", ".webp",
    ".heic", ".heif", ".raw", ".cr2", ".nef", ".arw", ".dng", ".svg",
})
DOCUMENT_EXTENSIONS: frozenset = frozenset({
    ".pdf", ".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx",
    ".odt", ".ods", ".odp", ".odg", ".odf",
    ".txt", ".csv", ".rtf", ".md", ".html", ".htm", ".xml",
    ".eml", ".msg",
})
# ============================================================================


def _color_tuple(val) -> tuple:
    """Convert a value to an RGB color tuple for HIGHLIGHT_COLOR.

    Accepts a comma-separated string "R,G,B", a list [R,G,B], or a tuple (R,G,B).
    All components must be floats in the range 0.0–1.0.

    Args:
        val: Input value — string "R,G,B", list, or tuple.

    Returns:
        Tuple of three floats (r, g, b), each 0.0–1.0.

    Raises:
        ValueError: If the input cannot be parsed or values are out of range.
    """
    if isinstance(val, (list, tuple)):
        parts = [float(x) for x in val]
    else:
        parts = [float(x.strip()) for x in str(val).split(",")]
    if len(parts) != 3 or not all(0.0 <= p <= 1.0 for p in parts):
        raise ValueError("Expected R,G,B floats between 0.0 and 1.0")
    return tuple(parts)


def _parse_extensions(val) -> frozenset:
    """Convert a value to a frozenset of lowercase file extensions.

    Accepts a comma-separated string (e.g. ``".mp4,.mkv,.jpg"``), a list, a
    tuple, or an existing ``frozenset`` / ``set``.  Each token is stripped,
    lowercased, and must begin with a dot; tokens that do not start with a dot
    receive one automatically.

    Args:
        val: Input value — comma-separated string, list, tuple, set, or frozenset.

    Returns:
        A ``frozenset`` of normalised extension strings, e.g. ``frozenset({".mp4", ".mkv"})``.

    Raises:
        ValueError: If the input is empty after normalisation.
    """
    if isinstance(val, frozenset):
        return val
    if isinstance(val, (set, list, tuple)):
        tokens = [str(x).strip().lower() for x in val]
    else:
        tokens = [t.strip().lower() for t in str(val).split(",")]
    result = frozenset(t if t.startswith(".") else f".{t}" for t in tokens if t)
    if not result:
        raise ValueError("Extension list must contain at least one entry.")
    return result


def _parse_bool(val: object) -> bool:
    """Convert a value to bool, accepting common string representations.

    Args:
        val: Input value — ``bool``, or a string such as ``"true"``/``"false"``,
            ``"1"``/``"0"``, ``"yes"``/``"no"`` (case-insensitive).

    Returns:
        Parsed boolean value.

    Raises:
        ValueError: If the string cannot be interpreted as a boolean.
    """
    if isinstance(val, bool):
        return val
    s = str(val).strip().lower()
    if s in ("true", "1", "yes"):
        return True
    if s in ("false", "0", "no"):
        return False
    raise ValueError(f"Cannot parse {val!r} as bool — use true/false/1/0/yes/no")


@dataclass
class FileEntry:
    """A single entry from filenames.json."""
    filename: str   # bare filename only, e.g. "document.pdf"
    size: int
    sha256: str
    md5: str
    filepath: str = ""  # relative path from scan root, e.g. "subdir/document.pdf"

    @classmethod
    def from_dict(cls, d: dict) -> "FileEntry":
        """Construct a FileEntry from a raw dictionary (e.g. from JSON).

        Args:
            d: Dictionary with keys ``filename``, and optionally ``size``,
                ``sha256``, ``md5``, and ``filepath``.

        Returns:
            A populated FileEntry instance.
        """
        return cls(
            filename=d["filename"],
            size=d.get("size", 0),
            sha256=d.get("sha256", "N/A"),
            md5=d.get("md5", "N/A"),
            filepath=d.get("filepath", ""),
        )


@dataclass
class FileSearchResult:
    """Result of searching a single file."""
    file: Path
    matches: int
    output: Path | None
    report: Path | None
    error: str | None = None


def _supports_color() -> bool:
    """Check if the terminal supports ANSI color codes."""
    if os.environ.get("NO_COLOR"):
        return False
    if not hasattr(sys.stdout, "isatty") or not sys.stdout.isatty():
        return False
    if platform.system() == "Windows":
        return os.environ.get("WT_SESSION") or os.environ.get("ANSICON") or "xterm" in os.environ.get("TERM", "")
    return True


class _C:
    """ANSI color codes with graceful fallback."""
    _on = _supports_color()
    RESET = "\033[0m" if _on else ""
    BOLD = "\033[1m" if _on else ""
    DIM = "\033[2m" if _on else ""
    CYAN = "\033[36m" if _on else ""
    GREEN = "\033[32m" if _on else ""
    YELLOW = "\033[33m" if _on else ""
    ORANGE = "\033[38;5;208m" if _on else ""  # 256-color orange
    RED = "\033[31m" if _on else ""


# ============================================================================
# CONFIGURATION SECTION - Tesseract OCR Path Setup
# ============================================================================
# IMPORTANT FOR WINDOWS USERS:
# If Tesseract is installed but not found, uncomment and set the correct path:
# TESSERACT_PATH = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
#
# Common Windows locations:
# - C:\Program Files\Tesseract-OCR\tesseract.exe
# - C:\Program Files (x86)\Tesseract-OCR\tesseract.exe
# - C:\Users\YourName\AppData\Local\Programs\Tesseract-OCR\tesseract.exe
#
# Linux/Mac: Leave as None (auto-detected from PATH)
TESSERACT_PATH = None  # Set to explicit path if auto-detection fails
# ============================================================================

app = typer.Typer(add_completion=False)


# Seconds without a progress update before the spinner switches to a "may be
# slow / hung?" warning colour.  Applies only when update_message() is used
# (i.e. the caller is reporting page-level progress).
HANG_THRESHOLD = 30


def _fmt_elapsed(seconds: float) -> str:
    """Format elapsed seconds as M:SS or H:MM:SS.

    Args:
        seconds: Non-negative elapsed time in seconds.

    Returns:
        A human-readable duration string in ``M:SS`` or ``H:MM:SS`` format.
    """
    s = int(seconds)
    if s < 3600:
        return f"{s // 60}:{s % 60:02d}"
    return f"{s // 3600}:{(s % 3600) // 60:02d}:{s % 60:02d}"


class Spinner:
    """Animated spinner context manager for long-running operations.

    Displays a rotating frame and a live elapsed-time counter while work runs
    in the foreground thread.  If ``update_message()`` has not been called for
    ``HANG_THRESHOLD`` seconds the frame switches to yellow ``[?]`` to signal
    that progress has stalled.  In non-TTY environments a status line is
    printed every 15 seconds instead of animating in place.

    Usage::

        with Spinner("Loading data") as sp:
            do_work()
        # [✓] Loading data — 0:05

        with Spinner("Loading data") as sp:
            sp.warn("3 warnings")
        # [!] Loading data: 3 warnings — 0:05   (orange)

        with Spinner("Loading data") as sp:
            sp.fail("timeout")
        # [!] Loading data: timeout — 0:05       (red)
    """

    _FRAMES = ("[/]", "[-]", "[\\]", "[|]")
    _INTERVAL = 0.12
    _NON_TTY_INTERVAL = 15.0   # seconds between periodic status lines in non-TTY mode

    def __init__(self, message: str) -> None:
        """Initialise the spinner with a display message.

        Args:
            message: Status text shown next to the spinning frame.
        """
        self._lock = threading.Lock()
        self._message = message
        self._stop = threading.Event()
        self._failed: str | None = None
        self._warned: str | None = None
        self._thread: threading.Thread | None = None
        self._is_tty = hasattr(sys.stdout, "isatty") and sys.stdout.isatty()
        self._stream_handlers: list[logging.StreamHandler] = []
        self._start_time: float = 0.0
        self._last_activity: float = 0.0   # reset by update_message(); drives hang detection
        self._tracks_activity: bool = False # True once update_message() has been called

    def _animate(self) -> None:
        """Background thread loop that redraws the spinner frame each tick."""
        idx = 0
        last_non_tty_print = self._start_time
        while not self._stop.is_set():
            now = time.monotonic()
            elapsed = now - self._start_time
            with self._lock:
                msg = self._message
                last_act = self._last_activity
                tracks = self._tracks_activity

            if self._is_tty:
                # Hang detection: switch to yellow [?] when progress has stalled
                if tracks and (now - last_act) >= HANG_THRESHOLD:
                    idle = now - last_act
                    frame = f"[{_C.YELLOW}?{_C.RESET}]"
                    hint = f"  {_C.DIM}(no activity {_fmt_elapsed(idle)}){_C.RESET}"
                else:
                    frame = f"{_C.CYAN}{self._FRAMES[idx % len(self._FRAMES)]}{_C.RESET}"
                    hint = ""
                sys.stdout.write(
                    f"\r\033[K{frame} {msg}  {_C.DIM}({_fmt_elapsed(elapsed)}){_C.RESET}{hint}"
                )
                sys.stdout.flush()
            else:
                # Non-TTY: print a timestamped status line periodically
                if now - last_non_tty_print >= self._NON_TTY_INTERVAL:
                    sys.stdout.write(f"  [{_fmt_elapsed(elapsed)}] {msg}\n")
                    sys.stdout.flush()
                    last_non_tty_print = now

            idx += 1
            self._stop.wait(timeout=self._INTERVAL)

    def fail(self, reason: str = "") -> None:
        """Mark the spinner as failed with an optional reason.

        Args:
            reason: Short description of the failure shown in the exit line.
        """
        self._failed = reason

    def update_message(self, message: str) -> None:
        """Update the spinner message and reset the hang-detection timer.

        Args:
            message: New status text to display next to the spinning frame.
        """
        with self._lock:
            self._message = message
            self._last_activity = time.monotonic()
            self._tracks_activity = True

    def warn(self, reason: str = "") -> None:
        """Mark the spinner as completed with warnings.

        Args:
            reason: Short description of the warning shown in the exit line.
        """
        self._warned = reason

    def __enter__(self) -> "Spinner":
        """Start the spinner animation thread and suppress stream log handlers.

        Returns:
            This spinner instance, enabling ``with Spinner(...) as sp:`` syntax.
        """
        self._start_time = time.monotonic()
        self._last_activity = self._start_time
        if self._is_tty:
            # Suppress log StreamHandlers to prevent log lines from breaking the spinner
            for handler in logging.root.handlers:
                if isinstance(handler, logging.StreamHandler) and not isinstance(handler, logging.FileHandler):
                    self._stream_handlers.append(handler)
                    logging.root.removeHandler(handler)
            self._thread = threading.Thread(target=self._animate, daemon=True)
            self._thread.start()
        else:
            sys.stdout.write(f"[ ] {self._message}\n")
            sys.stdout.flush()
            # Start animation thread for non-TTY periodic updates as well
            self._thread = threading.Thread(target=self._animate, daemon=True)
            self._thread.start()
        return self

    def __exit__(self, exc_type, exc_val, exc_tb) -> None:
        """Stop the spinner thread, restore log handlers, and print the final status line.

        Args:
            exc_type: Exception type, or None if no exception occurred.
            exc_val: Exception value, or None.
            exc_tb: Exception traceback, or None.
        """
        self._stop.set()
        if self._thread is not None:
            self._thread.join(timeout=1.0)
        # Restore StreamHandlers
        for handler in self._stream_handlers:
            logging.root.addHandler(handler)
        self._stream_handlers.clear()
        elapsed_str = _fmt_elapsed(time.monotonic() - self._start_time)
        if not self._is_tty:
            with self._lock:
                msg = self._message
            if exc_type is not None:
                detail = f": {exc_val}" if exc_val else ""
                sys.stdout.write(f"[!] {msg}{detail} — {elapsed_str}\n")
            elif self._failed is not None:
                detail = f": {self._failed}" if self._failed else ""
                sys.stdout.write(f"[!] {msg}{detail} — {elapsed_str}\n")
            else:
                sys.stdout.write(f"[✓] {msg} — {elapsed_str}\n")
            sys.stdout.flush()
            return
        # TTY: \033[K = erase to end of line (clears leftover spinner frame)
        with self._lock:
            msg = self._message
        if exc_type is not None:
            detail = f": {exc_val}" if exc_val else ""
            sys.stdout.write(f"\r\033[K[{_C.RED}!{_C.RESET}] {msg}{detail} — {elapsed_str}\n")
        elif self._failed is not None:
            detail = f": {self._failed}" if self._failed else ""
            sys.stdout.write(f"\r\033[K[{_C.RED}!{_C.RESET}] {msg}{detail} — {elapsed_str}\n")
        elif self._warned is not None:
            detail = f": {self._warned}" if self._warned else ""
            sys.stdout.write(f"\r\033[K[{_C.ORANGE}!{_C.RESET}] {msg}{detail} — {elapsed_str}\n")
        else:
            sys.stdout.write(f"\r\033[K{_C.GREEN}[✓]{_C.RESET} {msg} — {elapsed_str}\n")
        sys.stdout.flush()


class ProgressBar:
    """Animated progress bar context manager for batch operations.

    Displays a block-fill bar with a live count, elapsed time, and estimated
    time remaining.  Exits with ``[!]`` on exception or interruption and
    ``[✓]`` on normal completion.  In non-TTY environments a status line is
    printed every 15 seconds.

    Usage::

        with ProgressBar(10, "PDF files") as pb:
            for f in files:
                process(f)
                pb.advance()
        # [✓] 10/10 PDF files — 0:08
    """

    _INTERVAL = 0.25
    _NON_TTY_INTERVAL = 15.0  # seconds between periodic status lines in non-TTY mode

    def __init__(self, total: int, message: str) -> None:
        """Initialise the progress bar.

        Args:
            total: Total number of items to process.
            message: Label displayed after the ``done/total`` counter.
        """
        self._total = total
        self._message = message
        self._completed = 0
        self._lock = threading.Lock()
        self._stop = threading.Event()
        self._thread: threading.Thread | None = None
        self._is_tty = hasattr(sys.stdout, "isatty") and sys.stdout.isatty()
        self._stream_handlers: list[logging.StreamHandler] = []
        self._interrupted = False
        self._start_time: float = 0.0

    def advance(self) -> None:
        """Increment completed count (thread-safe)."""
        with self._lock:
            self._completed += 1

    def mark_interrupted(self) -> None:
        """Mark the progress bar as interrupted."""
        self._interrupted = True

    def _render(self, done: int, now: float) -> str:
        """Build the progress bar string with elapsed time and ETA.

        Args:
            done: Number of items completed so far.
            now: Current monotonic timestamp (seconds).

        Returns:
            A formatted string ready for writing to stdout.
        """
        total = self._total
        pct = done / total if total > 0 else 0
        bar_width = 20
        filled = int(bar_width * pct)
        bar = "\u2588" * filled + "\u2591" * (bar_width - filled)
        elapsed = now - self._start_time
        elapsed_str = _fmt_elapsed(elapsed)
        # ETA is only meaningful once at least one item has completed
        if done > 0 and done < total:
            rate = elapsed / done          # seconds per item
            remaining = rate * (total - done)
            eta_str = f"  ~{_fmt_elapsed(remaining)} remaining"
        elif done >= total:
            eta_str = ""
        else:
            eta_str = ""
        return (
            f"\r\033[K[{bar}] {done}/{total} {self._message}"
            f"  {_C.DIM}{elapsed_str} elapsed{eta_str}{_C.RESET}"
        )

    def _animate(self) -> None:
        """Background thread loop that redraws the progress bar each tick."""
        last_non_tty_print = self._start_time
        while not self._stop.is_set():
            now = time.monotonic()
            with self._lock:
                done = self._completed
            if self._is_tty:
                sys.stdout.write(self._render(done, now))
                sys.stdout.flush()
            else:
                # Non-TTY: print a timestamped status line periodically
                if now - last_non_tty_print >= self._NON_TTY_INTERVAL:
                    elapsed_str = _fmt_elapsed(now - self._start_time)
                    sys.stdout.write(f"  [{elapsed_str}] {done}/{self._total} {self._message}\n")
                    sys.stdout.flush()
                    last_non_tty_print = now
            self._stop.wait(timeout=self._INTERVAL)

    def __enter__(self) -> "ProgressBar":
        """Start the progress animation thread and suppress stream log handlers.

        Returns:
            This progress bar instance, enabling ``with ProgressBar(...) as pb:`` syntax.
        """
        self._start_time = time.monotonic()
        if self._is_tty:
            for handler in logging.root.handlers:
                if isinstance(handler, logging.StreamHandler) and not isinstance(handler, logging.FileHandler):
                    self._stream_handlers.append(handler)
                    logging.root.removeHandler(handler)
        else:
            sys.stdout.write(f"[ ] {self._message} (0/{self._total})\n")
            sys.stdout.flush()
        # Animation thread runs for both TTY and non-TTY modes
        self._thread = threading.Thread(target=self._animate, daemon=True)
        self._thread.start()
        return self

    def __exit__(self, exc_type, exc_val, exc_tb) -> None:
        """Stop the progress thread, restore log handlers, and print the final status line.

        Args:
            exc_type: Exception type, or None if no exception occurred.
            exc_val: Exception value, or None.
            exc_tb: Exception traceback, or None.
        """
        self._stop.set()
        if self._thread is not None:
            self._thread.join(timeout=1.0)
        for handler in self._stream_handlers:
            logging.root.addHandler(handler)
        self._stream_handlers.clear()
        elapsed_str = _fmt_elapsed(time.monotonic() - self._start_time)
        with self._lock:
            done = self._completed

        if not self._is_tty:
            if exc_type is not None:
                detail = f": {exc_val}" if exc_val else ""
                sys.stdout.write(f"[!] {done}/{self._total} {self._message}{detail} — {elapsed_str}\n")
            elif self._interrupted:
                sys.stdout.write(f"[!] {done}/{self._total} {self._message} (interrupted) — {elapsed_str}\n")
            else:
                sys.stdout.write(f"[✓] {self._total}/{self._total} {self._message} — {elapsed_str}\n")
            sys.stdout.flush()
            return

        if exc_type is not None:
            detail = f": {exc_val}" if exc_val else ""
            sys.stdout.write(
                f"\r\033[K[{_C.RED}!{_C.RESET}] {done}/{self._total} {self._message}{detail}"
                f" — {elapsed_str}\n"
            )
        elif self._interrupted:
            sys.stdout.write(
                f"\r\033[K[{_C.ORANGE}!{_C.RESET}] {done}/{self._total} {self._message}"
                f" {_C.YELLOW}(interrupted){_C.RESET} — {elapsed_str}\n"
            )
        else:
            sys.stdout.write(
                f"\r\033[K{_C.GREEN}[\u2713]{_C.RESET} {self._total}/{self._total} {self._message}"
                f" — {elapsed_str}\n"
            )
        sys.stdout.flush()


def _output_path(source: Path, suffix: str, ext: str | None = None) -> Path:
    """Build an output file path respecting the OUTPUT_DIR setting.

    Args:
        source: The source file being processed
        suffix: Stem suffix to append (e.g. "_checked", "_found", "_report")
        ext: File extension including dot (default: same as source)

    Returns:
        Output Path in OUTPUT_DIR if set, otherwise next to source file
    """
    out_ext = ext if ext is not None else source.suffix
    name = f"{source.stem}{suffix}{out_ext}"
    if OUTPUT_DIR:
        out_dir = Path(OUTPUT_DIR)
        out_dir.mkdir(parents=True, exist_ok=True)
        return out_dir / name
    return source.parent / name


def clear_screen() -> None:
    """Clear the terminal screen in a cross-platform manner."""
    print("\033c", end="", flush=True)


def print_header() -> None:
    """Clear the terminal screen and print the persistent application header."""
    clear_screen()
    W = 44  # total number of ═ characters (visible inner width = W - 2)
    title    = f"txtFinder v{APP_VERSION} by ot2i7ba"
    subtitle = "Search \u00b7 Highlight \u00b7 Report"
    t_pad = W - 2 - len(title)
    s_pad = W - 2 - len(subtitle)
    print(f"{_C.CYAN}\u2554{chr(0x2550) * W}\u2557{_C.RESET}")
    print(f"{_C.CYAN}\u2551{_C.RESET} {_C.BOLD}{title}{_C.RESET}{' ' * t_pad} {_C.CYAN}\u2551{_C.RESET}")
    print(f"{_C.CYAN}\u2551{_C.RESET} {_C.DIM}{subtitle}{_C.RESET}{' ' * s_pad} {_C.CYAN}\u2551{_C.RESET}")
    print(f"{_C.CYAN}\u255a{chr(0x2550) * W}\u255d{_C.RESET}")
    cwd_name  = Path.cwd().name or "/"
    ocr_label = f"{_C.GREEN}ON{_C.RESET}" if TESSERACT_AVAILABLE else f"{_C.RED}OFF{_C.RESET}"
    # filenames.json status line
    _json_path = Path(FILENAMES_JSON)
    _meta_path = Path(FILENAMES_META_JSON)
    if _json_path.exists():
        _count_str = "?"
        _date_str = ""
        if _meta_path.exists():
            try:
                _meta = json.loads(_meta_path.read_text(encoding="utf-8"))
                _count_str = f"{_meta.get('file_count', '?'):,}"
                _date_str = str(_meta.get("generated_at", ""))[:10]
            except (OSError, json.JSONDecodeError, ValueError, KeyError):
                pass
        _json_info = f"JSON: {_count_str} entries" + (f" ({_date_str})" if _date_str else "")
        print(f"  Dir: {_C.BOLD}{cwd_name}{_C.RESET}   OCR: {ocr_label}   {_json_info}")
    else:
        print(f"  Dir: {_C.BOLD}{cwd_name}{_C.RESET}   OCR: {ocr_label}   {_C.DIM}JSON: not found{_C.RESET}")
    if not TESSERACT_AVAILABLE:
        _sys = platform.system()
        if _sys == "Linux":
            _hint = "sudo apt install tesseract-ocr tesseract-ocr-eng tesseract-ocr-deu"
        elif _sys == "Windows":
            _hint = "https://github.com/UB-Mannheim/tesseract/wiki"
        elif _sys == "Darwin":
            _hint = "brew install tesseract tesseract-lang"
        else:
            _hint = "Install Tesseract OCR for your platform"
        print(f"  {_C.YELLOW}OCR unavailable \u2014 install Tesseract:{_C.RESET}")
        print(f"  {_C.DIM}{_hint}{_C.RESET}")
    print()


def wait_for_enter(msg: str = "") -> None:
    """Prompt user to press Enter to continue.

    Args:
        msg: Custom prompt text; defaults to ``"\\nPress Enter to continue..."``.
    """
    input(msg or "\nPress Enter to continue...")


def ensure_filenames_loaded() -> list[str] | None:
    """Load filenames from JSON, offering inline generation if the file is missing.

    If ``filenames.json`` does not exist the user is asked whether to generate
    it now (scanning the current directory).  A ``y`` answer triggers
    :func:`generate_filenames_list` immediately so the search can continue
    without returning to the main menu.

    Returns:
        List of filenames, or None if loading failed or was cancelled.
    """
    json_path = Path(FILENAMES_JSON)
    if not json_path.exists():
        print(f"  {_C.YELLOW}{FILENAMES_JSON} not found.{_C.RESET}")
        answer = input("  Generate it now from the current directory? [y/n]: ").strip().lower()
        if answer != "y":
            return None
        try:
            with Spinner("Generating filenames list"):
                entries, _ = generate_filenames_list(scan_dir=Path("."))
            if not entries:
                print(f"  {_C.YELLOW}No files found — {FILENAMES_JSON} not created.{_C.RESET}")
                return None
            print(f"  {_C.GREEN}{FILENAMES_JSON} generated ({len(entries)} entries).{_C.RESET}")
        except (OSError, RuntimeError, ValueError) as e:
            print(f"  {_C.RED}Generation failed: {e}{_C.RESET}")
            return None
    try:
        return load_filenames_from_json(json_path)
    except (OSError, json.JSONDecodeError, ValueError) as e:
        print(f"  {_C.RED}Error loading {FILENAMES_JSON}: {e}{_C.RESET}")
        return None


def _ensure_json_data_loaded() -> "list[FileEntry] | None":
    """Load filenames.json as FileEntry objects, offering inline generation if missing.

    Mirrors the UX of :func:`ensure_filenames_loaded` but returns the full
    ``list[FileEntry]`` needed by handlers that require hash/path metadata (hash
    search, hashlist comparison, JSON search, hashlist generation).  If the file
    does not exist the user is asked whether to generate it now before loading.

    Returns:
        List of :class:`FileEntry` objects, or ``None`` if loading failed or
        was cancelled by the user.
    """
    json_path = Path(FILENAMES_JSON)
    if not json_path.exists():
        print(f"  {_C.YELLOW}{FILENAMES_JSON} not found.{_C.RESET}")
        answer = input("  Generate it now from the current directory? [y/n]: ").strip().lower()
        if answer != "y":
            return None
        try:
            with Spinner("Generating filenames list"):
                entries, _ = generate_filenames_list(scan_dir=Path("."))
            if not entries:
                print(f"  {_C.YELLOW}No files found — {FILENAMES_JSON} not created.{_C.RESET}")
                return None
            print(f"  {_C.GREEN}{FILENAMES_JSON} generated ({len(entries):,} entries).{_C.RESET}")
        except (OSError, RuntimeError, ValueError) as e:
            print(f"  {_C.RED}Generation failed: {e}{_C.RESET}")
            return None
    try:
        return load_json_data(json_path)
    except (OSError, json.JSONDecodeError, ValueError) as e:
        print(f"  {_C.RED}Error loading {FILENAMES_JSON}: {e}{_C.RESET}")
        return None


def print_batch_summary(file_type: str, processed: int, with_matches: int,
                        total_matches: int, mode: str = "exact") -> None:
    """Print a formatted summary box after batch processing and record to history.

    Args:
        file_type: Human-readable label for the file type (e.g. ``"PDF"``).
        processed: Total number of files examined.
        with_matches: Number of files that contained at least one match.
        total_matches: Cumulative match count across all processed files.
        mode: Search mode string recorded in the history entry.
    """
    w = 50
    print(f"\n{_C.CYAN}{'=' * w}{_C.RESET}")
    print(f"  {_C.BOLD}{file_type} Search Results{_C.RESET}")
    print(f"{_C.CYAN}{'=' * w}{_C.RESET}")
    print(f"  Files processed:    {processed}")
    print(f"  Files with matches: {with_matches}")
    print(f"  Total matches:      {total_matches}")
    print(f"{_C.CYAN}{'=' * w}{_C.RESET}\n")
    _append_history(file_type, mode, processed, with_matches, total_matches)


def _ask_search_mode() -> str:
    """Prompt for search mode selection.

    Returns:
        "exact", "fuzzy", "regex", "stems", "back", or raises SystemExit on q
    """
    profiles = list_search_profiles()
    print()
    print(f"  {_C.BOLD}1{_C.RESET}  Exact matching (fast, with extension variants)")
    print(f"  {_C.BOLD}2{_C.RESET}  Fuzzy matching (slower, catches typos/OCR errors)")
    print(f"  {_C.BOLD}3{_C.RESET}  Stems only (fastest, bare name without extension)")
    print(f"  {_C.BOLD}4{_C.RESET}  Regex matching (custom pattern)")
    if profiles:
        print(f"  {_C.BOLD}p{_C.RESET}  Load search profile ({', '.join(profiles)})")
    print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}")
    while True:
        choice = input("  Search mode [1/2/3/4/b/q]: ").strip().lower()
        if choice == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        if choice == "b":
            return "back"
        if choice == "1":
            return "exact"
        if choice == "2":
            return "fuzzy"
        if choice == "3":
            return "stems"
        if choice == "4":
            return "regex"
        if choice == "p" and profiles:
            return "profile"
        print(f"  {_C.RED}Invalid input.{_C.RESET}")


def _ask_stems_min_length() -> int:
    """Prompt for the minimum stem length used in stems-only search.

    Returns:
        Minimum length as int; 0 means no restriction.
    """
    default = STEMS_MIN_LENGTH
    hint = "0 = no limit" if default == 0 else f"0 = no limit, Enter = {default}"
    raw = input(f"  Min. stem length [{hint}]: ").strip()
    if not raw:
        return default
    try:
        val = int(raw)
        if val < 0:
            val = 0
        return val
    except ValueError:
        return default


def _ask_regex_patterns() -> list[re.Pattern]:
    """Prompt for regex patterns (comma-separated).

    Returns:
        List of compiled regex patterns (may be empty).
    """
    print()
    print(f"  {_C.DIM}Enter regex pattern(s), comma-separated{_C.RESET}")
    print(f"  {_C.DIM}Example: report.*2024, IMG_\\d+\\.jpg{_C.RESET}")
    raw = input("  Regex pattern(s): ").strip()
    if not raw:
        return []
    patterns = []
    for p in raw.split(","):
        p = p.strip()
        if p:
            try:
                patterns.append(re.compile(p, re.IGNORECASE))
            except re.error as e:
                print(f"  {_C.RED}Invalid regex '{p}': {e}{_C.RESET}")
    return patterns


def regex_text_search(text: str, regex_patterns: list[re.Pattern]) -> dict[str, int]:
    """Search text with regex patterns.

    Args:
        text: Text to search in
        regex_patterns: List of compiled regex patterns

    Returns:
        Dict mapping pattern string to match count
    """
    results = {}
    for rx in regex_patterns:
        matches = rx.findall(text)
        results[rx.pattern] = len(matches)
    return results


def _search_file_with_regex(
    file_path: Path, _filenames: list[str],
    regex_patterns: list[re.Pattern] | None = None,
    pattern_set: "set[str] | None" = None,
) -> tuple[int, Path | None]:
    """Search a file using regex patterns on extracted text.

    Args:
        file_path: Path to the file to search
        _filenames: Unused (kept for signature compatibility with process_batch)
        regex_patterns: Compiled regex patterns to search for
        pattern_set: Unused (kept for signature compatibility with process_batch)

    Returns:
        Tuple of (total_matches, output_path or None)
    """
    if not regex_patterns:
        return 0, None

    text = extract_searchable_text(file_path)
    _cache_extracted_text(file_path, text)

    results = regex_text_search(text, regex_patterns)
    total = sum(results.values())

    if total == 0:
        return 0, None

    ext = file_path.suffix.lower()

    # For PDFs, highlight regex matches
    if ext == '.pdf':
        try:
            fitz.TOOLS.mupdf_display_errors(False)
            doc = fitz.open(file_path)
            if doc.is_encrypted:
                doc.close()
                return total, None
            for page in doc:
                page_text = page.get_text("text")
                for rx in regex_patterns:
                    for m in rx.finditer(page_text):
                        match_text = m.group()
                        if match_text.strip():
                            instances = page.search_for(match_text)
                            for bbox in instances:
                                highlight = page.add_highlight_annot(bbox)
                                highlight.set_colors(stroke=(1, 0, 0))
                                highlight.update()
            output_path = _output_path(file_path, OUTPUT_SUFFIX)
            doc.save(output_path)
            doc.close()
            fitz.TOOLS.mupdf_display_errors(True)
            return total, output_path
        except (OSError, RuntimeError, ValueError) as e:
            logger.error(f"Regex PDF highlight failed: {e}")
            fitz.TOOLS.mupdf_display_errors(True)
            return total, None

    # For TXT, add markers at regex match positions
    if ext == '.txt':
        try:
            content = file_path.read_text(encoding='utf-8', errors='ignore')
            modified = content
            for rx in regex_patterns:
                matches = list(rx.finditer(modified))
                for match in reversed(matches):
                    s, e = match.span()
                    matched_text = modified[s:e]
                    modified = modified[:s] + f"{TXT_MARKER}{matched_text}" + modified[e:]
            output_path = _output_path(file_path, OUTPUT_SUFFIX)
            output_path.write_text(modified, encoding='utf-8')
            return total, output_path
        except (OSError, UnicodeDecodeError) as e:
            logger.error(f"Regex TXT marking failed: {e}")
            return total, None

    # For other types, just copy if matches found
    suffix = OUTPUT_SUFFIX if ext != '.jpg' and ext != '.jpeg' else IMAGE_FOUND_SUFFIX
    output_path = file_path.parent / f"{file_path.stem}{suffix}{file_path.suffix}"
    shutil.copy2(file_path, output_path)
    return total, output_path


def _ask_extra_search_words() -> list[str]:
    """Prompt user for optional additional search words (comma-separated).

    Input format: "Word1, Word 2, Word3, This is a keyword"
    - Commas separate individual search terms
    - Leading/trailing whitespace per term is stripped
    - Empty terms are ignored
    - Each term is used as-is (no pattern generation)

    Returns:
        List of extra search terms (may be empty).
    """
    print()
    print(f"  {_C.DIM}Optional: Enter additional search words (comma-separated){_C.RESET}")
    print(f"  {_C.DIM}Example: Word1, Word 2, Another keyword{_C.RESET}")
    raw = input("  Extra words [none]: ").strip()
    if not raw:
        return []
    words = [w.strip() for w in raw.split(",") if w.strip()]
    if words:
        logger.info(f"Manual search words added: {words}")
    return words


class _SearchConfig:
    """Immutable search configuration passed to file-type handlers.

    Attributes:
        mode: Search mode string ("exact", "fuzzy", "regex", "stems").
        regex_patterns: Compiled regex patterns (non-empty only for "regex" mode).
        extra_words: User-supplied additional search terms.
        stems_min_len: Minimum stem length (0 = no restriction, "stems" mode only).
    """

    __slots__ = ("mode", "regex_patterns", "extra_words", "stems_min_len")

    def __init__(
        self,
        mode: str,
        regex_patterns: list[re.Pattern],
        extra_words: list[str],
        stems_min_len: int,
    ) -> None:
        """Initialise a _SearchConfig.

        Args:
            mode: Search mode ("exact", "fuzzy", "regex", "stems").
            regex_patterns: Pre-compiled regex patterns.
            extra_words: Additional user-supplied search terms.
            stems_min_len: Minimum stem length (0 = no restriction).
        """
        self.mode = mode
        self.regex_patterns = regex_patterns
        self.extra_words = extra_words
        self.stems_min_len = stems_min_len


def _ask_search_config() -> "_SearchConfig | None":
    """Prompt for the full search configuration (mode, patterns, extra words).

    Asks for search mode, additional inputs depending on the mode, and offers
    to save the configuration as a named profile.  Used by the "All of the
    above" branch in :func:`handle_search_files` so the four file-type
    handlers share a single configuration dialog.

    Returns:
        A :class:`_SearchConfig` instance, or ``None`` if the user cancels
        (back / quit).
    """
    mode = _ask_search_mode()
    if mode == "back":
        return None
    if mode == "stems":
        regex_patterns: list[re.Pattern] = []
        extra_words: list[str] = []
        stems_min_len = _ask_stems_min_length()
    elif mode == "profile":
        loaded = _load_profile_interactive()
        if loaded:
            mode, regex_patterns, extra_words = loaded
        else:
            mode, regex_patterns, extra_words = "exact", [], []
        stems_min_len = 0
    else:
        regex_patterns = _ask_regex_patterns() if mode == "regex" else []
        extra_words = _ask_extra_search_words()
        stems_min_len = 0

    # Offer to save profile (skipped for stems mode)
    if mode != "stems":
        save_name = input("  Save as profile? (name or Enter to skip): ").strip()
        if save_name:
            save_search_profile(save_name, mode, extra_words,
                                [rx.pattern for rx in regex_patterns])
            print(f"  {_C.GREEN}Profile '{save_name}' saved{_C.RESET}")

    return _SearchConfig(mode, regex_patterns, extra_words, stems_min_len)


def save_search_profile(name: str, mode: str, extra_words: list[str],
                        regex_patterns: list[str]) -> Path:
    """Save search configuration to JSON profile.

    Args:
        name: Profile name (used as filename)
        mode: "exact", "fuzzy", or "regex"
        extra_words: List of manual search words
        regex_patterns: List of regex pattern strings

    Returns:
        Path to the saved profile file
    """
    Path(PROFILES_DIR).mkdir(exist_ok=True)
    profile = {
        "name": name,
        "mode": mode,
        "extra_words": extra_words,
        "regex_patterns": regex_patterns,
        "created": datetime.now().isoformat(),
    }
    path = Path(PROFILES_DIR) / f"{name}.json"
    path.write_text(json.dumps(profile, indent=2, ensure_ascii=False), encoding="utf-8")
    logger.info(f"Profile saved: {path}")
    return path


def load_search_profile(name: str) -> dict:
    """Load a search profile from JSON.

    Args:
        name: Profile name (without .json extension)

    Returns:
        Dict with profile data (mode, extra_words, regex_patterns)
    """
    path = Path(PROFILES_DIR) / f"{name}.json"
    return json.loads(path.read_text(encoding="utf-8"))


def list_search_profiles() -> list[str]:
    """List available profile names.

    Returns:
        Sorted list of profile names (without .json extension)
    """
    d = Path(PROFILES_DIR)
    if not d.exists():
        return []
    return [f.stem for f in sorted(d.glob("*.json"))]


def _load_profile_interactive() -> tuple[str, list[re.Pattern], list[str]] | None:
    """Interactive profile loading.

    Returns:
        Tuple of (mode, regex_patterns, extra_words) or None if cancelled
    """
    profiles = list_search_profiles()
    if not profiles:
        print(f"  {_C.DIM}No profiles found{_C.RESET}")
        return None

    print(f"\n  Available profiles:")
    for i, name in enumerate(profiles, 1):
        print(f"    {i}. {name}")
    print(f"    {_C.DIM}q. Cancel{_C.RESET}")

    choice = input("  Select profile: ").strip()
    if choice.lower() == "q":
        return None

    try:
        idx = int(choice) - 1
        if 0 <= idx < len(profiles):
            profile = load_search_profile(profiles[idx])
            mode = profile.get("mode", "exact")
            extra_words = profile.get("extra_words", [])
            regex_strs = profile.get("regex_patterns", [])
            regex_patterns = []
            for p in regex_strs:
                try:
                    regex_patterns.append(re.compile(p, re.IGNORECASE))
                except re.error:
                    logger.warning(f"Skipping invalid regex from profile: {p}")
            print(f"  {_C.GREEN}Loaded profile '{profiles[idx]}' (mode: {mode}){_C.RESET}")
            return mode, regex_patterns, extra_words
    except ValueError:
        pass

    print(f"  {_C.RED}Invalid selection{_C.RESET}")
    return None


def validate_safe_path(file_path: Path, base_dir: Path) -> bool:
    """
    Validate that a path is within the base directory (prevents path traversal).

    Args:
        file_path: Path to validate
        base_dir: Base directory to check against

    Returns:
        True if path is safe, False otherwise
    """
    try:
        resolved_path = file_path.resolve()
        resolved_base = base_dir.resolve()
        resolved_path.relative_to(resolved_base)
        return True
    except ValueError:
        return False
    except OSError as e:
        logger.error(f"Path validation error: {e}")
        return False


def should_exclude_file(filename: str) -> bool:
    """
    Check if file should be excluded from scanning (case-insensitive).

    Args:
        filename: Name of the file to check

    Returns:
        True if file should be excluded, False otherwise
    """
    name_lower = Path(filename).name.lower()
    return any(name_lower == excl.lower() for excl in EXCLUDED_FILES)


def calculate_hashes(file_path: Path, hash_mode: str = "both") -> dict[str, str]:
    """Calculate cryptographic hashes of a file in a single pass (32 KB chunks).

    Only the hashes requested by ``hash_mode`` are computed, so callers that
    need just SHA-256 or just MD5 avoid the overhead of the unused algorithm.

    Args:
        file_path: Path to the file to hash.
        hash_mode: Which hashes to compute.  ``"both"`` computes SHA-256 and
            MD5 (default), ``"sha256"`` computes only SHA-256, ``"md5"``
            computes only MD5.

    Returns:
        Dict with keys present for each computed algorithm:
        ``"sha256"`` and/or ``"md5"`` containing lowercase hex digest strings.

    Raises:
        IOError: If file cannot be read.
    """
    do_sha256 = hash_mode in ("both", "sha256")
    do_md5    = hash_mode in ("both", "md5")
    sha256_h  = hashlib.sha256() if do_sha256 else None
    md5_h     = hashlib.md5()    if do_md5    else None
    try:
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(HASH_CHUNK_SIZE), b""):
                if sha256_h:
                    sha256_h.update(chunk)
                if md5_h:
                    md5_h.update(chunk)
        result: dict[str, str] = {}
        if sha256_h:
            result["sha256"] = sha256_h.hexdigest()
        if md5_h:
            result["md5"] = md5_h.hexdigest()
        return result
    except OSError as e:
        logger.error(f"Failed to calculate hashes for {file_path}: {e}")
        raise


def _count_files_in_dir(scan_dir: Path, extensions: frozenset | None = None) -> int:
    """Quick file count without hashing (for preview), using scandir for speed.

    Args:
        scan_dir: Root directory to count files in (recursive).
        extensions: Optional frozenset of lowercase extensions to count (e.g.
            ``frozenset({".mp4", ".jpg"})``).  When ``None`` all files are counted.

    Returns:
        Total number of regular files found under ``scan_dir``, optionally
        filtered by extension.
    """
    count = 0
    stack = [scan_dir]
    while stack:
        current = stack.pop()
        try:
            with os.scandir(current) as it:
                for entry in it:
                    if entry.is_file(follow_symlinks=False):
                        if extensions is None or Path(entry.name).suffix.lower() in extensions:
                            count += 1
                    elif entry.is_dir(follow_symlinks=False):
                        stack.append(entry.path)
        except PermissionError:
            pass
    return count


def _validate_scan_directory(dir_path: Path) -> tuple[bool, str]:
    """Validate that a directory is suitable for scanning.

    Args:
        dir_path: Filesystem path to validate.

    Returns:
        A ``(ok, error_message)`` tuple where ``ok`` is ``True`` when the
        directory exists, is accessible, and contains at least one entry.
    """
    if not dir_path.exists():
        return False, f"Directory does not exist: {dir_path}"
    if not dir_path.is_dir():
        return False, f"Path is not a directory: {dir_path}"
    try:
        next(dir_path.iterdir())
    except StopIteration:
        return False, f"Directory is empty: {dir_path}"
    except PermissionError:
        return False, f"Permission denied: {dir_path}"
    return True, ""


def generate_filenames_list(
    scan_dir: Path = Path("."),
    output_dir: Path | None = None,
    progress_callback=None,
    exclude_dirs: bool = False,
    stems_only: bool = False,
    no_hashes: bool = False,
    deduplicate: bool = False,
    extensions: frozenset | None = None,
    hash_mode: str = "both",
) -> tuple[int, int]:
    """Recursively scan a directory and generate a filenames list with metadata.

    Excludes system files: filenames.json, txtfinder.log, txtfinder.py,
    txtfinder.exe.  When ``extensions`` is provided only files whose suffix
    (case-insensitive) is in the set are included — all others are skipped.
    The lookup is O(1) per file regardless of set size.

    Args:
        scan_dir: Directory to scan (default: current directory).
        output_dir: Where to write the JSON (default: CWD).
        progress_callback: Called with (file_count) after each file for live
            updates.
        exclude_dirs: When ``True``, skip ``EXCLUDED_DIRS`` during walk (for
            CWD scans).
        stems_only: When ``True``, store only bare filename stems (no hashing).
        no_hashes: When ``True``, store full filename + filepath + size but
            skip hash calculation entirely.
        deduplicate: When ``True``, keep only the first occurrence of each
            filename (case-insensitive).
        extensions: Optional frozenset of lowercase extensions to include
            (e.g. ``frozenset({".mp4", ".jpg"})``).  ``None`` = include all.
        hash_mode: Which hashes to compute when ``stems_only`` and
            ``no_hashes`` are both ``False``.  ``"both"`` computes SHA-256
            and MD5 (default), ``"sha256"`` computes only SHA-256, ``"md5"``
            computes only MD5.  Ignored when ``stems_only`` or ``no_hashes``
            is ``True``.

    Returns:
        Tuple of (entries_written, duplicates_removed).

    Raises:
        IOError: If the JSON file cannot be written.
    """
    files_data: list[dict[str, Any]] = []
    file_count = 0
    excluded_count = 0

    logger.info(f"Starting file scan in directory: {scan_dir}")

    try:
        for root, dirs, files in os.walk(scan_dir):
            # Skip tool-artifact directories when scanning CWD
            if exclude_dirs:
                dirs[:] = [d for d in dirs if d not in EXCLUDED_DIRS]

            root_path = Path(root)
            for file in files:
                file_path = root_path / file

                # Check if file should be excluded (case-insensitive)
                if should_exclude_file(file):
                    excluded_count += 1
                    logger.info(f"Excluded file: {file}")
                    continue

                # Apply optional extension filter (O(1) frozenset lookup)
                if extensions is not None:
                    if Path(file).suffix.lower() not in extensions:
                        excluded_count += 1
                        continue

                # Skip tool output files when scanning CWD
                if exclude_dirs:
                    stem_lower = Path(file).stem.lower()
                    if (OUTPUT_SUFFIX in stem_lower
                            or REPORT_SUFFIX in stem_lower
                            or HASHES_SUFFIX in stem_lower
                            or IMAGE_FOUND_SUFFIX in stem_lower):
                        excluded_count += 1
                        continue

                try:
                    if stems_only:
                        # Stems-only mode: store only the bare filename stem.
                        # No hashing, no size calculation — maximum speed.
                        files_data.append({"filename": Path(file).stem})
                    elif no_hashes:
                        # Full filename mode without hashing — fast, no sha256/md5.
                        relative_path = file_path.relative_to(scan_dir)
                        file_size = file_path.stat().st_size
                        files_data.append(
                            {
                                "filename": file,
                                "filepath": relative_path.as_posix(),
                                "size": file_size,
                            }
                        )
                    else:
                        relative_path = file_path.relative_to(scan_dir)
                        file_size = file_path.stat().st_size
                        hashes = calculate_hashes(file_path, hash_mode=hash_mode)
                        entry: dict[str, Any] = {
                            "filename": file,
                            "filepath": relative_path.as_posix(),
                            "size": file_size,
                        }
                        if "sha256" in hashes:
                            entry["sha256"] = hashes["sha256"]
                        if "md5" in hashes:
                            entry["md5"] = hashes["md5"]
                        files_data.append(entry)
                    file_count += 1

                    if progress_callback:
                        progress_callback(file_count)
                    elif file_count % 100 == 0:
                        logger.info(f"Processed {file_count} files...")

                except (OSError, UnicodeDecodeError, ValueError) as e:
                    logger.error(f"Failed to process {file_path}: {e}")
                    continue

        # Deduplicate: keep only the first occurrence of each filename (case-insensitive).
        # For stems-only mode the key is the stem; for full mode it is the filename.
        duplicates_removed = 0
        if deduplicate and files_data:
            seen_keys: set[str] = set()
            unique_data: list[dict[str, Any]] = []
            for entry in files_data:
                key = entry["filename"].lower()
                if key not in seen_keys:
                    seen_keys.add(key)
                    unique_data.append(entry)
                else:
                    duplicates_removed += 1
            files_data = unique_data
            if duplicates_removed:
                logger.info(f"Deduplication removed {duplicates_removed} duplicate filename(s)")

        # Always write JSON into output_dir (defaults to CWD)
        out_dir = output_dir or Path(".")
        json_path = out_dir / FILENAMES_JSON
        with open(json_path, "w", encoding="utf-8") as f:
            json.dump(files_data, f, indent=2, ensure_ascii=False)

        entries_written = file_count - duplicates_removed

        # Write companion metadata so the export feature can suggest scan_root
        effective_hash_mode = "none" if (stems_only or no_hashes) else hash_mode
        meta = {
            "scan_root": str(scan_dir.resolve()),
            "hash_mode": effective_hash_mode,
            "generated_at": datetime.now().isoformat(timespec="seconds"),
            "file_count": entries_written,
        }
        try:
            meta_path = out_dir / FILENAMES_META_JSON
            with open(meta_path, "w", encoding="utf-8") as f:
                json.dump(meta, f, indent=2, ensure_ascii=False)
        except OSError as e:
            logger.warning(f"Could not write {FILENAMES_META_JSON}: {e}")
        logger.info(
            f"File scan complete: {file_count} files processed, "
            f"{excluded_count} excluded, {duplicates_removed} duplicates removed"
        )
        return entries_written, duplicates_removed

    except (OSError, ValueError) as e:
        logger.error(f"Failed to generate filenames list: {e}")
        raise


def load_json_data(json_path: Path = Path(FILENAMES_JSON)) -> list[FileEntry]:
    """
    Load and validate JSON data, returning FileEntry objects.

    Args:
        json_path: Path to the JSON file

    Returns:
        List of validated FileEntry objects

    Raises:
        FileNotFoundError, json.JSONDecodeError, ValueError
    """
    try:
        with open(json_path, "r", encoding="utf-8") as f:
            data = json.load(f)

        if not isinstance(data, list):
            raise ValueError(f"JSON must contain a list, got {type(data).__name__}")

        entries: list[FileEntry] = []
        base_dir = json_path.parent
        rejected_count = 0

        for idx, entry in enumerate(data):
            if not isinstance(entry, dict):
                logger.warning(f"Skipping entry {idx}: not a dictionary")
                continue

            if 'filename' not in entry:
                logger.warning(f"Skipping entry {idx}: missing 'filename' key")
                continue

            filename = entry["filename"]

            # filename must be a bare name — reject if it contains path separators
            if Path(filename).is_absolute() or "/" in filename or "\\" in filename:
                logger.warning(f"Skipping entry {idx}: filename must not contain path separators: {filename}")
                rejected_count += 1
                continue

            # validate filepath (relative path from scan root) if present
            filepath = entry.get("filepath", "")
            if filepath:
                if Path(filepath).is_absolute():
                    logger.warning(f"Skipping entry {idx}: absolute filepath not allowed: {filepath}")
                    rejected_count += 1
                    continue
                full_path = base_dir / filepath
                if not validate_safe_path(full_path, base_dir):
                    logger.warning(f"Skipping entry {idx}: filepath outside working directory: {filepath}")
                    rejected_count += 1
                    continue

            entries.append(FileEntry.from_dict(entry))

        if not entries:
            raise ValueError("No valid filenames found in JSON")

        if rejected_count > 0:
            logger.warning(f"Rejected {rejected_count} unsafe paths from JSON")

        logger.info(f"Loaded {len(entries)} valid entries from JSON")
        return entries

    except FileNotFoundError:
        logger.error(f"File {json_path} not found")
        raise
    except json.JSONDecodeError as e:
        logger.error(f"Invalid JSON in {json_path}: {e}")
        raise
    except ValueError as e:
        logger.error(f"Invalid JSON structure: {e}")
        raise
    except OSError as e:
        logger.error(f"Failed to load JSON data: {e}")
        raise


def load_filenames_from_json(json_path: Path = Path(FILENAMES_JSON)) -> list[str]:
    """
    Load filenames from JSON file with validation.

    Args:
        json_path: Path to the JSON file

    Returns:
        List of filenames from the JSON file
    """
    entries = load_json_data(json_path)
    return [entry.filename for entry in entries]


def _is_output_file(stem_lower: str) -> bool:
    """Return True if the filename stem matches any tool-generated output suffix.

    Checks OUTPUT_SUFFIX, REPORT_SUFFIX, HASHES_SUFFIX, IMAGE_FOUND_SUFFIX,
    HASHLIST_REPORT_SUFFIX, FILELIST_REPORT_SUFFIX, and HITLIST_SUFFIX so that all get_*_files functions and
    generate_filenames_list exclude tool-generated files from search and list results.

    Args:
        stem_lower: Lowercased filename stem (without extension).

    Returns:
        True if the stem contains any defined output suffix.
    """
    return (
        OUTPUT_SUFFIX in stem_lower
        or REPORT_SUFFIX in stem_lower
        or HASHES_SUFFIX in stem_lower
        or IMAGE_FOUND_SUFFIX in stem_lower
        or HASHLIST_REPORT_SUFFIX in stem_lower
        or FILELIST_REPORT_SUFFIX in stem_lower
        or HITLIST_SUFFIX in stem_lower
    )


def get_pdf_files(directory: Path = Path(".")) -> list[Path]:
    """
    Get all PDF files in the specified directory (non-recursive).
    Excludes PDFs whose filename stem matches any tool-generated output suffix.

    Args:
        directory: Directory to search for PDF files

    Returns:
        List of Path objects for PDF files (sorted, excluding _checked PDFs)
    """
    try:
        pdf_files = [
            f for f in directory.iterdir()
            if f.is_file()
            and f.suffix.lower() == ".pdf"
            and not _is_output_file(f.stem.lower())
        ]
        logger.info(f"Found {len(pdf_files)} PDF files (excluding output files)")
        return sorted(pdf_files)
    except PermissionError as e:
        logger.error(f"Permission denied scanning {directory}: {e}")
        return []
    except OSError as e:
        logger.error(f"Failed to get PDF files from {directory}: {e}")
        return []


@lru_cache(maxsize=512)
def generate_search_patterns(filename: str) -> list[str]:
    """
    Generate comprehensive search patterns with extension and separator variations.

    This function creates multiple pattern variations to maximize recognition rate:
    - Original filename (as-is)
    - All-lowercase variant
    - All-uppercase variant
    - Extension variations (.jpg vs .jpeg)
    - Separator variations (underscore, hyphen, space)
    - Individual words from multi-word stems (split by space/underscore/hyphen)
    - Stem-only patterns (without extension)

    All search paths (PyMuPDF, OCR, fuzzy, report counting) perform
    case-insensitive matching. Patterns are deduplicated by their lowercase
    form to prevent inflated match counts.

    Args:
        filename: The filename to generate patterns for (e.g., "IMG_1234.jpg")

    Returns:
        List of unique patterns (case-insensitive deduplicated)
    """
    patterns = []
    stem = Path(filename).stem
    ext = Path(filename).suffix.lower()

    # 1. Original filename + case variants
    patterns.append(filename)
    patterns.append(filename.lower())
    patterns.append(filename.upper())

    # 2. Stem (without extension) - only if long enough to avoid false positives
    if len(stem) >= MIN_STEM_LENGTH:
        patterns.append(stem)
        patterns.append(stem.lower())
        patterns.append(stem.upper())

    # 3. Extension variations
    ext_map = {
        '.jpg': ['.jpg', '.jpeg'],
        '.jpeg': ['.jpg', '.jpeg'],
        '.pdf': ['.pdf'],
        '.txt': ['.txt'],
        '.doc': ['.doc', '.docx'],
        '.docx': ['.doc', '.docx'],
        '.png': ['.png'],
        '.gif': ['.gif'],
        '.bmp': ['.bmp'],
        '.tif': ['.tif', '.tiff'],
        '.tiff': ['.tif', '.tiff'],
    }

    if ext in ext_map:
        for var_ext in ext_map[ext]:
            patterns.append(stem + var_ext)

    # 4. Separator variations (if filename has underscores or hyphens)
    if '_' in stem or '-' in stem:
        stem_normalized = stem.replace('-', '_')
        patterns.append(stem_normalized.replace('_', '-') + ext)
        patterns.append(stem_normalized.replace('_', ' ') + ext)

    # 5. Individual words from filename stem (split by spaces, underscores, hyphens)
    words = re.split(r'[\s_\-]+', stem)
    if len(words) > 1:
        for word in words:
            if len(word) >= MIN_STEM_LENGTH:
                patterns.append(word)

    # 6. Case-insensitive deduplication (prevents inflated match counts)
    seen: set[str] = set()
    unique = []
    for p in patterns:
        key = p.lower()
        if key not in seen:
            seen.add(key)
            unique.append(p)

    logger.debug(f"Generated {len(unique)} patterns for '{filename}': {unique[:5]}...")
    return unique


def build_pattern_set(filenames: list[str],
                      strip_single_words: bool = False) -> set[str]:
    """Build a deduplicated set of all search patterns for a list of filenames.

    Generates patterns once for all filenames and deduplicates case-insensitively.
    Use this before batch processing to avoid repeated per-file pattern generation.

    Args:
        filenames: List of filenames to generate patterns for
        strip_single_words: If True, remove patterns that are single words without
            a file extension (e.g. "test", "report"). Recommended for TXT/Docs
            searches to reduce false positives from short stem fragments.

    Returns:
        Set of unique patterns (case-insensitive deduplicated, capped at MAX_SEARCH_PATTERNS)
    """
    all_patterns: set[str] = set()
    seen_lower: set[str] = set()
    for filename in filenames:
        for p in generate_search_patterns(filename):
            key = p.lower()
            if key not in seen_lower:
                seen_lower.add(key)
                all_patterns.add(p)

    if strip_single_words:
        # Remove bare word patterns (no dot, no separator) shorter than a useful threshold.
        # Keeps full filenames (contain "."), separator-variants (contain "-" or " "),
        # and stems that are at least MIN_STEM_LENGTH*2 characters to avoid false positives.
        min_bare_len = max(MIN_STEM_LENGTH * 2, 8)
        filtered = {
            p for p in all_patterns
            if '.' in p or '-' in p or ' ' in p or len(p) >= min_bare_len
        }
        removed = len(all_patterns) - len(filtered)
        if removed:
            logger.info(f"Stripped {removed} short bare-word patterns (false-positive reduction)")
        all_patterns = filtered

    if len(all_patterns) > MAX_SEARCH_PATTERNS:
        total = len(all_patterns)
        logger.warning(
            f"Pattern count {total} exceeds limit {MAX_SEARCH_PATTERNS}, truncating"
        )
        all_patterns = set(list(all_patterns)[:MAX_SEARCH_PATTERNS])
        print(
            f"  {_C.YELLOW}⚠ Pattern limit: only {MAX_SEARCH_PATTERNS} of {total} patterns used. "
            f"Increase MAX_SEARCH_PATTERNS in Settings.{_C.RESET}"
        )
    logger.info(f"Built {len(all_patterns)} unique search patterns for {len(filenames)} filenames")
    return all_patterns


def build_stem_pattern_set(filenames: list[str], min_length: int = 0) -> set[str]:
    """Build a minimal pattern set containing only stems (filename without extension).

    Each entry is the bare stem of the filename. If the filename already has no
    extension (e.g. from a stems-only JSON), it is used as-is.
    Patterns are deduplicated case-insensitively.

    This is significantly faster than build_pattern_set() because it generates
    exactly one pattern per unique stem instead of 10-20 variants.

    Args:
        filenames: List of filenames (with or without extension)
        min_length: Minimum stem length to include (0 = no restriction).
            Stems shorter than this value are silently skipped, preventing
            very short names like "9" or "a" from producing false positives.

    Returns:
        Set of unique stems (case-insensitive deduplicated)
    """
    seen: set[str] = set()
    result: set[str] = set()
    skipped = 0
    for f in filenames:
        stem = Path(f).stem or f  # stem is empty only for dotfiles like ".gitignore"
        if min_length > 0 and len(stem) < min_length:
            skipped += 1
            continue
        key = stem.lower()
        if key not in seen:
            seen.add(key)
            result.add(stem)
    if skipped:
        logger.info(f"Skipped {skipped} stems shorter than {min_length} chars")
    logger.info(f"Built {len(result)} stem patterns for {len(filenames)} filenames")
    return result


def check_pdf_has_native_text(doc: fitz.Document) -> bool:
    """
    Check if PDF has extractable native text (not image-only).

    Args:
        doc: PyMuPDF document object

    Returns:
        True if any page has text content, False if image-only PDF
    """
    logger.info("Checking for native text in PDF...")
    for page_num, page in enumerate(doc, 1):
        text = page.get_text("text").strip()
        if text:
            logger.info(f"Native text found on page {page_num}")
            return True
    logger.info("No native text found - PDF is image-only")
    return False


def setup_tesseract_path() -> str | None:
    """
    Setup Tesseract path with auto-detection and explicit configuration.

    Priority order:
    1. Explicit TESSERACT_PATH from config section
    2. Common Windows installation paths
    3. System PATH (Linux/Mac default)

    Returns:
        Path to tesseract executable if found, None otherwise
    """
    # Priority 1: Check explicit path from config
    if TESSERACT_PATH and Path(TESSERACT_PATH).exists():
        logger.info(f"Using configured Tesseract path: {TESSERACT_PATH}")
        return str(TESSERACT_PATH)

    # Priority 2: Auto-detect common Windows locations
    if platform.system() == "Windows":
        common_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            Path.home() / 'AppData' / 'Local' / 'Programs' / 'Tesseract-OCR' / 'tesseract.exe',
            r'C:\Tesseract-OCR\tesseract.exe',
        ]

        for path in common_paths:
            path_obj = Path(path)
            if path_obj.exists():
                logger.info(f"Auto-detected Tesseract at: {path}")
                return str(path_obj)

    # Priority 3: Check system PATH
    tesseract_cmd = shutil.which('tesseract')
    if tesseract_cmd:
        logger.info(f"Found Tesseract in system PATH: {tesseract_cmd}")
        return tesseract_cmd

    # Not found
    logger.warning("Tesseract not found. Checked:")
    logger.warning("  1. Configured TESSERACT_PATH")
    logger.warning("  2. Common Windows install locations")
    logger.warning("  3. System PATH")
    return None


# Cache for Tesseract path to avoid repeated lookups
_TESSERACT_PATH_CACHE = None


def get_tesseract_path() -> str | None:
    """
    Get cached Tesseract path or perform lookup if not cached.

    Returns:
        Path to tesseract executable if found, None otherwise
    """
    global _TESSERACT_PATH_CACHE
    if _TESSERACT_PATH_CACHE is None:
        _TESSERACT_PATH_CACHE = setup_tesseract_path()
    return _TESSERACT_PATH_CACHE


def check_tesseract_installed() -> bool:
    """
    Check if Tesseract OCR is installed and configure pytesseract.

    Returns:
        True if Tesseract is installed and accessible, False otherwise
    """
    try:
        import pytesseract

        # Setup Tesseract path
        tesseract_path = get_tesseract_path()
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path

        # Verify it works
        version = pytesseract.get_tesseract_version()
        logger.info(f"Tesseract OCR v{version} is available")
        return True

    except ImportError:
        logger.warning("pytesseract is not installed — OCR unavailable. "
                       "Install with: pip install pytesseract")
        return False
    except (EnvironmentError, FileNotFoundError, OSError) as e:
        logger.warning(f"Tesseract binary not found or not accessible: {e}")
        return False


# Check Tesseract availability at module load
TESSERACT_AVAILABLE = check_tesseract_installed()

if not TESSERACT_AVAILABLE:
    logger.warning("Tesseract OCR not installed. OCR fallback will be disabled.")
    logger.warning("Install Tesseract: https://github.com/tesseract-ocr/tesseract")


def extract_text_and_boxes_via_ocr(page: fitz.Page, dpi: int = 200) -> dict:
    """
    Extract text and bounding boxes from page using OCR (for image-only PDFs).

    Uses pytesseract.image_to_data() to get word-level bounding boxes for
    precise text highlighting.

    Args:
        page: PyMuPDF page object
        dpi: Resolution for OCR (default 200)

    Returns:
        Dictionary with OCR data including text and bounding boxes
        Format: {
            'text': [...],      # List of words
            'left': [...],      # Left coordinate (pixmap)
            'top': [...],       # Top coordinate (pixmap)
            'width': [...],     # Width (pixmap)
            'height': [...],    # Height (pixmap)
            'conf': [...],      # Confidence scores
            'scale_x': float,   # Scale factor for X (pixmap to PDF)
            'scale_y': float    # Scale factor for Y (pixmap to PDF)
        }
    """
    try:
        import pytesseract
        from PIL import Image

        # Setup Tesseract path (critical for Windows)
        tesseract_path = get_tesseract_path()
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        elif not TESSERACT_AVAILABLE:
            logger.error("Tesseract not available - cannot perform OCR")
            return {}

        # Convert page to pixmap at specified DPI
        matrix = fitz.Matrix(dpi / 72, dpi / 72)
        pixmap = page.get_pixmap(matrix=matrix)

        # Convert to PIL Image
        img_data = pixmap.tobytes("png")
        img = Image.open(io.BytesIO(img_data))

        # Preprocessing for better OCR results
        img = img.convert('L')  # Convert to grayscale
        # Contrast enhancement
        from PIL import ImageEnhance, ImageFilter
        enhancer = ImageEnhance.Contrast(img)
        img = enhancer.enhance(2.0)  # Double contrast
        # Optional: Light sharpening
        img = img.filter(ImageFilter.SHARPEN)

        # Perform OCR with word-level bounding boxes
        # output_type=Output.DICT returns a dictionary with word positions
        from pytesseract import Output
        ocr_data = pytesseract.image_to_data(
            img, lang=OCR_LANGUAGES,
            output_type=Output.DICT,
            config=f"--psm {OCR_PSM} --oem 3"
        )

        # Calculate scale factors to convert pixmap coords to PDF coords
        pdf_rect = page.rect
        scale_x = pdf_rect.width / pixmap.width
        scale_y = pdf_rect.height / pixmap.height

        ocr_data['scale_x'] = scale_x
        ocr_data['scale_y'] = scale_y

        # Cleanup — release pixmap and PIL image to reduce peak RAM during batch OCR
        img.close()
        del pixmap

        # Count valid words (confidence > 0)
        word_count = sum(1 for conf in ocr_data.get('conf', []) if int(conf) > 0)
        logger.info(f"OCR extracted {word_count} words with bounding boxes")

        return ocr_data

    except ImportError:
        logger.error("pytesseract not installed. Install: pip install pytesseract")
        logger.error("Also install Tesseract: https://github.com/tesseract-ocr/tesseract")
        return {}
    except (OSError, RuntimeError, EnvironmentError) as e:
        logger.error(f"OCR failed: {e}")
        return {}


def search_ocr_data_and_highlight(page: fitz.Page, ocr_data: dict, patterns: list[str]) -> dict[str, int]:
    """
    Search OCR data for patterns and add precise highlights using bounding boxes.

    Uses a two-pass approach:
    1. Fast exact word matching for single words
    2. Fulltext substring search with position tracking for multi-character patterns
       (handles Tesseract tokenization like "report_2023.pdf" → ["report", "_", "2023", ".", "pdf"])

    Args:
        page: PyMuPDF page object
        ocr_data: Dictionary from pytesseract.image_to_data() with bounding boxes
        patterns: List of search patterns

    Returns:
        Dict mapping each matched pattern (lowercased) to its hit count
    """
    if not ocr_data or 'text' not in ocr_data:
        logger.warning("No OCR data available for highlighting")
        return {}

    hit_counts: dict[str, int] = {}
    scale_x = ocr_data.get('scale_x', 1.0)
    scale_y = ocr_data.get('scale_y', 1.0)
    pdf_rect = page.rect

    # Build a list of (word, bbox) tuples from OCR data
    words_with_boxes = []
    for i, word in enumerate(ocr_data['text']):
        if not word or not word.strip():
            continue

        conf = int(ocr_data['conf'][i])
        if conf < OCR_MIN_CONFIDENCE:  # Skip low confidence detections
            continue

        # Get pixmap coordinates
        x = ocr_data['left'][i]
        y = ocr_data['top'][i]
        w = ocr_data['width'][i]
        h = ocr_data['height'][i]

        # Convert to PDF coordinates
        pdf_x0 = pdf_rect.x0 + (x * scale_x)
        pdf_y0 = pdf_rect.y0 + (y * scale_y)
        pdf_x1 = pdf_x0 + (w * scale_x)
        pdf_y1 = pdf_y0 + (h * scale_y)

        bbox = fitz.Rect(pdf_x0, pdf_y0, pdf_x1, pdf_y1)
        words_with_boxes.append((word.strip(), bbox))

    # Build fulltext with position tracking (char offset → word index).
    # Words are separated by a single space so that multi-token patterns like
    # "report 2024" are found even when OCR tokenised them into separate words.
    # The space character is mapped to the preceding word's index so that
    # substring matches spanning a space boundary resolve to the correct boxes.
    fulltext = ""
    char_to_word_map = []  # Maps each character position to word index
    for word_idx, (word, _) in enumerate(words_with_boxes):
        fulltext += word
        for _ in range(len(word)):
            char_to_word_map.append(word_idx)
        # Space separator (except after the last word)
        if word_idx < len(words_with_boxes) - 1:
            fulltext += " "
            char_to_word_map.append(word_idx)  # space belongs to preceding word

    fulltext_lower = fulltext.lower()

    # Track already highlighted word indices to avoid duplicates
    highlighted_words = set()

    # Search for patterns
    for pattern in patterns:
        pattern_lower = pattern.lower()

        # Pass 1: Fast exact word matching (for single words)
        for word_idx, (word, bbox) in enumerate(words_with_boxes):
            if word_idx in highlighted_words:
                continue
            if word.lower() == pattern_lower:
                try:
                    highlight = page.add_highlight_annot(bbox)
                    highlight.set_colors(stroke=HIGHLIGHT_COLOR)
                    highlight.update()
                    hit_counts[pattern_lower] = hit_counts.get(pattern_lower, 0) + 1
                    highlighted_words.add(word_idx)
                    logger.info(f"OCR exact match: '{pattern}' at bbox {bbox}")
                except (RuntimeError, TypeError, ValueError) as e:
                    logger.error(f"Failed to add OCR highlight annotation: {e}")

        # Pass 2: Fulltext substring search (for multi-token patterns)
        # Find all occurrences of pattern in fulltext
        search_start = 0
        while True:
            pos = fulltext_lower.find(pattern_lower, search_start)
            if pos == -1:
                break

            # Determine which words are covered by this match
            start_char = pos
            end_char = pos + len(pattern_lower)

            # Boundary check
            if start_char >= len(char_to_word_map) or end_char > len(char_to_word_map):
                search_start = pos + 1
                continue

            # Get word indices covered by this match
            start_word_idx = char_to_word_map[start_char]
            end_word_idx = char_to_word_map[end_char - 1]

            # Collect bboxes for all words in this range
            boxes_to_combine = []
            for word_idx in range(start_word_idx, end_word_idx + 1):
                if word_idx not in highlighted_words:
                    boxes_to_combine.append(words_with_boxes[word_idx][1])
                    highlighted_words.add(word_idx)

            if boxes_to_combine:
                # Create union of all boxes
                combined_bbox = boxes_to_combine[0]
                for box in boxes_to_combine[1:]:
                    combined_bbox = combined_bbox | box  # Union operator

                try:
                    highlight = page.add_highlight_annot(combined_bbox)
                    highlight.set_colors(stroke=HIGHLIGHT_COLOR)
                    highlight.update()
                    hit_counts[pattern_lower] = hit_counts.get(pattern_lower, 0) + 1
                    logger.info(f"OCR substring match: '{pattern}' at bbox {combined_bbox}")
                except (RuntimeError, TypeError, ValueError) as e:
                    logger.error(f"Failed to add OCR highlight annotation: {e}")

            search_start = pos + 1

    return hit_counts


def _build_words_with_boxes_from_ocr(
    ocr_data: dict, page_rect: fitz.Rect
) -> list[tuple[str, fitz.Rect]]:
    """Extract (word, bbox) pairs from OCR data for fuzzy matching.

    Args:
        ocr_data: Dictionary from pytesseract.image_to_data()
        page_rect: PDF page rectangle for coordinate conversion

    Returns:
        List of (word, fitz.Rect) tuples
    """
    scale_x = ocr_data.get('scale_x', 1.0)
    scale_y = ocr_data.get('scale_y', 1.0)
    result = []
    for i, word in enumerate(ocr_data['text']):
        if not word or not word.strip():
            continue
        conf = int(ocr_data['conf'][i])
        if conf < OCR_MIN_CONFIDENCE:
            continue
        x = ocr_data['left'][i]
        y = ocr_data['top'][i]
        w = ocr_data['width'][i]
        h = ocr_data['height'][i]
        pdf_x0 = page_rect.x0 + (x * scale_x)
        pdf_y0 = page_rect.y0 + (y * scale_y)
        pdf_x1 = pdf_x0 + (w * scale_x)
        pdf_y1 = pdf_y0 + (h * scale_y)
        result.append((word.strip(), fitz.Rect(pdf_x0, pdf_y0, pdf_x1, pdf_y1)))
    return result


def fuzzy_search_words_and_highlight(
    page: fitz.Page,
    words_with_boxes: list[tuple[str, fitz.Rect]],
    filenames: list[str],
    threshold: float = FUZZY_THRESHOLD,
) -> dict[str, int]:
    """Fuzzy-search filenames in word+bbox data and add highlights.

    Uses rapidfuzz to find approximate matches. Compares each filename
    (and its stem) against individual words and sliding windows of 2-5
    consecutive words.

    Args:
        page: PyMuPDF page object
        words_with_boxes: List of (word, bounding_rect) tuples
        filenames: Original filenames to search for
        threshold: Minimum similarity ratio (0.0-1.0)

    Returns:
        Dict mapping each matched target (lowercased) to its hit count
    """
    if not words_with_boxes:
        return {}

    # Build search targets from filenames (deduplicated)
    targets: list[str] = []
    for fn in filenames:
        targets.append(fn.lower())
        stem = Path(fn).stem.lower()
        if len(stem) >= MIN_STEM_LENGTH and stem != fn.lower():
            targets.append(stem)
    targets = list(dict.fromkeys(targets))

    hit_counts: dict[str, int] = {}
    highlighted: set[int] = set()

    for target in targets:
        target_len = len(target)
        max_window = min(5, len(words_with_boxes))

        for i in range(len(words_with_boxes)):
            if i in highlighted:
                continue

            matched = False
            for win_size in range(1, max_window + 1):
                if i + win_size > len(words_with_boxes):
                    break

                window = words_with_boxes[i:i + win_size]
                combined = "".join(w[0] for w in window).lower()

                # Length pre-filter to skip obviously non-matching windows
                if not combined:
                    continue
                comb_len = len(combined)
                if comb_len < target_len * 0.5 or comb_len > target_len * 2.0:
                    continue

                ratio = _rfuzz.ratio(target, combined) / 100.0
                if ratio >= threshold:
                    indices = set(range(i, i + win_size))
                    if indices & highlighted:
                        break  # already highlighted — skip larger windows too

                    # Combine bounding boxes
                    bbox = window[0][1]
                    for _, r in window[1:]:
                        bbox = bbox | r

                    try:
                        annot = page.add_highlight_annot(bbox)
                        annot.set_colors(stroke=HIGHLIGHT_COLOR)
                        annot.update()
                        hit_counts[target] = hit_counts.get(target, 0) + 1
                        highlighted.update(indices)
                        logger.info(f"Fuzzy match ({ratio:.0%}): '{target}' ~ '{combined}'")
                    except (RuntimeError, TypeError, ValueError) as e:
                        logger.error(f"Fuzzy highlight failed: {e}")
                    matched = True
                    break  # Match found at this position, skip larger windows

            if matched:
                continue  # move to next position

    return hit_counts


def _normalize_text(s: str) -> str:
    """Normalize text for matching by resolving Unicode and OCR artefacts.

    Applies NFC Unicode normalization, replaces common typographic ligatures
    with their ASCII equivalents, and collapses runs of whitespace.  Call this
    on both the source text and any search pattern before comparing them so that
    differences in Unicode representation or OCR output do not cause missed
    matches.

    Common ligatures handled: ﬁ→fi, ﬂ→fl, ﬀ→ff, ﬃ→ffi, ﬄ→ffl, ﬆ→st.

    Args:
        s: Raw input string to normalize.

    Returns:
        The normalized string with ligatures expanded and whitespace collapsed.
    """
    s = unicodedata.normalize('NFC', s)
    ligatures = {'ﬁ': 'fi', 'ﬂ': 'fl', 'ﬀ': 'ff', 'ﬃ': 'ffi', 'ﬄ': 'ffl', 'ﬆ': 'st'}
    for lig, repl in ligatures.items():
        s = s.replace(lig, repl)
    return re.sub(r'[ \t]+', ' ', s).strip()


def _normalize_sep(s: str) -> str:
    """Normalize filename separators to spaces for fuzzy comparison.

    Replaces hyphens, underscores, and dots with spaces, then collapses
    runs of whitespace.  Applied to both the search pattern and candidate
    window before similarity scoring so that e.g. ``"report-2024"`` and
    ``"report 2024"`` are treated as equivalent.

    Args:
        s: Input string whose separators should be normalised.

    Returns:
        The string with ``-``, ``_``, and ``.`` replaced by spaces and
        consecutive spaces collapsed.
    """
    return re.sub(r'[\-_.]+', ' ', s).strip()


def _is_short_bare(pattern: str) -> bool:
    """Return True if pattern is a short bare word needing word-boundary guarding.

    A pattern qualifies when it contains no dot, hyphen, or space (i.e. it is a
    plain alphanumeric token without separators) AND its length is below the
    threshold used in build_pattern_set's strip_single_words filter.  Such
    patterns risk substring false positives and should only match at word
    boundaries.

    Args:
        pattern: A single search pattern string to evaluate.

    Returns:
        ``True`` if the pattern is a short bare word; ``False`` otherwise.
    """
    min_len = max(MIN_STEM_LENGTH * 2, 8)
    return '.' not in pattern and '-' not in pattern and ' ' not in pattern and len(pattern) < min_len


def fuzzy_text_search(text: str, patterns: list[str],
                      threshold: float = FUZZY_THRESHOLD) -> list[tuple[str, str]]:
    """Fuzzy-match patterns against text using sliding windows.

    For non-PDF file types where bounding boxes aren't needed.
    Finds ALL non-overlapping occurrences of each pattern, not just the first.

    Before computing similarity, both the target pattern and the candidate window
    are normalized via ``_normalize_sep`` so that separators ``-``, ``_``, ``.``
    are treated the same as spaces.  The original (non-normalized) window text is
    returned so that callers can place markers at the exact source location.

    Args:
        text: Full text to search in
        patterns: List of pattern strings to match
        threshold: Minimum similarity ratio (0.0-1.0)

    Returns:
        List of (pattern, matched_window) tuples, one per occurrence.
        `pattern`        — the search pattern that matched
        `matched_window` — the actual text fragment from `text` that matched
    """
    words = text.split()
    text_lower_words = [w.lower() for w in words]
    n_words = len(text_lower_words)
    matched: list[tuple[str, str]] = []

    targets = set()
    for p in patterns:
        targets.add(p.lower())

    for target in targets:
        # Normalize separators for comparison; use normalized form for word count
        # so that "report-2024" (1 raw word) is treated as 2 words for window sizing.
        target_norm = _normalize_sep(target)
        target_word_count = len(target_norm.split())
        max_window = min(target_word_count + 3, 6)
        skip_until = 0  # first word-index not yet consumed by a match

        for i in range(n_words):
            if i < skip_until:
                continue
            for window_size in range(1, max_window + 1):
                if i + window_size > n_words:
                    break
                combined = " ".join(text_lower_words[i:i + window_size])
                combined_norm = _normalize_sep(combined)
                # Length guard on normalized forms
                if abs(len(combined_norm) - len(target_norm)) > len(target_norm) * 0.5:
                    continue
                ratio = _rfuzz.ratio(target_norm, combined_norm) / 100.0
                if ratio >= threshold:
                    # Reconstruct the original-case window from the source text
                    original_window = " ".join(words[i:i + window_size])
                    matched.append((target, original_window))
                    skip_until = i + window_size  # advance past matched window
                    break  # move on to next start position

    return matched


def search_and_highlight_pdf(
    pdf_path: Path, filenames: list[str], output_suffix: str = OUTPUT_SUFFIX,
    fuzzy: bool = False,
    progress_callback=None,
    pattern_set: set[str] | None = None,
) -> tuple[int, Path | None, bool, bool]:
    """
    Search for filenames in PDF with OCR fallback and highlight matches.
    Only creates output PDF if matches are found.

    For native text PDFs: Uses fast PyMuPDF text search with exact positioning
    For image-only PDFs: Falls back to OCR with word-level bounding boxes for precise highlighting

    When fuzzy=True, uses Levenshtein-based approximate matching (slower but
    catches OCR errors and typos).

    Args:
        pdf_path: Path to the PDF file
        filenames: List of filenames from JSON to search for
        output_suffix: Suffix to append to output filename (before .pdf)
        fuzzy: If True, use fuzzy matching instead of exact pattern matching
        progress_callback: Called with (page_num, total_pages) for progress updates

    Returns:
        Tuple of (total_hits, output_path or None, has_native_text, ocr_used)

    Raises:
        IOError: If PDF cannot be read or written
    """
    try:
        logger.info(f"Processing PDF: {pdf_path.name} (fuzzy={'ON' if fuzzy else 'OFF'})")

        # Suppress MuPDF C-level warnings from polluting stdout (e.g. fill color warnings)
        # Warnings are still collected via fitz.TOOLS.mupdf_warnings() afterwards.
        fitz.TOOLS.mupdf_display_errors(False)
        fitz.TOOLS.mupdf_warnings()  # Clear any accumulated warnings

        # Attempt to open PDF with validation
        try:
            doc = fitz.open(pdf_path)
        except (RuntimeError, IOError, OSError) as e:
            logger.error(f"Cannot open PDF {pdf_path.name}: {e}")
            raise IOError(f"PDF appears to be corrupt or invalid: {e}")

        # Check for encrypted PDF
        if doc.is_encrypted:
            logger.warning(f"PDF {pdf_path.name} is encrypted/password-protected")
            doc.close()
            raise IOError(f"PDF is encrypted/password-protected: {pdf_path.name}. Please provide an unencrypted version.")

        # Validate PDF has pages
        if doc.page_count == 0:
            doc.close()
            raise ValueError(f"PDF has no pages (may be corrupt): {pdf_path.name}")

        # Step 1: Check for native text
        has_native_text = check_pdf_has_native_text(doc)
        ocr_used = False
        total_hits = 0
        collected_text_parts: list[str] = []

        if not fuzzy:
            # Use pre-built pattern set if provided, otherwise build it now
            if pattern_set is not None:
                all_patterns = pattern_set
            else:
                all_patterns = build_pattern_set(filenames)
            logger.info(f"Total unique search patterns: {len(all_patterns)}")

        # Step 2: Process based on native text availability
        _search_hit_counts: dict[str, int] = {}
        if has_native_text:
            if fuzzy:
                # NATIVE TEXT + FUZZY - word-level fuzzy matching
                logger.info("Using native text fuzzy search")
                for page_num, page in enumerate(doc, 1):
                    if progress_callback and page_num % 5 == 0:
                        progress_callback(page_num, len(doc))
                    collected_text_parts.append(page.get_text("text"))
                    words_data = page.get_text("words")
                    words_with_boxes = [
                        (wd[4].strip(), fitz.Rect(wd[0], wd[1], wd[2], wd[3]))
                        for wd in words_data if wd[4].strip()
                    ]
                    page_counts = fuzzy_search_words_and_highlight(page, words_with_boxes, filenames)
                    for tgt, cnt in page_counts.items():
                        _search_hit_counts[tgt] = _search_hit_counts.get(tgt, 0) + cnt
                    total_hits += sum(page_counts.values())
                    if page_num % 10 == 0 or page_num == len(doc):
                        logger.info(f"Fuzzy page {page_num}/{len(doc)} - {total_hits} matches so far")
            else:
                # NATIVE TEXT + EXACT - fast PyMuPDF search
                logger.info("Using native text search (fast)")
                found_patterns = set()

                for page_num, page in enumerate(doc, 1):
                    if progress_callback and page_num % 5 == 0:
                        progress_callback(page_num, len(doc))
                    page_text = page.get_text("text")
                    collected_text_parts.append(page_text)
                    highlighted_regions = set()

                    for pattern in all_patterns:
                        # Word-boundary guard: short bare-word patterns (no dot/hyphen/space,
                        # length < 8) are only highlighted when the pattern appears as a
                        # complete token in the page text, preventing substring false positives
                        # such as "abc" matching inside "fabricate".
                        if _is_short_bare(pattern) and not re.search(
                            r'(?<!\w)' + re.escape(pattern) + r'(?!\w)',
                            page_text, re.IGNORECASE
                        ):
                            continue
                        instances = page.search_for(pattern)
                        for bbox in instances:
                            bbox_key = (round(bbox.x0, 2), round(bbox.y0, 2),
                                        round(bbox.x1, 2), round(bbox.y1, 2))
                            if bbox_key not in highlighted_regions:
                                highlighted_regions.add(bbox_key)
                                highlight = page.add_highlight_annot(bbox)
                                highlight.set_colors(stroke=HIGHLIGHT_COLOR)
                                highlight.update()
                                total_hits += 1
                                found_patterns.add(pattern)
                                _search_hit_counts[pattern.lower()] = _search_hit_counts.get(pattern.lower(), 0) + 1

                    if page_num % 10 == 0 or page_num == len(doc):
                        logger.info(f"Page {page_num}/{len(doc)} - {total_hits} matches, {len(found_patterns)}/{len(all_patterns)} unique patterns found")

        else:
            # IMAGE-ONLY PATH - OCR fallback
            if not TESSERACT_AVAILABLE:
                logger.warning("PDF has no native text, but Tesseract not available. Skipping OCR.")
                ocr_used = False
            else:
                logger.info(f"No native text - using OCR fallback (fuzzy={'ON' if fuzzy else 'OFF'})")
                ocr_used = True

                for page_num, page in enumerate(doc, 1):
                    if progress_callback and page_num % 5 == 0:
                        progress_callback(page_num, len(doc))
                    logger.info(f"OCR processing page {page_num}/{len(doc)}...")
                    ocr_data = extract_text_and_boxes_via_ocr(page, dpi=OCR_DPI)

                    if ocr_data and 'text' in ocr_data:
                        collected_text_parts.append(
                            " ".join(w for w in ocr_data['text'] if w and w.strip()))

                        if fuzzy:
                            words_boxes = _build_words_with_boxes_from_ocr(ocr_data, page.rect)
                            page_counts = fuzzy_search_words_and_highlight(page, words_boxes, filenames)
                            for tgt, cnt in page_counts.items():
                                _search_hit_counts[tgt] = _search_hit_counts.get(tgt, 0) + cnt
                            hits = sum(page_counts.values())
                        else:
                            page_counts = search_ocr_data_and_highlight(page, ocr_data, list(all_patterns))
                            for pat, cnt in page_counts.items():
                                _search_hit_counts[pat] = _search_hit_counts.get(pat, 0) + cnt
                            hits = sum(page_counts.values())
                        total_hits += hits

                    ocr_data = None

        # Step 3: Save if hits found
        output_path = None
        if total_hits > 0:
            output_path = _output_path(pdf_path, output_suffix)
            doc.save(output_path)
            logger.info(f"Saved: {output_path.name} with {total_hits} highlights")
        else:
            logger.info("No matches found - no output PDF created")

        doc.close()

        # Cache extracted text for report generation (avoids double OCR)
        _cache_extracted_text(pdf_path, "\n".join(collected_text_parts))
        _cache_match_counts(pdf_path, _search_hit_counts)

        return total_hits, output_path, has_native_text, ocr_used

    except (OSError, RuntimeError, ValueError) as e:
        logger.error(f"Error processing PDF {pdf_path.name}: {e}")
        return 0, None, False, False
    finally:
        # Restore MuPDF error display so other operations show errors normally
        fitz.TOOLS.mupdf_display_errors(True)


def get_jpg_files(directory: Path = Path(".")) -> list[Path]:
    """
    Get all JPG files in directory (excluding *_found.jpg files).

    Args:
        directory: Directory to search

    Returns:
        List of Path objects for JPG files (sorted, excluding _found files)
    """
    try:
        jpg_extensions = ['.jpg', '.jpeg', '.JPG', '.JPEG']
        jpg_files = [
            f for f in directory.iterdir()
            if f.is_file()
            and f.suffix in jpg_extensions
            and not _is_output_file(f.stem.lower())
        ]
        logger.info(f"Found {len(jpg_files)} JPG files (excluding output files)")
        return sorted(jpg_files)
    except PermissionError as e:
        logger.error(f"Permission denied scanning {directory}: {e}")
        return []
    except OSError as e:
        logger.error(f"Failed to get JPG files from {directory}: {e}")
        return []


def search_filenames_in_jpg(jpg_path: Path, filenames: list[str],
                           fuzzy: bool = False,
                           pattern_set: set[str] | None = None) -> tuple[int, Path | None]:
    """
    Search for filenames in JPG metadata and OCR text.

    Args:
        jpg_path: Path to JPG file
        filenames: List of filenames to search for
        fuzzy: If True, use fuzzy matching instead of exact

    Returns:
        Tuple of (match_count, output_path or None)
    """
    try:
        from PIL import Image

        # Use pre-built pattern set if provided, otherwise build it now
        all_patterns = pattern_set if pattern_set is not None else build_pattern_set(filenames)

        # Build reverse map pattern_lower → original filename for per-filename counting.
        # Each filename may generate many patterns; we want to count distinct filenames,
        # not distinct patterns (which would inflate the count by 10-20×).
        pattern_to_filename: dict[str, str] = {}
        for fn in filenames:
            for p in generate_search_patterns(fn):
                pattern_to_filename.setdefault(p.lower(), fn)
            # Also map the stem and the bare name for stems-mode patterns
            pattern_to_filename.setdefault(fn.lower(), fn)
            pattern_to_filename.setdefault(Path(fn).stem.lower(), fn)

        matched_filenames: set[str] = set()  # unique logical filenames found
        collected_text_parts: list[str] = []
        pattern_hit_counts: dict[str, int] = {}

        # Step 1: Check EXIF metadata (fast)
        try:
            img = Image.open(jpg_path)
            exif_data = img.getexif()
            if exif_data:
                exif_text = ' '.join(str(v) for v in exif_data.values() if v)
                collected_text_parts.append(exif_text)
                if fuzzy:
                    for pattern, _window in fuzzy_text_search(exif_text, list(all_patterns)):
                        matched_filenames.add(pattern_to_filename.get(pattern, pattern))
                        pattern_hit_counts[pattern] = pattern_hit_counts.get(pattern, 0) + 1
                else:
                    for pattern in all_patterns:
                        if pattern.lower() in exif_text.lower():
                            matched_filenames.add(
                                pattern_to_filename.get(pattern.lower(), pattern))
                            logger.info(f"JPG EXIF match: '{pattern}' in {jpg_path.name}")
                            pattern_hit_counts[pattern.lower()] = pattern_hit_counts.get(pattern.lower(), 0) + 1
        except (OSError, AttributeError, KeyError) as e:
            logger.warning(f"EXIF read failed for {jpg_path.name}: {e}")

        # Step 2: Perform OCR on image (if Tesseract available)
        if TESSERACT_AVAILABLE:
            try:
                import pytesseract

                tesseract_path = get_tesseract_path()
                if tesseract_path:
                    pytesseract.pytesseract.tesseract_cmd = tesseract_path

                img = Image.open(jpg_path)
                ocr_text = pytesseract.image_to_string(img, lang=OCR_LANGUAGES, config=f"--psm {OCR_PSM} --oem 3")
                collected_text_parts.append(ocr_text)

                if fuzzy:
                    for pattern, _window in fuzzy_text_search(ocr_text, list(all_patterns)):
                        matched_filenames.add(pattern_to_filename.get(pattern, pattern))
                        pattern_hit_counts[pattern] = pattern_hit_counts.get(pattern, 0) + 1
                else:
                    for pattern in all_patterns:
                        if pattern.lower() in ocr_text.lower():
                            matched_filenames.add(
                                pattern_to_filename.get(pattern.lower(), pattern))
                            logger.info(f"JPG OCR match: '{pattern}' in {jpg_path.name}")
                            pattern_hit_counts[pattern.lower()] = pattern_hit_counts.get(pattern.lower(), 0) + 1
            except (OSError, RuntimeError, EnvironmentError) as e:
                logger.warning(f"OCR failed for {jpg_path.name}: {e}")

        # Cache extracted text for report generation (avoids double OCR)
        _cache_extracted_text(jpg_path, " ".join(collected_text_parts))
        _cache_match_counts(jpg_path, pattern_hit_counts)

        # Step 3: If matches found, copy to *_found.jpg
        matches_found = len(matched_filenames)
        if matches_found > 0:
            logger.info(f"JPG {jpg_path.name}: {matches_found} unique filename(s) found: "
                        f"{sorted(matched_filenames)}")
            output_path = _output_path(jpg_path, IMAGE_FOUND_SUFFIX)
            shutil.copy2(jpg_path, output_path)
            logger.info(f"Copied {jpg_path.name} → {output_path.name}")
            return matches_found, output_path

        return 0, None

    except ImportError:
        logger.error("Pillow not installed. Install: pip install Pillow")
        return 0, None
    except (OSError, RuntimeError, EnvironmentError) as e:
        logger.error(f"Error processing JPG {jpg_path.name}: {e}")
        return 0, None


def get_doc_files(directory: Path = Path("."),
                  extensions: list[str] | None = None) -> list[Path]:
    """
    Get all document files in directory (excluding *_checked files).

    Args:
        directory: Directory to search
        extensions: Optional list of extensions to include (e.g. ['.docx', '.xlsx']).
            Defaults to all supported document formats.

    Returns:
        List of Path objects for document files (sorted)
    """
    try:
        doc_extensions = extensions or ['.docx', '.xlsx', '.pptx', '.odt', '.ods', '.odp']
        doc_files = [
            f for f in directory.iterdir()
            if f.is_file()
            and f.suffix.lower() in doc_extensions
            and not _is_output_file(f.stem.lower())
        ]
        logger.info(f"Found {len(doc_files)} document files (excluding output files)")
        return sorted(doc_files)
    except PermissionError as e:
        logger.error(f"Permission denied scanning {directory}: {e}")
        return []
    except OSError as e:
        logger.error(f"Failed to get document files from {directory}: {e}")
        return []


def _color_paragraph_runs(paragraph, pattern_lower: str, color) -> int:
    """Color runs in a paragraph that cover pattern_lower, handling split runs.

    Builds a character→run index map from the concatenated paragraph text so
    that patterns spanning multiple runs (e.g. "report_" in run 1, "2024" in
    run 2) are found and both runs are colored.

    Works for both python-docx and python-pptx paragraph objects.

    Args:
        paragraph:     Paragraph object with `.text` and `.runs` attributes
        pattern_lower: Lower-cased search string
        color:         RGBColor object to apply (docx or pptx variant)

    Returns:
        Number of distinct occurrences colored
    """
    para_text = paragraph.text
    if not para_text:
        return 0
    para_lower = para_text.lower()
    if pattern_lower not in para_lower:
        return 0

    # Map each character position in para_text to its run index
    char_to_run_idx: list[int] = []
    for run_idx, run in enumerate(paragraph.runs):
        for _ in run.text:
            char_to_run_idx.append(run_idx)

    if not char_to_run_idx:
        return 0

    count = 0
    pos = 0
    while True:
        idx = para_lower.find(pattern_lower, pos)
        if idx == -1:
            break
        end = idx + len(pattern_lower)
        if end > len(char_to_run_idx):
            break
        first_run = char_to_run_idx[idx]
        last_run = char_to_run_idx[end - 1]
        for run in paragraph.runs[first_run:last_run + 1]:
            run.font.color.rgb = color
        count += 1
        pos = idx + 1

    return count


def search_in_docx(docx_path: Path, filenames: list[str],
                   fuzzy: bool = False,
                   pattern_set: set[str] | None = None) -> tuple[int, Path | None]:
    """
    Search Word document and add red highlights.

    Args:
        docx_path: Path to .docx file
        filenames: List of filenames to search for
        fuzzy: If True, use fuzzy matching

    Returns:
        Tuple of (match_count, output_path or None)
    """
    try:
        from docx import Document
        from docx.shared import RGBColor

        all_patterns = pattern_set if pattern_set is not None else build_pattern_set(filenames)
        red = RGBColor(255, 0, 0)

        doc = Document(docx_path)
        total_matches = 0
        pattern_hit_counts: dict[str, int] = {}

        def _search_paragraphs(paragraphs) -> int:
            """Search a list of paragraphs for all patterns and return the hit count.

            Args:
                paragraphs: Iterable of paragraph objects with a ``runs`` attribute.
            """
            count = 0
            for paragraph in paragraphs:
                for pattern in all_patterns:
                    n = _color_paragraph_runs(paragraph, pattern.lower(), red)
                    if n:
                        count += n
                        pattern_hit_counts[pattern.lower()] = pattern_hit_counts.get(pattern.lower(), 0) + n
                        logger.info(f"DOCX match: '{pattern}' ({n}×) in paragraph")
            return count

        if fuzzy:
            # Collect full document text and locate fuzzy matches
            body_parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        body_parts.append(cell.text)
            full_text = "\n".join(body_parts)
            fuzzy_hits = fuzzy_text_search(full_text, list(all_patterns))
            total_matches = len(fuzzy_hits)
            for pat, _win in fuzzy_hits:
                pattern_hit_counts[pat] = pattern_hit_counts.get(pat, 0) + 1
            # Mark runs using the actual matched window text (not the pattern)
            if total_matches > 0:
                windows = {window.lower() for _pat, window in fuzzy_hits}
                for paragraph in doc.paragraphs:
                    for window_lower in windows:
                        _color_paragraph_runs(paragraph, window_lower, red)
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for paragraph in cell.paragraphs:
                                for window_lower in windows:
                                    _color_paragraph_runs(paragraph, window_lower, red)
        else:
            # Exact: paragraph-level detection handles split runs via _color_paragraph_runs
            total_matches += _search_paragraphs(doc.paragraphs)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        total_matches += _search_paragraphs(cell.paragraphs)

        # Save if matches found
        _cache_match_counts(docx_path, pattern_hit_counts)
        if total_matches > 0:
            output_path = _output_path(docx_path, OUTPUT_SUFFIX)
            doc.save(output_path)
            return total_matches, output_path

        return 0, None

    except ImportError:
        logger.error("python-docx not installed. Install: pip install python-docx")
        return 0, None
    except (OSError, RuntimeError) as e:
        logger.error(f"Error processing DOCX {docx_path.name}: {e}")
        return 0, None


def search_in_xlsx(xlsx_path: Path, filenames: list[str],
                   fuzzy: bool = False,
                   pattern_set: set[str] | None = None) -> tuple[int, Path | None]:
    """
    Search Excel spreadsheet and highlight cells.

    Args:
        xlsx_path: Path to .xlsx file
        filenames: List of filenames to search for
        fuzzy: If True, use fuzzy matching

    Returns:
        Tuple of (match_count, output_path or None)
    """
    try:
        from openpyxl import load_workbook
        from openpyxl.styles import PatternFill

        all_patterns = pattern_set if pattern_set is not None else build_pattern_set(filenames)

        wb = load_workbook(xlsx_path)
        total_matches = 0
        pattern_hit_counts: dict[str, int] = {}

        red_fill = PatternFill(start_color='FF0000', end_color='FF0000', fill_type='solid')

        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        cell_text = str(cell.value)
                        if fuzzy:
                            fuzzy_hits = fuzzy_text_search(cell_text, list(all_patterns))
                            if fuzzy_hits:
                                cell.fill = red_fill
                                total_matches += len(fuzzy_hits)
                                for pat, _win in fuzzy_hits:
                                    pattern_hit_counts[pat] = pattern_hit_counts.get(pat, 0) + 1
                        else:
                            cell_lower = cell_text.lower()
                            for pattern in all_patterns:
                                if pattern.lower() in cell_lower:
                                    cell.fill = red_fill
                                    total_matches += 1
                                    logger.info(f"XLSX match: '{pattern}' in {sheet.title}!{cell.coordinate}")
                                    pattern_hit_counts[pattern.lower()] = pattern_hit_counts.get(pattern.lower(), 0) + 1

        _cache_match_counts(xlsx_path, pattern_hit_counts)
        if total_matches > 0:
            output_path = _output_path(xlsx_path, OUTPUT_SUFFIX)
            wb.save(output_path)
            return total_matches, output_path

        return 0, None

    except ImportError:
        logger.error("openpyxl not installed. Install: pip install openpyxl")
        return 0, None
    except (OSError, RuntimeError) as e:
        logger.error(f"Error processing XLSX {xlsx_path.name}: {e}")
        return 0, None


def search_in_pptx(pptx_path: Path, filenames: list[str],
                   fuzzy: bool = False,
                   pattern_set: set[str] | None = None) -> tuple[int, Path | None]:
    """
    Search PowerPoint presentation and highlight text.

    Args:
        pptx_path: Path to .pptx file
        filenames: List of filenames to search for
        fuzzy: If True, use fuzzy matching

    Returns:
        Tuple of (match_count, output_path or None)
    """
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor

        all_patterns = pattern_set if pattern_set is not None else build_pattern_set(filenames)
        red = RGBColor(255, 0, 0)

        prs = Presentation(pptx_path)
        total_matches = 0
        pattern_hit_counts: dict[str, int] = {}

        for slide_num, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                if fuzzy:
                    # Fuzzy: search the full shape text, then mark using window text
                    fuzzy_hits = fuzzy_text_search(shape.text, list(all_patterns))
                    if fuzzy_hits:
                        total_matches += len(fuzzy_hits)
                        for pat, _win in fuzzy_hits:
                            pattern_hit_counts[pat] = pattern_hit_counts.get(pat, 0) + 1
                        windows = {window.lower() for _pat, window in fuzzy_hits}
                        for paragraph in shape.text_frame.paragraphs:
                            for window_lower in windows:
                                _color_paragraph_runs(paragraph, window_lower, red)
                else:
                    # Exact: paragraph-level search handles split runs
                    for paragraph in shape.text_frame.paragraphs:
                        for pattern in all_patterns:
                            n = _color_paragraph_runs(paragraph, pattern.lower(), red)
                            if n:
                                total_matches += n
                                pattern_hit_counts[pattern.lower()] = pattern_hit_counts.get(pattern.lower(), 0) + n
                                logger.info(f"PPTX match: '{pattern}' ({n}×) on slide {slide_num}")

        _cache_match_counts(pptx_path, pattern_hit_counts)
        if total_matches > 0:
            output_path = _output_path(pptx_path, OUTPUT_SUFFIX)
            prs.save(output_path)
            return total_matches, output_path

        return 0, None

    except ImportError:
        logger.error("python-pptx not installed. Install: pip install python-pptx")
        return 0, None
    except (OSError, RuntimeError) as e:
        logger.error(f"Error processing PPTX {pptx_path.name}: {e}")
        return 0, None


def search_in_odt(odt_path: Path, filenames: list[str],
                  fuzzy: bool = False,
                  pattern_set: set[str] | None = None) -> tuple[int, Path | None]:
    """
    Search OpenDocument Text file for filenames.
    Note: odfpy does not support easy highlighting, so we extract text and report matches.
    Creates a copy with OUTPUT_SUFFIX if matches found.

    Args:
        odt_path: Path to .odt file
        filenames: List of filenames to search for
        fuzzy: If True, use fuzzy matching

    Returns:
        Tuple of (match_count, output_path or None)
    """
    try:
        from odf.opendocument import load
        from odf.text import P
        from odf import teletype

        all_patterns = pattern_set if pattern_set is not None else build_pattern_set(filenames)

        doc = load(str(odt_path))
        total_matches = 0
        pattern_hit_counts: dict[str, int] = {}

        paragraphs = doc.getElementsByType(P)
        if fuzzy:
            full_text = "\n".join(teletype.extractText(p) for p in paragraphs)
            fuzzy_hits = fuzzy_text_search(full_text, list(all_patterns))
            total_matches = len(fuzzy_hits)
            for pat, _win in fuzzy_hits:
                pattern_hit_counts[pat] = pattern_hit_counts.get(pat, 0) + 1
        else:
            for para in paragraphs:
                text = teletype.extractText(para).lower()
                for pattern in all_patterns:
                    if pattern.lower() in text:
                        total_matches += 1
                        pattern_hit_counts[pattern.lower()] = pattern_hit_counts.get(pattern.lower(), 0) + 1
                        logger.info(f"ODT match: '{pattern}' in {odt_path.name}")

        _cache_match_counts(odt_path, pattern_hit_counts)
        if total_matches > 0:
            output_path = _output_path(odt_path, OUTPUT_SUFFIX)
            shutil.copy2(odt_path, output_path)
            return total_matches, output_path

        return 0, None

    except ImportError:
        logger.error("odfpy not installed. Install: pip install odfpy")
        return 0, None
    except (OSError, RuntimeError) as e:
        logger.error(f"Error processing ODT {odt_path.name}: {e}")
        return 0, None


def get_txt_files(directory: Path = Path(".")) -> list[Path]:
    """
    Get all TXT files in directory (excluding *_checked.txt files).

    Args:
        directory: Directory to search

    Returns:
        List of Path objects for TXT files (sorted)
    """
    try:
        txt_files = [
            f for f in directory.iterdir()
            if f.is_file()
            and f.suffix.lower() == '.txt'
            and not _is_output_file(f.stem.lower())
        ]
        logger.info(f"Found {len(txt_files)} TXT files (excluding output files)")
        return sorted(txt_files)
    except PermissionError as e:
        logger.error(f"Permission denied scanning {directory}: {e}")
        return []
    except OSError as e:
        logger.error(f"Failed to get TXT files from {directory}: {e}")
        return []


def search_in_txt(txt_path: Path, filenames: list[str],
                  marker: str = TXT_MARKER,
                  position: str = TXT_MARKER_POSITION,
                  fuzzy: bool = False,
                  pattern_set: set[str] | None = None) -> tuple[int, Path | None]:
    """
    Search text file and add markers.

    Args:
        txt_path: Path to .txt file
        filenames: List of filenames to search for
        marker: Marker text to add (default: "[[FOUND]]")
        position: "before" or "after" the match
        fuzzy: If True, use fuzzy matching

    Returns:
        Tuple of (match_count, output_path or None)
    """
    try:
        all_patterns = pattern_set if pattern_set is not None else build_pattern_set(filenames)

        # Read file
        with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()

        total_matches = 0
        modified_content = content
        pattern_hit_counts: dict[str, int] = {}

        if fuzzy:
            fuzzy_hits = fuzzy_text_search(content, list(all_patterns))
            total_matches = len(fuzzy_hits)
            for pat, _win in fuzzy_hits:
                pattern_hit_counts[pat] = pattern_hit_counts.get(pat, 0) + 1
            # Use the actual matched window text for marker placement, not the pattern.
            # This ensures markers are placed at the real match location even when the
            # fuzzy window differs from the pattern (e.g. "report 2024" vs "report_2024").
            content_lower = content.lower()
            marked_positions: set[int] = set()
            for _pattern, window in fuzzy_hits:
                window_lower = window.lower()
                pos = 0
                while True:
                    idx = content_lower.find(window_lower, pos)
                    if idx == -1:
                        break
                    if idx not in marked_positions:
                        matched_text = modified_content[idx:idx + len(window)]
                        if position == "before":
                            replacement = f"{marker}{matched_text}"
                        else:
                            replacement = f"{matched_text}{marker}"
                        modified_content = (modified_content[:idx]
                                            + replacement
                                            + modified_content[idx + len(window):])
                        content_lower = modified_content.lower()
                        marked_positions.add(idx)
                        pos = idx + len(replacement)
                    else:
                        pos = idx + 1
        else:
            # Search and mark each pattern.
            # Track already-marked source positions to prevent overlapping patterns
            # (e.g. "testfile", "testfile_a", "testfile_a.txt") from inserting multiple
            # markers at the same location.
            marked_ranges: list[tuple[int, int]] = []

            def _overlaps(start: int, end: int) -> bool:
                """Return True if the span [start, end) overlaps any already-marked range.

                Args:
                    start: Inclusive start index of the span to check.
                    end: Exclusive end index of the span to check.
                """
                return any(s <= start < e or s < end <= e for s, e in marked_ranges)

            for pattern in sorted(all_patterns, key=len, reverse=True):  # longest first
                # Word-boundary guard: short bare-word patterns (no dot/hyphen/space,
                # length < 8) must appear as complete tokens, not substrings of longer words.
                if _is_short_bare(pattern):
                    regex = re.compile(r'(?<!\w)' + re.escape(pattern) + r'(?!\w)', re.IGNORECASE)
                else:
                    regex = re.compile(re.escape(pattern), re.IGNORECASE)
                raw_matches = list(regex.finditer(content))  # search in original content

                hit_count = 0
                for raw_match in raw_matches:
                    orig_start, orig_end = raw_match.span()
                    if _overlaps(orig_start, orig_end):
                        continue  # skip: a longer pattern already marked this span
                    marked_ranges.append((orig_start, orig_end))
                    hit_count += 1

                if hit_count:
                    total_matches += hit_count
                    pattern_hit_counts[pattern.lower()] = pattern_hit_counts.get(pattern.lower(), 0) + hit_count
                    logger.info(f"TXT match: '{pattern}' found {hit_count} times in {txt_path.name}")

            # Apply markers in a second pass over original content using marked_ranges
            marked_ranges.sort(key=lambda x: x[0])
            result_parts: list[str] = []
            prev = 0
            for orig_start, orig_end in marked_ranges:
                result_parts.append(content[prev:orig_start])
                matched_text = content[orig_start:orig_end]
                if position == "before":
                    result_parts.append(f"{marker}{matched_text}")
                else:
                    result_parts.append(f"{matched_text}{marker}")
                prev = orig_end
            result_parts.append(content[prev:])
            modified_content = "".join(result_parts)

        # Save if matches found
        _cache_match_counts(txt_path, pattern_hit_counts)
        if total_matches > 0:
            output_path = _output_path(txt_path, OUTPUT_SUFFIX)
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(modified_content)
            return total_matches, output_path

        return 0, None

    except (OSError, UnicodeDecodeError, RuntimeError) as e:
        logger.error(f"Error processing TXT {txt_path.name}: {e}")
        return 0, None


# ============================================================================
# REPORT & HASHLIST UTILITIES
# ============================================================================


def generate_hashlist(hash_type: str) -> int:
    """Generate a text file with one hash value per line.

    Args:
        hash_type: "md5" or "sha256"

    Returns:
        Number of entries written
    """
    data = load_json_data()
    output_file = f"filenames_{hash_type}.txt"
    count = 0
    with open(output_file, "w", encoding="utf-8") as f:
        for entry in data:
            hash_val = getattr(entry, hash_type, "")
            if hash_val and hash_val != "N/A":
                f.write(hash_val + "\n")
                count += 1
    logger.info(f"Hash list generated: {output_file} with {count} entries")
    return count


def _extract_searchable_text_uncached(file_path: Path) -> str:
    """Extract text from *file_path* without consulting the disk cache.

    Args:
        file_path: Path to the source file (PDF, JPG, DOCX, XLSX, PPTX, ODT, or TXT).

    Returns:
        Normalized extracted text, or an empty string on failure or unsupported type.
    """
    ext = file_path.suffix.lower()
    try:
        if ext == '.pdf':
            doc = fitz.open(file_path)
            if doc.is_encrypted:
                doc.close()
                return ""
            text = "\n".join(page.get_text("text") for page in doc)
            doc.close()
            if text.strip():
                return _normalize_text(text)
            if TESSERACT_AVAILABLE:
                import pytesseract
                from PIL import Image
                doc = fitz.open(file_path)
                parts = []
                for page in doc:
                    pix = page.get_pixmap(matrix=fitz.Matrix(OCR_DPI / 72, OCR_DPI / 72))
                    img = Image.open(io.BytesIO(pix.tobytes("png")))
                    parts.append(pytesseract.image_to_string(
                        img, lang=OCR_LANGUAGES, config=f"--psm {OCR_PSM} --oem 3"))
                    img.close()
                doc.close()
                return _normalize_text("\n".join(parts))
            return ""

        if ext in ('.jpg', '.jpeg'):
            from PIL import Image
            parts = []
            img = Image.open(file_path)
            exif = img.getexif()
            if exif:
                parts.append(' '.join(str(v) for v in exif.values() if v))
            if TESSERACT_AVAILABLE:
                import pytesseract
                parts.append(pytesseract.image_to_string(
                    img, lang=OCR_LANGUAGES, config=f"--psm {OCR_PSM} --oem 3"))
            img.close()
            return _normalize_text(" ".join(parts))

        if ext == '.docx':
            from docx import Document
            doc = Document(file_path)
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            return _normalize_text("\n".join(parts))

        if ext == '.xlsx':
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True)
            parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value:
                            parts.append(str(cell.value))
            wb.close()
            return _normalize_text("\n".join(parts))

        if ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        parts.append(shape.text)
            return _normalize_text("\n".join(parts))

        if ext in ('.odt', '.ods', '.odp'):
            from odf.opendocument import load
            from odf.text import P
            from odf import teletype
            doc = load(str(file_path))
            return _normalize_text("\n".join(teletype.extractText(p)
                                             for p in doc.getElementsByType(P)))

        if ext == '.txt':
            return _normalize_text(file_path.read_text(encoding='utf-8', errors='ignore'))

    except (OSError, ImportError, RuntimeError, UnicodeDecodeError) as e:
        logger.error(f"Text extraction failed for {file_path}: {e}")
    return ""


def extract_searchable_text(file_path: Path) -> str:
    """Extract all searchable text from a file for report match counting.

    Checks the disk cache first (keyed by file SHA-256) to avoid re-running
    Tesseract OCR on unchanged files.  The returned string is Unicode-
    normalized (NFC) and has common OCR ligatures resolved.

    Args:
        file_path: Path to the source file to extract text from.

    Returns:
        Normalized extracted text, or an empty string if extraction fails.
    """
    cached = _load_disk_cache(file_path)
    if cached is not None:
        logger.debug(f"Disk cache hit: {file_path.name}")
        return cached
    result = _extract_searchable_text_uncached(file_path)
    if result:
        _save_disk_cache(file_path, result)
    return result


def count_matches_per_pattern(
    text: str, filenames: list[str]
) -> dict[str, dict[str, int]]:
    """Count occurrences of each search pattern per filename in text.

    For each filename, generates search patterns via generate_search_patterns()
    and counts how many times each pattern appears (case-insensitive).

    Each pattern is also searched in a separator-normalized copy of the text
    (hyphens, underscores, and dots replaced with spaces).  The higher of the
    two counts is used, catching cases where the document uses different
    separators than the pattern (e.g. pattern ``"report-2024"`` but the text
    contains ``"report 2024"``, or vice versa).

    Args:
        text: The full extracted text to search in
        filenames: List of original filenames

    Returns:
        Nested dict: {filename: {pattern: count, ...}, ...}
    """
    text_lower = _normalize_text(text).lower()
    text_norm = _normalize_sep(text_lower)   # separator-normalized copy for fallback counting
    result: dict[str, dict[str, int]] = {}
    for filename in filenames:
        patterns = generate_search_patterns(filename)
        pattern_counts: dict[str, int] = {}
        for pattern in patterns:
            p_lower = _normalize_text(pattern).lower()
            if _is_short_bare(p_lower):
                count_raw = len(re.findall(r'(?<!\w)' + re.escape(p_lower) + r'(?!\w)', text_lower))
                count_norm = 0  # boundary guard takes precedence; skip sep-normalized fallback
            else:
                count_raw = text_lower.count(p_lower)
                p_norm = _normalize_sep(p_lower)
                # Only search the separator-normalized copy when the form actually differs
                count_norm = text_norm.count(p_norm) if p_norm != p_lower else 0
            pattern_counts[pattern] = max(count_raw, count_norm)
        result[filename] = pattern_counts
    return result


def count_regex_matches_per_pattern(
    text: str, regex_patterns: list[re.Pattern]
) -> dict[str, dict[str, int]]:
    """Count regex match occurrences in text for report generation.

    Used in regex search mode so that the PDF report reflects actual regex
    hits rather than filename-pattern counts.

    Args:
        text: The full extracted text to search in
        regex_patterns: Compiled regex patterns from the user

    Returns:
        Nested dict keyed by regex pattern string: {pattern: {pattern: count}}
    """
    result: dict[str, dict[str, int]] = {}
    for rx in regex_patterns:
        count = len(rx.findall(text))
        result[rx.pattern] = {rx.pattern: count}
    return result


def _get_match_context(text: str, pattern: str, context_chars: int = 70) -> str | None:
    """Return a short text snippet surrounding the first occurrence of pattern.

    Provides the reader with a quick preview of where a match appears in the
    document.  Returns ``None`` when the pattern is not found in the text.

    Args:
        text: Full extracted text to search in (original case)
        pattern: Pattern to locate (case-insensitive)
        context_chars: Number of characters to show on each side of the match

    Returns:
        A single-line snippet string, or None if the pattern is not found.
    """
    idx = text.lower().find(pattern.lower())
    if idx == -1:
        return None
    start = max(0, idx - context_chars // 2)
    end = min(len(text), idx + len(pattern) + context_chars // 2)
    snippet = text[start:end].replace('\n', ' ').replace('\r', ' ')
    # Collapse multiple spaces that result from newline replacement
    snippet = re.sub(r' {2,}', ' ', snippet).strip()
    prefix = '…' if start > 0 else ''
    suffix = '…' if end < len(text) else ''
    return f"{prefix}{snippet}{suffix}"


def _wrap_text_lines(
    text: str,
    fontname: str,
    fontsize: float,
    max_width: float,
) -> list[str]:
    """Split *text* into lines that each fit within *max_width* points.

    Words are preserved where possible.  Single words that exceed *max_width*
    on their own are broken at the character level as a last resort.

    Args:
        text: Input string to wrap.
        fontname: PyMuPDF font name used for width measurement (e.g. ``"helv"``).
        fontsize: Font size in points.
        max_width: Maximum line width in points.

    Returns:
        List of line strings; always contains at least one element.
    """
    if not text:
        return [""]
    if fitz.get_text_length(text, fontname=fontname, fontsize=fontsize) <= max_width:
        return [text]
    lines: list[str] = []
    current = ""
    for word in text.split():
        candidate = f"{current} {word}".lstrip() if current else word
        if fitz.get_text_length(candidate, fontname=fontname, fontsize=fontsize) <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            if fitz.get_text_length(word, fontname=fontname, fontsize=fontsize) > max_width:
                # Hard-break oversized single word at character level
                partial = ""
                for ch in word:
                    if fitz.get_text_length(partial + ch, fontname=fontname, fontsize=fontsize) <= max_width:
                        partial += ch
                    else:
                        if partial:
                            lines.append(partial)
                        partial = ch
                current = partial
            else:
                current = word
    if current:
        lines.append(current)
    return lines or [""]


def _insert_text_wrapped(
    page: "fitz.Page",
    x: float,
    y: float,
    text: str,
    fontsize: float,
    fontname: str,
    color: tuple,
    max_x: float = 523.0,
    line_height: float | None = None,
) -> float:
    """Insert text with automatic word-wrap and return the y after the last line.

    Combines :func:`_wrap_text_lines` with :meth:`fitz.Page.insert_text` to
    produce multi-line output whenever the text would otherwise overflow the
    right margin.  Suitable for filenames, file paths, and any other field
    whose length is not bounded.

    Args:
        page: Target PDF page.
        x: Left x coordinate (baseline anchor for the first line).
        y: Baseline y coordinate of the first line.
        text: Text to render.
        fontsize: Font size in points.
        fontname: PyMuPDF font name (e.g. ``"helv"``, ``"hebo"``).
        color: RGB color tuple with float values 0.0–1.0.
        max_x: Right boundary in points; defaults to the standard report right
            margin (523).
        line_height: Vertical advance per line in points; defaults to
            ``fontsize * 1.35``.

    Returns:
        y coordinate of the baseline one line below the last rendered line,
        ready for the next element.
    """
    if line_height is None:
        line_height = fontsize * 1.35
    for line in _wrap_text_lines(text, fontname, fontsize, max_x - x):
        page.insert_text((x, y), line, fontsize=fontsize, fontname=fontname, color=color)
        y += line_height
    return y


def generate_search_report(
    searched_file: str,
    json_data: list[FileEntry],
    pattern_matches: dict[str, dict[str, int]],
    output_dir: Path = Path("."),
    extra_words: list[str] | None = None,
    source_text: str | None = None,
    search_mode: str = "exact",
    stems_min_length: int = 0,
    ocr_used: bool = False,
) -> Path:
    """Generate a PDF report with per-pattern match breakdown for each filename.

    For each JSON entry, the report shows:
    - Filename, sha256, md5
    - Each search pattern generated from that filename with its match count
    - A text snippet around the first occurrence (when source_text is provided)
    - Total matches across all patterns

    If extra_words are provided, they appear in a separate "Manual Search Words"
    section after the JSON entries.

    Args:
        searched_file: Name of the file that was searched
        json_data: List of FileEntry objects
        pattern_matches: Nested dict {filename: {pattern: count, ...}, ...}
        output_dir: Directory to write the report to
        extra_words: Optional user-defined extra search terms

    Returns:
        Path to the generated report PDF
    """
    report_name = f"{Path(searched_file).stem}{REPORT_SUFFIX}.pdf"
    if OUTPUT_DIR:
        report_dir = Path(OUTPUT_DIR)
        report_dir.mkdir(parents=True, exist_ok=True)
        report_path = report_dir / report_name
    else:
        report_path = output_dir / report_name

    doc = fitz.open()
    page = doc.new_page()
    y = 72

    # ── Title ────────────────────────────────────────────────────────────────
    y = _insert_text_wrapped(page, 72, y, f"Search Report: {searched_file}",
                              14, "helv", (0, 0, 0), line_height=18.0)
    y += 6
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0, 0, 0), width=0.5)
    y += 14

    # ── NOTE (disclaimer with bold prefix) ───────────────────────────────────
    note_html = (
        "<b>NOTE:</b> This report, as well as the keyword search, has been "
        "generated automatically. It cannot be guaranteed that all keywords "
        "have been accurately identified. The results are intended primarily "
        "as a preliminary assessment and reference point, requiring further "
        "verification. All pattern matching is case-insensitive. "
        "Depending on the search method applied, false positive matches may occur: "
        "exact and fuzzy pattern matching may match substrings in unrelated "
        "contexts; fuzzy matching intentionally accepts near-matches up to the "
        "configured similarity threshold. In OCR-based searches, character "
        "recognition errors may produce both false positives and false negatives. "
        "These occurrences are not errors but inherent technical characteristics "
        "of the respective search method. "
        "This report does not claim to be exhaustive; matches may be absent if "
        "text is embedded in images without OCR, stored in unsupported formats, "
        "or located in password-protected or corrupted content. "
        "The original source file has not been modified. "
        "All results must be reviewed manually."
    )
    note_rect = fitz.Rect(72, y, 523, y + 130)
    page.insert_htmlbox(note_rect, note_html,
                        css="body { font-size: 7pt; color: #666666; }")
    y += 134

    # ── Two blank lines ───────────────────────────────────────────────────────
    y += 18

    # ── Context paragraph ─────────────────────────────────────────────────────
    _now = datetime.now()
    _date_str = _now.strftime("%d.%m.%Y")
    _time_str = _now.strftime("%H:%M:%S")

    _mode_labels = {
        "exact": "Exact matching",
        "fuzzy": "Fuzzy matching",
        "stems": "Stems only",
        "regex": "Regex matching",
    }
    _mode_label = _mode_labels.get(search_mode, search_mode.capitalize())

    if search_mode == "regex" and extra_words:
        _words_phrase = (
            f"The following regex patterns were used: "
            f"<b>{', '.join(extra_words)}</b>."
        )
    elif extra_words:
        _words_phrase = (
            f"The following extra words were added manually: "
            f"<b>{', '.join(extra_words)}</b>."
        )
    else:
        _words_phrase = "No extra words were added manually."

    if search_mode == "stems":
        _stem_phrase = (
            f" STEMS_MIN_LENGTH was set to <b>{stems_min_length}</b> "
            f"character{'s' if stems_min_length != 1 else ''}; "
            f"MIN_STEM_LENGTH (internal pattern minimum) is "
            f"<b>{MIN_STEM_LENGTH}</b> characters."
        )
    else:
        _stem_phrase = ""

    if ocr_used:
        try:
            import pytesseract as _pt
            _tess_ver = str(_pt.get_tesseract_version())
        except (ImportError, OSError, RuntimeError):
            _tess_ver = "unknown"
        _ocr_phrase = (
            f" OCR was performed using Tesseract <b>{_tess_ver}</b>"
            f" (languages: <b>{OCR_LANGUAGES}</b>,"
            f" PSM: <b>{OCR_PSM}</b>,"
            f" DPI: <b>{OCR_DPI}</b>)."
        )
    else:
        _ocr_phrase = ""

    context_html = (
        f"This report was generated on <b>{_date_str}</b> at "
        f"<b>{_time_str}</b>. "
        f"The search mode used was <b>{_mode_label}</b>. "
        f"{_words_phrase}"
        f"{_stem_phrase}"
        f"{_ocr_phrase} "
        f"For additional details see <b>txtfinder.log</b>."
    )
    context_rect = fitz.Rect(72, y, 523, y + 70)
    page.insert_htmlbox(context_rect, context_html,
                        css="body { font-size: 7pt; color: #444444; }")
    y += 74

    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0.8, 0.8, 0.8), width=0.3)
    y += 14

    # ── Summary section ──────────────────────────────────────────────────────
    # Build a flat list of (filename, total_matches) for all entries with hits.
    # extra_words are included under their own keyword label.
    # summary_hits: (filename, filepath, size_str, total)
    summary_hits: list[tuple[str, str, str, int]] = []
    for entry in json_data:
        total = sum(pattern_matches.get(entry.filename, {}).values())
        if total > 0:
            if entry.size:
                _sz = (f"{entry.size / 1_048_576:.2f} MB"
                       if entry.size >= 1_048_576 else f"{entry.size / 1024:.1f} KB")
            else:
                _sz = ""
            summary_hits.append((entry.filename, entry.filepath or "", _sz, total))
    if extra_words:
        for word in extra_words:
            total = sum(pattern_matches.get(word, {}).values())
            if total > 0:
                summary_hits.append((word, "", "", total))

    page.insert_text((72, y), "Summary — Matched Files",
                      fontsize=11, fontname="hebo", color=(0, 0, 0))
    y += 16

    if summary_hits:
        for name, filepath, size_str, total in summary_hits:
            _name_lines = _wrap_text_lines(f"  {name}", "helv", 9, 380 - 72)
            _fp_lines = _wrap_text_lines(filepath, "helv", 7, 523 - 84) if filepath else []
            needed = (len(_name_lines) * 13
                      + len(_fp_lines) * 10
                      + (10 if size_str else 0))
            if y + needed > 750:
                page = doc.new_page()
                y = 72
            count_str = f"{total} match{'es' if total != 1 else ''}"
            y_start = y
            y = _insert_text_wrapped(page, 72, y, f"  {name}",
                                      9, "helv", (0, 0, 0),
                                      max_x=380.0, line_height=13.0)
            page.insert_text((390, y_start), count_str,
                              fontsize=9, fontname="hebo", color=(0.8, 0, 0))
            if filepath:
                y = _insert_text_wrapped(page, 84, y, filepath,
                                          7, "helv", (0.5, 0.5, 0.5),
                                          max_x=523.0, line_height=10.0)
            if size_str:
                page.insert_text((84, y), size_str,
                                  fontsize=7, fontname="helv", color=(0.5, 0.5, 0.5))
                y += 10
    else:
        page.insert_text((72, y), "  No matches found.",
                          fontsize=9, fontname="helv", color=(0.5, 0.5, 0.5))
        y += 13

    y += 6
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0, 0, 0), width=0.5)
    y += 16

    # ── Detail section ───────────────────────────────────────────────────────
    page.insert_text((72, y), "Details — Pattern Breakdown",
                      fontsize=11, fontname="hebo", color=(0, 0, 0))
    y += 18

    for entry in json_data:
        patterns = pattern_matches.get(entry.filename, {})
        # Calculate space needed: account for potential wrapping of long names/paths
        _fn_lines = len(_wrap_text_lines(f"Filename: {entry.filename}", "helv", 10, 451))
        _fp_lines = len(_wrap_text_lines(entry.filepath or "", "helv", 7, 451)) if entry.filepath else 0
        _extra_wrap = max(0, _fn_lines - 1) * 13 + max(0, _fp_lines - 1) * 10
        _meta_lines = _fp_lines + (1 if entry.size else 0)
        needed = _fn_lines * 13 + _meta_lines * 10 + 23 + len(patterns) * 11 + 18 + 8 + _extra_wrap
        if y + needed > 750:
            page = doc.new_page()
            y = 72

        # Filename header
        y = _insert_text_wrapped(page, 72, y, f"Filename: {entry.filename}",
                                  10, "helv", (0, 0, 0), line_height=13.0)
        if entry.sha256 != "N/A":
            if entry.filepath:
                y = _insert_text_wrapped(page, 72, y, f"filepath: {entry.filepath}",
                                          7, "helv", (0.4, 0.4, 0.4), line_height=10.0)
            if entry.size:
                _sz = (f"{entry.size / 1_048_576:.2f} MB"
                       if entry.size >= 1_048_576 else f"{entry.size / 1024:.1f} KB")
                page.insert_text((72, y), f"size: {_sz}",
                                  fontsize=7, fontname="helv", color=(0.4, 0.4, 0.4))
                y += 10
            page.insert_text((72, y), f"sha256: {entry.sha256}",
                              fontsize=7, fontname="helv", color=(0.4, 0.4, 0.4))
            y += 10
            page.insert_text((72, y), f"md5: {entry.md5}",
                              fontsize=7, fontname="helv", color=(0.4, 0.4, 0.4))
            y += 13
        else:
            y += 4  # small gap before patterns in stems-only mode

        # Pattern breakdown
        total = 0
        for pattern, count in patterns.items():
            total += count
            c_color = (0.8, 0, 0) if count > 0 else (0.6, 0.6, 0.6)
            y_start = y
            y = _insert_text_wrapped(page, 88, y, f'"{pattern}"',
                                      7, "helv", (0.3, 0.3, 0.3),
                                      max_x=330.0, line_height=11.0)
            page.insert_text((340, y_start), f"{count} match{'es' if count != 1 else ''}",
                              fontsize=7, fontname="helv", color=c_color)
            # Show a text excerpt around the first hit for quick context
            if count > 0 and source_text:
                snippet = _get_match_context(source_text, pattern)
                if snippet:
                    if y + 10 > 750:
                        page = doc.new_page()
                        y = 72
                    rect = fitz.Rect(96, y, 523, y + 10)
                    page.insert_textbox(rect, snippet, fontsize=6, fontname="helv",
                                        color=(0.35, 0.35, 0.35))
                    y += 10

        # Total line
        t_color = (0.8, 0, 0) if total > 0 else (0.5, 0.5, 0.5)
        page.insert_text((72, y), f"Total: {total} match{'es' if total != 1 else ''}",
                          fontsize=9, fontname="hebo", color=t_color)
        y += 18
        # Separator line
        page.draw_line(fitz.Point(72, y - 4), fitz.Point(420, y - 4),
                        color=(0.9, 0.9, 0.9), width=0.2)

    # Extra manual search words section
    if extra_words:
        # Section header
        if y + 30 > 750:
            page = doc.new_page()
            y = 72
        y += 10
        page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                        color=(0, 0, 0), width=0.5)
        y += 16
        page.insert_text((72, y), "Manual Search Words",
                          fontsize=11, fontname="hebo", color=(0, 0, 0))
        y += 18

        for word in extra_words:
            patterns = pattern_matches.get(word, {})
            needed = 16 + len(patterns) * 11 + 16 + 8
            if y + needed > 750:
                page = doc.new_page()
                y = 72

            y = _insert_text_wrapped(page, 72, y, f"Keyword: {word}",
                                      10, "helv", (0, 0, 0), line_height=13.0)

            total = 0
            for pattern, count in patterns.items():
                total += count
                c_color = (0.8, 0, 0) if count > 0 else (0.6, 0.6, 0.6)
                y_start = y
                y = _insert_text_wrapped(page, 88, y, f'"{pattern}"',
                                          7, "helv", (0.3, 0.3, 0.3),
                                          max_x=330.0, line_height=11.0)
                page.insert_text((340, y_start), f"{count} match{'es' if count != 1 else ''}",
                                  fontsize=7, fontname="helv", color=c_color)

            t_color = (0.8, 0, 0) if total > 0 else (0.5, 0.5, 0.5)
            page.insert_text((72, y), f"Total: {total} match{'es' if total != 1 else ''}",
                              fontsize=9, fontname="hebo", color=t_color)
            y += 18
            page.draw_line(fitz.Point(72, y - 4), fitz.Point(420, y - 4),
                            color=(0.9, 0.9, 0.9), width=0.2)

    doc.set_metadata({
        "author": "ot2i7ba",
        "subject": "github.com/ot2i7ba",
    })
    doc.save(report_path)
    doc.close()
    logger.info(f"Report generated: {report_path}")
    return report_path


def _generate_report_for_file(
    file_path: Path,
    json_data: list[FileEntry] | None = None,
    extracted_text: str | None = None,
    extra_words: list[str] | None = None,
    regex_patterns: list[re.Pattern] | None = None,
    stems_mode: bool = False,
    search_mode: str = "exact",
    stems_min_length: int = 0,
    ocr_used: bool = False,
) -> Path | None:
    """Generate a search report for a single searched file.

    Args:
        file_path: The file that was searched
        json_data: Pre-loaded JSON data (avoids reloading per file in batch)
        extracted_text: Pre-extracted text (avoids double OCR extraction)
        extra_words: Optional user-defined extra search terms
        regex_patterns: If provided, report shows regex hit counts instead of
            filename-pattern counts (used in regex search mode)
        stems_mode: If True, report only the bare stem of each filename (no
            pattern variants). True when the user selected "stems" search mode
            OR when the JSON was generated in stems-only mode (sha256="N/A").

    Returns:
        Report path or None on error.
    """
    try:
        if json_data is None:
            json_data = load_json_data()
        text = extracted_text if extracted_text is not None else extract_searchable_text(file_path)

        if regex_patterns:
            # Regex mode: report reflects actual regex hits, not filename patterns
            pattern_matches = count_regex_matches_per_pattern(text, regex_patterns)
            # Pass an empty json_data list so the report only shows the regex section
            return generate_search_report(file_path.name, [], pattern_matches,
                                          file_path.parent,
                                          extra_words=[rx.pattern for rx in regex_patterns],
                                          source_text=text,
                                          search_mode="regex",
                                          stems_min_length=stems_min_length,
                                          ocr_used=ocr_used)

        # stems_only is True when explicitly requested OR when the JSON was built
        # without hashes (sha256 = "N/A"), which is the stems-only JSON mode.
        stems_only = stems_mode or (bool(json_data) and all(e.sha256 == "N/A" for e in json_data))

        cached_counts = _pop_cached_match_counts(file_path)

        if stems_only:
            # Stems mode: each filename is represented by exactly one pattern —
            # its bare stem (no extension, no separator variants).  The search
            # already used build_stem_pattern_set so cached_counts keys are stems.
            pattern_matches = {}
            for entry in json_data:
                stem = Path(entry.filename).stem or entry.filename
                if cached_counts is not None:
                    count = cached_counts.get(stem.lower(), 0)
                else:
                    count = text.lower().count(stem.lower())
                pattern_matches[entry.filename] = {stem: count}
            if extra_words:
                for word in extra_words:
                    stem_w = Path(word).stem or word
                    if cached_counts is not None:
                        count = cached_counts.get(stem_w.lower(), 0)
                    else:
                        count = text.lower().count(stem_w.lower())
                    pattern_matches[word] = {stem_w: count}
        elif cached_counts is not None:
            # Normal mode with cached counts: use actual hits from the search run.
            pattern_matches = _build_pattern_matches_from_counts(
                json_data, cached_counts, extra_words
            )
        else:
            # Fallback: recount from extracted text (e.g. when report is called
            # standalone without a prior search run having populated the cache).
            all_filenames = [e.filename for e in json_data]
            if extra_words:
                all_filenames = all_filenames + extra_words
            pattern_matches = count_matches_per_pattern(text, all_filenames)

        return generate_search_report(file_path.name, json_data, pattern_matches,
                                      file_path.parent, extra_words=extra_words,
                                      source_text=text,
                                      search_mode=search_mode,
                                      stems_min_length=stems_min_length,
                                      ocr_used=ocr_used)
    except (OSError, RuntimeError, ValueError) as e:
        logger.error(f"Report generation failed for {file_path.name}: {e}")
        return None


def generate_csv_summary(
    results: list[FileSearchResult],
    output_name: str | None = None,
) -> Path:
    """Generate a CSV summary of batch search results.

    Args:
        results: List of FileSearchResult from batch processing
        output_name: Output CSV filename. Defaults to
            ``search_results_YYYYMMDDHHMMSS.csv`` (timestamped to prevent
            overwriting results from previous search runs).

    Returns:
        Path to the generated CSV file
    """
    if output_name is None:
        output_name = f"search_results_{datetime.now().strftime('%Y%m%d%H%M%S')}.csv"
    output_path = Path(OUTPUT_DIR) / output_name if OUTPUT_DIR else Path(output_name)
    with open(output_path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f, delimiter=";")
        writer.writerow(["searched_file", "matches", "output_file", "report_file"])
        for r in results:
            writer.writerow([
                r.file.name,
                r.matches,
                r.output.name if r.output else "",
                r.report.name if r.report else "",
            ])
    logger.info(f"CSV summary: {output_path}")
    return output_path


# ============================================================================
# HASH LOOKUP
# ============================================================================

def _classify_hash(hash_str: str) -> str | None:
    """Classify a hash string by length and character set.

    Args:
        hash_str: The candidate hash value to inspect (whitespace is stripped).

    Returns:
        ``"MD5"`` for 32-char hex strings, ``"SHA256"`` for 64-char hex strings,
        or ``None`` if the input is not a valid hash.
    """
    h = hash_str.strip().lower()
    if not h:
        return None
    if not all(c in "0123456789abcdef" for c in h):
        return None
    if len(h) == 32:
        return "MD5"
    if len(h) == 64:
        return "SHA256"
    return None


def generate_hash_search_report(
    queried_hashes: list[tuple[str, str]],
    json_data: list[FileEntry],
    output_dir: Path = Path("."),
) -> Path:
    """Generate a PDF report for a hash lookup against filenames.json.

    Args:
        queried_hashes: List of (hash_value_lower, hash_type) tuples
        json_data: Entries loaded from filenames.json
        output_dir: Directory to write the report (overridden by OUTPUT_DIR)

    Returns:
        Path to the generated report PDF
    """
    now = datetime.now()
    report_name = f"{now.strftime('%Y%m%d%H%M%S')}{HASHES_SUFFIX}.pdf"
    if OUTPUT_DIR:
        report_dir = Path(OUTPUT_DIR)
        report_dir.mkdir(parents=True, exist_ok=True)
        report_path = report_dir / report_name
    else:
        report_path = output_dir / report_name

    # Build lookup indexes (case-insensitive)
    md5_index: dict[str, list[FileEntry]] = {}
    sha256_index: dict[str, list[FileEntry]] = {}
    for entry in json_data:
        if entry.md5 and entry.md5 != "N/A":
            md5_index.setdefault(entry.md5.lower(), []).append(entry)
        if entry.sha256 and entry.sha256 != "N/A":
            sha256_index.setdefault(entry.sha256.lower(), []).append(entry)

    doc = fitz.open()
    page = doc.new_page()
    y = 72

    # ── Title ────────────────────────────────────────────────────────────────
    page.insert_text((72, y), "Hash Lookup Report",
                      fontsize=14, fontname="helv", color=(0, 0, 0))
    y += 24
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0, 0, 0), width=0.5)
    y += 14

    # ── NOTE (disclaimer) ─────────────────────────────────────────────────────
    note_html = (
        "<b>NOTE:</b> This report, as well as the hash lookup, has been "
        "generated automatically. It cannot be guaranteed that all hash values "
        "have been accurately identified. The results are intended primarily "
        "as a preliminary assessment and reference point, requiring further "
        "verification. All hash comparisons are case-insensitive."
    )
    note_rect = fitz.Rect(72, y, 523, y + 52)
    page.insert_htmlbox(note_rect, note_html,
                        css="body { font-size: 7pt; color: #666666; }")
    y += 56

    # ── Two blank lines ───────────────────────────────────────────────────────
    y += 18

    # ── Context paragraph ─────────────────────────────────────────────────────
    context_html = (
        f"This report was generated on <b>{now.strftime('%d.%m.%Y')}</b> at "
        f"<b>{now.strftime('%H:%M:%S')}</b>. "
        f"Hash lookup performed against <b>{FILENAMES_JSON}</b> "
        f"({len(json_data):,} entries). "
        f"For additional details see <b>txtfinder.log</b>."
    )
    ctx_rect = fitz.Rect(72, y, 523, y + 40)
    page.insert_htmlbox(ctx_rect, context_html,
                        css="body { font-size: 7pt; color: #444444; }")
    y += 44
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0.8, 0.8, 0.8), width=0.3)
    y += 14

    # ── Per-hash results ──────────────────────────────────────────────────────
    found_total = 0
    not_found_total = 0

    for hash_val, hash_type in queried_hashes:
        if y + 44 > 750:
            page = doc.new_page()
            y = 72

        # Hash header line
        page.insert_text((72, y), f"[{hash_type}]",
                          fontsize=8, fontname="hebo", color=(0.2, 0.2, 0.6))
        page.insert_text((115, y), hash_val,
                          fontsize=8, fontname="helv", color=(0, 0, 0))
        y += 13

        matches = (md5_index.get(hash_val, []) if hash_type == "MD5"
                   else sha256_index.get(hash_val, []))

        if not matches:
            not_found_total += 1
            page.insert_text((88, y), "Not found in filenames.json",
                              fontsize=8, fontname="helv", color=(0.7, 0.0, 0.0))
            y += 12
        else:
            found_total += 1
            for entry in matches:
                if entry.size:
                    size_str = (f"{entry.size / 1_048_576:.2f} MB"
                                if entry.size >= 1_048_576
                                else f"{entry.size / 1024:.1f} KB")
                else:
                    size_str = "unknown"
                fields = [
                    ("Filename", entry.filename),
                    ("Filepath", entry.filepath or "—"),
                    ("Size",     size_str),
                    ("MD5",      entry.md5),
                    ("SHA256",   entry.sha256),
                ]
                for label, value in fields:
                    _vlines = _wrap_text_lines(value, "helv", 7, 523 - 155)
                    if y + len(_vlines) * 10 > 750:
                        page = doc.new_page()
                        y = 72
                    page.insert_text((88, y), f"{label}:",
                                      fontsize=7, fontname="hebo",
                                      color=(0.35, 0.35, 0.35))
                    y = _insert_text_wrapped(page, 155, y, value,
                                             7, "helv", (0, 0, 0),
                                             max_x=523.0, line_height=10.0)
                y += 4

        page.draw_line(fitz.Point(72, y), fitz.Point(420, y),
                        color=(0.9, 0.9, 0.9), width=0.2)
        y += 10

    # ── Summary footer ────────────────────────────────────────────────────────
    if y + 24 > 750:
        page = doc.new_page()
        y = 72
    y += 6
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0, 0, 0), width=0.5)
    y += 14
    summary = (f"Queried: {len(queried_hashes)}   "
               f"Found: {found_total}   "
               f"Not found: {not_found_total}")
    page.insert_text((72, y), summary,
                      fontsize=9, fontname="hebo", color=(0, 0, 0))

    doc.set_metadata({
        "author": "ot2i7ba",
        "subject": "github.com/ot2i7ba",
    })
    doc.save(str(report_path))
    doc.close()
    logger.info(f"Hash lookup report: {report_path}")
    return report_path


# ============================================================================
# HASH LIST COMPARISON (forensic external hash list match)
# ============================================================================

# Pre-compiled regex for extracting MD5 (32 hex chars) and SHA256 (64 hex chars)
# from arbitrary text.  The alternation orders 64 first so the longer pattern
# is tried before the shorter one, preventing a 64-char hash from being split
# into two 32-char matches.
_HASH_RE = re.compile(r'\b([0-9a-fA-F]{64}|[0-9a-fA-F]{32})\b')


def _extract_hashes_from_text(text: str) -> frozenset[str]:
    """Extract all MD5 and SHA256 hex strings from arbitrary text.

    Skips comment lines (starting with ``#``) and returns lowercase hash
    strings.  The regex matches only isolated hex sequences of exactly 32
    (MD5) or 64 (SHA256) characters to avoid partial matches inside longer
    strings.

    Args:
        text: Raw text content of an external hash list file.

    Returns:
        Frozenset of lowercase hex hash strings (MD5 and/or SHA256).
    """
    found: set[str] = set()
    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith("#") or not stripped:
            continue
        for m in _HASH_RE.finditer(stripped):
            found.add(m.group(1).lower())
    return frozenset(found)


_CSV_HASH_HEADERS: frozenset[str] = frozenset({"sha-256", "sha256", "md5", "hash"})


def _detect_csv_hash_cols(header_line: str) -> list[int]:
    """Return column indices that appear to contain hashes in a CSV header line.

    Args:
        header_line: First non-comment line of a file, potentially a CSV header.

    Returns:
        List of zero-based column indices whose header name matches a known hash
        column label.  Empty list when the line is not a recognised CSV header.
    """
    if "," not in header_line:
        return []
    cols = [c.strip().strip('"').lower() for c in header_line.split(",")]
    return [i for i, c in enumerate(cols) if c in _CSV_HASH_HEADERS]


def _extract_hashes_from_file(path: Path) -> frozenset[str]:
    """Extract MD5/SHA256 hashes from a file using line-by-line streaming.

    Reads the file incrementally to avoid loading large files (e.g. NSRL full
    set) entirely into RAM.  Automatically detects structured CSV files whose
    first content line is a recognised hash-column header (e.g. NSRL RDS
    format) and restricts extraction to the identified hash columns, preventing
    false positives from hex-looking filename fields.

    Lines starting with ``#`` are treated as comments and skipped.

    Args:
        path: Path to the hash list file.

    Returns:
        Frozenset of lowercase hex hash strings (MD5 and/or SHA256).
    """
    found: set[str] = set()
    for encoding in ("utf-8", "latin-1"):
        try:
            csv_hash_cols: list[int] = []
            csv_mode = False
            first_content = True
            with open(path, "r", encoding=encoding, errors="strict") as fh:
                for line in fh:
                    stripped = line.strip()
                    if not stripped or stripped.startswith("#"):
                        continue
                    if first_content:
                        first_content = False
                        cols = _detect_csv_hash_cols(stripped)
                        if cols:
                            csv_hash_cols = cols
                            csv_mode = True
                            continue  # skip header row
                    if csv_mode:
                        fields = [c.strip().strip('"') for c in stripped.split(",")]
                        for i in csv_hash_cols:
                            if i < len(fields):
                                for m in _HASH_RE.finditer(fields[i]):
                                    found.add(m.group(1).lower())
                    else:
                        for m in _HASH_RE.finditer(stripped):
                            found.add(m.group(1).lower())
            break  # encoding succeeded
        except UnicodeDecodeError:
            found.clear()
            continue
    return frozenset(found)


def _load_hashes_from_dir(
    dir_path: Path,
) -> tuple[frozenset[str], list[Path], dict[str, list[str]]]:
    """Load and merge all hashes from text files inside *dir_path*.

    Each file is read as UTF-8 text (with fallback to latin-1). Lines
    starting with ``#`` are treated as comments and skipped.  Recognised
    hash lengths are 32 hex chars (MD5) and 64 hex chars (SHA256).

    Args:
        dir_path: Path to the directory containing external hash list files.

    Returns:
        A 3-tuple of:
            - frozenset of lowercase hash strings extracted from all files
            - list of Path objects for each file that contributed at least one hash
            - dict mapping each lowercase hash string to the list of source
              filenames (``entry.name``) it was found in
    """
    all_hashes: set[str] = set()
    source_files: list[Path] = []
    hash_to_sources: dict[str, list[str]] = {}
    for entry in sorted(dir_path.rglob("*")):
        if not entry.is_file():
            continue
        try:
            hashes = _extract_hashes_from_file(entry)
            if hashes:
                all_hashes.update(hashes)
                source_files.append(entry)
                for h in hashes:
                    hash_to_sources.setdefault(h, []).append(entry.name)
                logger.debug(f"Loaded {len(hashes):,} hashes from {entry.name}")
        except OSError as e:
            logger.warning(f"Could not read hash list file {entry.name}: {e}")
    return frozenset(all_hashes), source_files, hash_to_sources


def generate_hashlist_comparison_report(
    matched_entries: list[tuple[str, list["FileEntry"]]],
    source_files: list[Path],
    total_external: int,
    json_entry_count: int,
    output_dir: Path = Path("."),
    hash_to_sources: dict[str, list[str]] | None = None,
    output_stem: str | None = None,
) -> Path:
    """Generate a PDF report for a hash list comparison against filenames.json.

    Args:
        matched_entries: List of (hash_value_lower, matching_FileEntry_list) tuples
            for every hash that was found in both the external list and filenames.json.
        source_files: List of external hash list files that were loaded.
        total_external: Total number of unique hashes extracted from external files.
        json_entry_count: Number of entries in filenames.json.
        output_dir: Directory to write the report (overridden by OUTPUT_DIR if set).
        hash_to_sources: Optional mapping of each hash to the list of source
            filenames it was found in. When provided, a "Source:" line is added
            to each match block in the report.
        output_stem: Optional file stem to use for the report name (without .pdf
            extension). When provided, the report is named ``{output_stem}.pdf``
            instead of a freshly computed timestamp. Use this to keep the PDF and
            CSV filenames in sync when both are generated together.

    Returns:
        Path to the generated report PDF.
    """
    now = datetime.now()
    if output_stem:
        report_name = f"{output_stem}.pdf"
    else:
        report_name = f"{now.strftime('%Y%m%d%H%M%S')}{HASHLIST_REPORT_SUFFIX}.pdf"
    if OUTPUT_DIR:
        report_dir = Path(OUTPUT_DIR)
        report_dir.mkdir(parents=True, exist_ok=True)
        report_path = report_dir / report_name
    else:
        report_path = output_dir / report_name

    doc = fitz.open()
    page = doc.new_page()
    y = 72

    # ── Title ──────────────────────────────────────────────────────────────────
    page.insert_text((72, y), "Hash List Comparison Report",
                      fontsize=14, fontname="helv", color=(0, 0, 0))
    y += 24
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0, 0, 0), width=0.5)
    y += 14

    # ── NOTE (disclaimer) ──────────────────────────────────────────────────────
    note_html = (
        "<b>NOTE:</b> This report has been generated automatically. "
        "It cannot be guaranteed that all hash values have been accurately "
        "identified. Results are intended as a preliminary assessment requiring "
        "further verification. All hash comparisons are case-insensitive. "
        "The original source files have not been modified. "
        "All results must be reviewed manually."
    )
    note_rect = fitz.Rect(72, y, 523, y + 52)
    page.insert_htmlbox(note_rect, note_html,
                        css="body { font-size: 7pt; color: #666666; }")
    y += 58

    # ── Context paragraph ──────────────────────────────────────────────────────
    source_names = ", ".join(f.name for f in source_files) if source_files else "—"
    context_html = (
        f"This report was generated on <b>{now.strftime('%d.%m.%Y')}</b> at "
        f"<b>{now.strftime('%H:%M:%S')}</b>. "
        f"Compared <b>{total_external:,}</b> external hashes from "
        f"<b>{len(source_files)}</b> file(s) in <b>{HASHLISTS_DIR}/</b> "
        f"({source_names}) against <b>{FILENAMES_JSON}</b> "
        f"({json_entry_count:,} entries). "
        f"Matches found: <b>{len(matched_entries)}</b>."
    )
    ctx_rect = fitz.Rect(72, y, 523, y + 52)
    page.insert_htmlbox(ctx_rect, context_html,
                        css="body { font-size: 7pt; color: #444444; }")
    y += 58
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0.8, 0.8, 0.8), width=0.3)
    y += 14

    # ── Per-match results ──────────────────────────────────────────────────────
    for hash_val, entries in matched_entries:
        hash_type = "SHA256" if len(hash_val) == 64 else "MD5"
        if y + 44 > 750:
            page = doc.new_page()
            y = 72

        page.insert_text((72, y), f"[{hash_type}]",
                          fontsize=8, fontname="hebo", color=(0.2, 0.2, 0.6))
        page.insert_text((115, y), hash_val,
                          fontsize=8, fontname="helv", color=(0, 0, 0))
        y += 13

        # Source file(s) this hash was found in
        if hash_to_sources is not None:
            sources_str = ", ".join(hash_to_sources.get(hash_val, [])) or "—"
            _slines = _wrap_text_lines(sources_str, "helv", 7, 523 - 155)
            if y + len(_slines) * 10 > 750:
                page = doc.new_page()
                y = 72
            page.insert_text((88, y), "Source:",
                              fontsize=7, fontname="hebo",
                              color=(0.35, 0.35, 0.35))
            y = _insert_text_wrapped(page, 155, y, sources_str,
                                     7, "helv", (0.1, 0.1, 0.5),
                                     max_x=523.0, line_height=10.0)

        for entry in entries:
            if entry.size:
                size_str = (f"{entry.size / 1_048_576:.2f} MB"
                            if entry.size >= 1_048_576
                            else f"{entry.size / 1024:.1f} KB")
            else:
                size_str = "unknown"
            fields = [
                ("Filename", entry.filename),
                ("Filepath", entry.filepath or "—"),
                ("Size",     size_str),
                ("MD5",      entry.md5),
                ("SHA256",   entry.sha256),
            ]
            for label, value in fields:
                _vlines = _wrap_text_lines(value, "helv", 7, 523 - 155)
                if y + len(_vlines) * 10 > 750:
                    page = doc.new_page()
                    y = 72
                page.insert_text((88, y), f"{label}:",
                                  fontsize=7, fontname="hebo",
                                  color=(0.35, 0.35, 0.35))
                y = _insert_text_wrapped(page, 155, y, value,
                                         7, "helv", (0, 0, 0),
                                         max_x=523.0, line_height=10.0)
            y += 4

        page.draw_line(fitz.Point(72, y), fitz.Point(420, y),
                        color=(0.9, 0.9, 0.9), width=0.2)
        y += 10

    # ── Summary footer ─────────────────────────────────────────────────────────
    if y + 24 > 750:
        page = doc.new_page()
        y = 72
    y += 6
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0, 0, 0), width=0.5)
    y += 14
    summary = (f"External hashes: {total_external:,}   "
               f"Entries in JSON: {json_entry_count:,}   "
               f"Matches: {len(matched_entries)}")
    page.insert_text((72, y), summary,
                      fontsize=9, fontname="hebo", color=(0, 0, 0))

    doc.set_metadata({
        "author": "ot2i7ba",
        "subject": "github.com/ot2i7ba",
    })
    doc.save(str(report_path))
    doc.close()
    logger.info(f"Hash list comparison report: {report_path}")
    return report_path


def generate_hashlist_comparison_csv(
    matched_entries: list[tuple[str, list["FileEntry"]]],
    hash_to_sources: dict[str, list[str]],
    output_path: Path,
) -> Path:
    """Generate a CSV export manifest for a hash list comparison result.

    Writes one row per (hash × source_hashlist × FileEntry) combination so
    that every matched file can be unambiguously attributed to the hash list
    it was found in.  The file is UTF-8 with BOM for broad tool compatibility
    (Excel, forensic suites).

    Args:
        matched_entries: List of (hash_value_lower, matching_FileEntry_list)
            tuples, as returned by the frozenset intersection in the handler.
        hash_to_sources: Mapping of each hash string to the list of source
            hashlist filenames it was extracted from.
        output_path: Destination path for the CSV file.

    Returns:
        The resolved output path.
    """
    with open(output_path, "w", newline="", encoding="utf-8-sig") as f:
        writer = csv.writer(f, delimiter=";")
        writer.writerow([
            "hash_value", "hash_type", "source_hashlist",
            "filename", "filepath", "size",
        ])
        for hash_val, entries in matched_entries:
            hash_type = "SHA256" if len(hash_val) == 64 else "MD5"
            sources = hash_to_sources.get(hash_val, ["unknown"])
            for source in sources:
                for entry in entries:
                    writer.writerow([
                        hash_val,
                        hash_type,
                        source,
                        entry.filename,
                        entry.filepath or "",
                        entry.size or 0,
                    ])
    logger.info(f"Hash list comparison CSV: {output_path}")
    return output_path


def _ask_source_root() -> Path | None:
    """Prompt the user for the source root directory used when filenames.json was created.

    The filepaths stored in ``filenames.json`` (and thus the export CSV) are
    relative to the directory that was scanned during list generation.  This
    function asks the user to re-enter that root so the export can reconstruct
    full absolute paths.  When ``filenames_meta.json`` exists in the current
    working directory, the ``scan_root`` stored there is shown as a suggestion
    so the user can confirm with Enter rather than retyping.

    Returns:
        Resolved ``Path`` to the source root directory, or ``None`` if the
        user cancels (enters ``b`` or presses Enter without a suggestion to
        confirm).
    """
    # Try to read scan_root from companion meta file
    suggestion: str = ""
    meta_path = Path(FILENAMES_META_JSON)
    if meta_path.is_file():
        try:
            meta = json.loads(meta_path.read_text(encoding="utf-8"))
            suggestion = meta.get("scan_root", "")
        except (OSError, json.JSONDecodeError):
            pass

    print()
    print(f"  {_C.BOLD}Source root directory{_C.RESET}")
    print(f"  {_C.DIM}Enter the directory you used when generating filenames.json.{_C.RESET}")
    if suggestion:
        print(f"  {_C.DIM}Press Enter to use the stored path: {suggestion}{_C.RESET}")
    else:
        print(f"  {_C.DIM}Example: /media/usb/evidence  or  D:\\Evidence{_C.RESET}")
    print(f"  {_C.DIM}(b = cancel){_C.RESET}\n")

    raw = input("  Source root: ").strip()

    # Empty input + suggestion → use suggestion
    if not raw and suggestion:
        raw = suggestion

    if not raw or raw.lower() == "b":
        return None
    p = Path(raw)
    if not p.is_dir():
        print(f"\n  {_C.RED}Directory not found: {p}{_C.RESET}")
        return None
    return p.resolve()


def _perform_export(
    csv_path: Path,
    source_root: Path | None = None,
) -> tuple[int, int, Path | None]:
    """Copy matched files from a hash list comparison CSV into an export directory.

    Reads the CSV manifest produced by :func:`generate_hashlist_comparison_csv`,
    creates the export directory tree, copies each file with ``shutil.copy2``
    (preserving timestamps and permissions), then moves the PDF report and the
    CSV itself into the export root.

    Directory layout::

        EXPORT_DIR/
            YYYYMMDD_HHMMSS_hashlists/
                <report>.pdf          (moved)
                <manifest>.csv        (moved)
                <source_hashlist_stem>/
                    matched_file.ext
                    ...

    Args:
        csv_path: Path to the ``*_hashlists.csv`` manifest file.
        source_root: Root directory that was used when generating
            ``filenames.json``.  When provided, each relative ``filepath``
            value from the CSV is joined to this root to form the absolute
            source path.  When ``None`` the raw ``filepath`` value is used
            as-is (which only works when it is already absolute).

    Returns:
        A 3-tuple of (files_copied, files_skipped, export_dir_path).
        ``export_dir_path`` is ``None`` if the export directory could not be
        created.
    """
    # Derive export directory name from CSV stem (e.g. 20260316_143022_hashlists)
    stem = csv_path.stem  # e.g. "20260316_143022_hashlists"
    base = Path(EXPORT_DIR) if EXPORT_DIR else Path("export")
    export_dir = base / stem
    try:
        export_dir.mkdir(parents=True, exist_ok=True)
    except OSError as e:
        logger.error(f"Cannot create export directory {export_dir}: {e}")
        return 0, 0, None

    copied = 0
    skipped = 0

    try:
        with open(csv_path, newline="", encoding="utf-8-sig") as f:
            reader = csv.DictReader(f, delimiter=";")
            for row in reader:
                filepath = row.get("filepath", "").strip()
                source = row.get("source_hashlist", "unknown").strip()
                filename = row.get("filename", "").strip()

                if not filepath or not filename:
                    skipped += 1
                    continue

                # Build full path: join source_root with the stored relative path.
                # The filepath in the CSV is always relative to the directory the
                # user selected when generating filenames.json.
                if source_root is not None:
                    src = source_root / filepath
                else:
                    src = Path(filepath)

                if not src.is_file():
                    logger.warning(f"Export: source file not found: {src}")
                    skipped += 1
                    continue

                # Subdir named after hashlist stem (strip extension)
                subdir_name = Path(source).stem if source else "unknown"
                dest_dir = export_dir / subdir_name
                try:
                    dest_dir.mkdir(parents=True, exist_ok=True)
                    dest = dest_dir / src.name
                    # Avoid overwriting: append counter suffix if name collision
                    counter = 1
                    while dest.exists():
                        dest = dest_dir / f"{src.stem}_{counter}{src.suffix}"
                        counter += 1
                    shutil.copy2(str(src), str(dest))
                    copied += 1
                    logger.debug(f"Exported: {src} → {dest}")
                except OSError as e:
                    logger.warning(f"Export: could not copy {src.name}: {e}")
                    skipped += 1

    except OSError as e:
        logger.error(f"Cannot read CSV {csv_path}: {e}")
        return copied, skipped, export_dir

    # Move PDF report (same stem, .pdf) and CSV into export root
    pdf_path = csv_path.with_suffix(".pdf")
    for artifact in (pdf_path, csv_path):
        if artifact.is_file():
            dest = export_dir / artifact.name
            try:
                shutil.move(str(artifact), str(dest))
                logger.debug(f"Moved {artifact.name} → {export_dir}")
            except OSError as e:
                logger.warning(f"Could not move {artifact.name}: {e}")
                print(f"  {_C.YELLOW}Warning: could not move {artifact.name} to export dir: {e}{_C.RESET}")

    return copied, skipped, export_dir


# ============================================================================
# SEARCH HISTORY
# ============================================================================

def _append_history(file_type: str, mode: str, files_processed: int,
                    files_with_matches: int, total_matches: int) -> None:
    """Append a search run entry to search_history.json.

    Args:
        file_type: Type label e.g. "PDF", "TXT", "Docs", "JPG"
        mode: Search mode used ("exact", "fuzzy", "regex")
        files_processed: Total files scanned
        files_with_matches: Files that had at least one match
        total_matches: Sum of all matches across all files
    """
    history_path = Path(HISTORY_FILE)
    try:
        if history_path.exists():
            history: list[dict] = json.loads(history_path.read_text(encoding="utf-8"))
        else:
            history = []
    except (OSError, json.JSONDecodeError, ValueError):
        history = []

    history.append({
        "timestamp": datetime.now().isoformat(timespec="seconds"),
        "file_type": file_type,
        "mode": mode,
        "files_processed": files_processed,
        "files_with_matches": files_with_matches,
        "total_matches": total_matches,
    })

    # Keep only the most recent entries
    if len(history) > HISTORY_MAX_ENTRIES:
        history = history[-HISTORY_MAX_ENTRIES:]

    try:
        history_path.write_text(json.dumps(history, indent=2, ensure_ascii=False), encoding="utf-8")
    except OSError as e:
        logger.warning(f"Could not write search history: {e}")


def handle_search_history() -> None:
    """Handle menu option 5: Show search history."""
    print_header()
    history_path = Path(HISTORY_FILE)

    if not history_path.exists():
        print("No search history yet. Run a search first.")
        wait_for_enter()
        return

    try:
        history: list[dict] = json.loads(history_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError, ValueError) as e:
        print(f"Could not read history: {e}")
        wait_for_enter()
        return

    if not history:
        print("Search history is empty.")
        wait_for_enter()
        return

    w = 68
    print(f"{_C.CYAN}{'=' * w}{_C.RESET}")
    print(f"  {_C.BOLD}Search History  ({len(history)} entries){_C.RESET}")
    print(f"{_C.CYAN}{'=' * w}{_C.RESET}")
    header = f"  {'Timestamp':<22} {'Type':<6} {'Mode':<7} {'Files':>5} {'Hits':>5} {'Match':>5}"
    print(f"{_C.DIM}{header}{_C.RESET}")
    print(f"  {'-' * (w - 2)}")

    for entry in reversed(history[-30:]):  # show last 30, newest first
        ts = entry.get("timestamp", "?")[:19]
        ft = entry.get("file_type", "?")[:5]
        md = entry.get("mode", "?")[:6]
        fp = entry.get("files_processed", 0)
        fm = entry.get("files_with_matches", 0)
        tm = entry.get("total_matches", 0)
        hit_color = _C.GREEN if tm > 0 else _C.DIM
        print(f"  {ts:<22} {ft:<6} {md:<7} {fp:>5} {fm:>5} {hit_color}{tm:>5}{_C.RESET}")

    print(f"{_C.CYAN}{'=' * w}{_C.RESET}")
    print(f"  {_C.DIM}Showing last {min(30, len(history))} of {len(history)} entries{_C.RESET}")
    print(f"\n  {_C.BOLD}c{_C.RESET}  Clear history   {_C.DIM}b  Back   q  Quit   Enter  Back{_C.RESET}\n")
    cmd = input("  [c/b/q/Enter]: ").strip().lower()
    if cmd == "q":
        clear_screen()
        print("Thank you for using txtFinder!")
        raise SystemExit(0)
    if cmd == "c":
        confirm = input("  Clear all search history? [y/N]: ").strip().lower()
        if confirm == "y":
            try:
                history_path.write_text("[]", encoding="utf-8")
                print(f"  {_C.GREEN}Search history cleared.{_C.RESET}")
            except OSError as e:
                print(f"  {_C.RED}Could not clear history: {e}{_C.RESET}")
            wait_for_enter()


# ============================================================================
# BATCH PROCESSING
# ============================================================================

# Text cache: search functions populate this so report generation
# can reuse extracted text without double OCR extraction.
# Lock ensures thread-safe access during parallel batch processing.
# Bounded to _EXTRACTED_TEXT_CACHE_MAX entries; oldest entry is evicted (FIFO)
# when the limit is reached so RAM usage stays bounded even with large batches.
_EXTRACTED_TEXT_CACHE_MAX = 200
_extracted_text_cache: dict[str, str] = {}
_extracted_text_cache_lock = threading.Lock()

# Second cache: actual pattern hit counts produced by the search function.
# Keyed by lowercase pattern so lookup is case-insensitive.
# Consumed by _generate_report_for_file to ensure the report reflects exactly
# the same matches that were highlighted/marked in the _checked output file.
_match_counts_cache: dict[str, dict[str, int]] = {}
_match_counts_cache_lock = threading.Lock()


def _cache_match_counts(file_path: Path, counts: dict[str, int]) -> None:
    """Cache actual per-pattern hit counts from a search run (thread-safe).

    Args:
        file_path: Source file whose counts are being stored (used as cache key).
        counts: Mapping of pattern string to number of matches found.
    """
    with _match_counts_cache_lock:
        _match_counts_cache[str(file_path)] = {k.lower(): v for k, v in counts.items()}


def _pop_cached_match_counts(file_path: Path) -> dict[str, int] | None:
    """Retrieve and remove cached match counts (thread-safe, None if absent).

    Args:
        file_path: Source file whose cached counts should be retrieved.

    Returns:
        The cached ``{pattern: count}`` dict, or ``None`` if no entry exists.
    """
    with _match_counts_cache_lock:
        return _match_counts_cache.pop(str(file_path), None)


def _build_pattern_matches_from_counts(
    json_data: list[FileEntry],
    pattern_counts: dict[str, int],
    extra_words: list[str] | None = None,
) -> dict[str, dict[str, int]]:
    """Build the nested {filename: {pattern: count}} structure from actual search counts.

    Maps each generated pattern back to its source filename using the same
    ``generate_search_patterns`` logic used during the search, so the report
    rows correspond 1-to-1 with what was highlighted in the output file.

    Args:
        json_data: FileEntry list from filenames.json
        pattern_counts: Flat {pattern_lower: count} dict from the search run
        extra_words: Optional extra search terms

    Returns:
        Nested dict suitable for ``generate_search_report``
    """
    result: dict[str, dict[str, int]] = {}
    for entry in json_data:
        fn_patterns = generate_search_patterns(entry.filename)
        result[entry.filename] = {p: pattern_counts.get(p.lower(), 0) for p in fn_patterns}
    if extra_words:
        for word in extra_words:
            fn_patterns = generate_search_patterns(word)
            result[word] = {p: pattern_counts.get(p.lower(), 0) for p in fn_patterns}
    return result


def _cache_extracted_text(file_path: Path, text: str) -> None:
    """Cache extracted text for later use by report generation (thread-safe).

    Text is normalized via ``_normalize_text`` before storage so that the
    report's match counting operates on the same Unicode/ligature-resolved
    form regardless of whether the text came from native extraction or OCR.

    Args:
        file_path: Source file that the text was extracted from (used as cache key).
        text: Raw extracted text to normalize and store.
    """
    normalized = _normalize_text(text)
    with _extracted_text_cache_lock:
        key = str(file_path)
        _extracted_text_cache[key] = normalized
        # FIFO eviction: keep cache bounded to avoid unbounded RAM growth
        while len(_extracted_text_cache) > _EXTRACTED_TEXT_CACHE_MAX:
            _extracted_text_cache.pop(next(iter(_extracted_text_cache)))
    _save_disk_cache(file_path, normalized)


def _pop_cached_text(file_path: Path) -> str | None:
    """Retrieve and remove cached text (thread-safe, returns None if not cached).

    Args:
        file_path: Source file whose cached text should be retrieved.

    Returns:
        The cached text string, or ``None`` if no entry exists.
    """
    with _extracted_text_cache_lock:
        return _extracted_text_cache.pop(str(file_path), None)


# ── Disk text cache ──────────────────────────────────────────────────────────
# Stores extracted/OCR'd text on disk so repeated searches and report
# generation do not re-run Tesseract on unchanged files.
# Cache directory: .txtfinder_cache/  (created on first write)
# Cache key: SHA-256 of the source file → .txtfinder_cache/{sha256}.txt
# Invalidation is implicit: a modified file has a different SHA-256, so its
# old cache entry is simply ignored (stale files accumulate but cause no
# incorrect behaviour).  All I/O errors are caught and silenced so that a
# missing or full disk never breaks the main search workflow.
_DISK_CACHE_DIR = Path(".txtfinder_cache")


def _disk_cache_path(file_path: Path) -> Path | None:
    """Return the disk-cache Path for *file_path*, or None on any I/O error.

    The cache key is derived from the absolute path, file size, and last-modified
    time (``mtime``).  This requires only a single ``stat()`` call (O(1)) instead
    of reading the entire file content to compute a SHA-256 hash, which makes
    cache lookups fast even for very large files.  The key is stable as long as
    the file is not modified; any content change updates ``mtime`` and therefore
    produces a different cache key, so stale entries are never served.

    Args:
        file_path: Source file to compute the cache key for.

    Returns:
        A ``Path`` inside ``_DISK_CACHE_DIR`` named ``{key}.txt``, or ``None``
        if the file's ``stat()`` cannot be read.
    """
    try:
        st = file_path.stat()
        raw = f"{file_path.resolve()}:{st.st_size}:{st.st_mtime}"
        key = hashlib.sha256(raw.encode()).hexdigest()
        return _DISK_CACHE_DIR / f"{key}.txt"
    except OSError:
        return None


def _load_disk_cache(file_path: Path) -> str | None:
    """Return cached extracted text for *file_path*, or None on cache miss.

    Args:
        file_path: Source file whose cached text is being requested.

    Returns:
        The cached text string, or ``None`` if no cache entry exists or the
        cache file cannot be read.
    """
    cache = _disk_cache_path(file_path)
    if cache is None:
        return None
    try:
        if cache.exists():
            return cache.read_text(encoding="utf-8")
    except OSError:
        pass
    return None


def _save_disk_cache(file_path: Path, text: str) -> None:
    """Persist extracted text for *file_path* to the disk cache (silent on error).

    Args:
        file_path: Source file the text was extracted from (determines cache key).
        text: Normalized extracted text to write to the cache file.
    """
    cache = _disk_cache_path(file_path)
    if cache is None:
        return
    try:
        _DISK_CACHE_DIR.mkdir(exist_ok=True)
        cache.write_text(text, encoding="utf-8")
    except OSError:
        pass


def _cleanup_disk_cache(max_age_days: int = 90) -> None:
    """Remove disk cache entries older than *max_age_days* (silent on errors).

    Called automatically at startup to prevent unbounded growth of
    ``.txtfinder_cache/``.  Only ``*.txt`` files inside the cache directory
    are considered; other files are never touched.

    Args:
        max_age_days: Entries whose last-modified time is older than this
            number of days are deleted.  Defaults to 90.
    """
    if not _DISK_CACHE_DIR.exists():
        return
    cutoff = time.time() - max_age_days * 86_400
    removed = 0
    for f in _DISK_CACHE_DIR.glob("*.txt"):
        try:
            if f.stat().st_mtime < cutoff:
                f.unlink()
                removed += 1
        except OSError:
            pass
    if removed:
        logger.debug(f"Disk cache auto-cleanup: removed {removed} entr{'y' if removed == 1 else 'ies'} older than {max_age_days} days")


def _search_doc_file(doc_path: Path, filenames: list[str],
                     fuzzy: bool = False) -> tuple[int, Path | None]:
    """Route document search to the appropriate handler by extension.

    Args:
        doc_path: Path to the document file (DOCX, XLSX, PPTX, ODT, ODS, or ODP).
        filenames: List of filenames to search for inside the document.
        fuzzy: When ``True`` use fuzzy matching instead of exact substring search.

    Returns:
        A ``(match_count, output_path)`` tuple where ``output_path`` is ``None``
        when no matches were found.
    """
    ext = doc_path.suffix.lower()
    if ext == '.docx':
        return search_in_docx(doc_path, filenames, fuzzy=fuzzy)
    if ext == '.xlsx':
        return search_in_xlsx(doc_path, filenames, fuzzy=fuzzy)
    if ext == '.pptx':
        return search_in_pptx(doc_path, filenames, fuzzy=fuzzy)
    if ext in ('.odt', '.ods', '.odp'):
        return search_in_odt(doc_path, filenames, fuzzy=fuzzy)
    return 0, None


def _search_pdf_file(pdf_path: Path, filenames: list[str], fuzzy: bool = False) -> tuple[int, Path | None, bool]:
    """Wrap search_and_highlight_pdf to return a (matches, output, ocr_used) tuple.

    Args:
        pdf_path: Path to the PDF file to search.
        filenames: List of filenames to search for.
        fuzzy: When ``True`` use fuzzy matching instead of exact search.

    Returns:
        A ``(match_count, output_path, ocr_used)`` tuple where ``output_path``
        is ``None`` when no matches were found and ``ocr_used`` indicates whether
        Tesseract OCR was invoked during processing.
    """
    hits, output, _, ocr_used = search_and_highlight_pdf(pdf_path, filenames, fuzzy=fuzzy)
    return hits, output, ocr_used


def _process_single_file(
    file_path: Path,
    filenames: list[str],
    search_fn,
    json_data: list[FileEntry],
    extra_words: list[str] | None = None,
    pattern_set: set[str] | None = None,
    stems_mode: bool = False,
    search_mode: str = "exact",
    stems_min_length: int = 0,
) -> FileSearchResult:
    """Process a single file: search then generate a report (pure logic, no UI).

    Args:
        file_path: Path to the file being processed.
        filenames: List of filenames to search for.
        search_fn: Callable ``(path, filenames, pattern_set=…) -> (matches, output)``.
        json_data: FileEntry list from filenames.json used for report generation.
        extra_words: Optional additional search terms to include in the report.
        pattern_set: Pre-built pattern set passed through to ``search_fn``.
        stems_mode: When ``True`` the report uses stem-only pattern variants.
        search_mode: Search mode label recorded in the report context paragraph.
        stems_min_length: Minimum stem length forwarded to the report generator.

    Returns:
        A populated ``FileSearchResult`` for ``file_path``.
    """
    ocr_used = False
    try:
        result = search_fn(file_path, filenames, pattern_set=pattern_set)
        matches, output = result[0], result[1]
        if len(result) > 2:
            ocr_used = bool(result[2])
    except (OSError, RuntimeError, ValueError) as e:
        logger.error(f"Search failed for {file_path.name}: {e}")
        return FileSearchResult(file_path, 0, None, None, error=str(e))

    # Use cached text from search function if available, else extract fresh
    cached = _pop_cached_text(file_path)
    report = _generate_report_for_file(file_path, json_data, cached,
                                       extra_words=extra_words, stems_mode=stems_mode,
                                       search_mode=search_mode,
                                       stems_min_length=stems_min_length,
                                       ocr_used=ocr_used) if (matches > 0 or not REPORT_ON_MATCH_ONLY) else None
    return FileSearchResult(file_path, matches, output, report)


def process_batch(
    files: list[Path],
    filenames: list[str],
    search_fn,
    json_data: list[FileEntry],
    max_workers: int | None = None,
    extra_words: list[str] | None = None,
    progress_callback=None,
    pattern_set: set[str] | None = None,
    stems_mode: bool = False,
    search_mode: str = "exact",
    stems_min_length: int = 0,
) -> tuple[list[FileSearchResult], bool]:
    """Process multiple files with optional parallelization.

    Args:
        files: List of files to search
        filenames: Filenames to search for
        search_fn: Search function (file, filenames) -> (matches, output)
        json_data: Pre-loaded JSON data for report generation
        max_workers: Maximum parallel threads.  Defaults to ``None``, which
            reads the current value of the ``MAX_WORKERS`` module constant at
            call time so that runtime changes via the Settings editor take
            effect immediately without restarting the application.
        extra_words: Optional user-defined extra search terms for report
        progress_callback: Called after each file completes (for progress bar)
        pattern_set: Pre-built pattern set (avoids repeated generation per thread)
        stems_mode: If True, report uses stem-only matching (no pattern variants)

    Returns:
        Tuple of (results list, interrupted flag)
    """
    effective_workers = max_workers if max_workers is not None else MAX_WORKERS
    logger.info(f"Batch processing {len(files)} files with {effective_workers} workers")
    interrupted = False
    worker = partial(
        _process_single_file,
        filenames=filenames,
        search_fn=search_fn,
        json_data=json_data,
        extra_words=extra_words,
        pattern_set=pattern_set,
        stems_mode=stems_mode,
        search_mode=search_mode,
        stems_min_length=stems_min_length,
    )
    if len(files) <= 1:
        results = [worker(f) for f in files]
        if progress_callback:
            for _ in results:
                progress_callback()
        return results, False

    results: list[FileSearchResult | None] = [None] * len(files)
    with ThreadPoolExecutor(max_workers=effective_workers) as executor:
        future_to_idx = {
            executor.submit(worker, f): i for i, f in enumerate(files)
        }
        try:
            for future in as_completed(future_to_idx):
                idx = future_to_idx[future]
                try:
                    results[idx] = future.result()
                except Exception as e:
                    logger.error(f"File processing failed: {e}")
                    results[idx] = FileSearchResult(files[idx], 0, None, None, error=str(e))
                if progress_callback:
                    progress_callback()
        except KeyboardInterrupt:
            interrupted = True
            for f in future_to_idx:
                f.cancel()
            logger.warning("Batch processing interrupted by user")

    return [
        r if r is not None else FileSearchResult(files[i], 0, None, None)
        for i, r in enumerate(results)
    ], interrupted


def _display_batch_results(results: list[FileSearchResult], interrupted: bool = False) -> None:
    """Display batch search results (pure UI, no return value).

    Args:
        results: List of FileSearchResult objects from a completed batch run.
        interrupted: When ``True`` an interruption notice is appended after the results.
    """
    for r in results:
        if r.error is not None:
            short = r.error[:80] + ("…" if len(r.error) > 80 else "")
            print(f"  {_C.RED}[!]{_C.RESET} {r.file.name}: {_C.RED}error — {short}{_C.RESET}")
        elif r.matches > 0:
            print(f"  {_C.GREEN}[\u2713]{_C.RESET} {r.file.name}: {r.matches} matches -> {r.output.name}")
        else:
            print(f"  {_C.DIM}[-]{_C.RESET} {r.file.name}: no matches")
        if r.report:
            logger.info(f"Report: {r.report.name}")
    if interrupted:
        print(f"  {_C.YELLOW}(Interrupted - showing partial results){_C.RESET}")


def handle_search_jpg(
    search_dir: Path = Path("."),
    config: "_SearchConfig | None" = None,
) -> None:
    """Search JPG and JPEG files for filenames from filenames.json via OCR.

    Args:
        search_dir: Directory to scan for JPG/image files.
        config: Pre-built search configuration; when provided the interactive
            mode/extra-words dialog is skipped (used by "All of the above").
    """
    print_header()

    if not TESSERACT_AVAILABLE:
        print("JPG search requires Tesseract OCR.")
        print("Install Tesseract: https://github.com/tesseract-ocr/tesseract")
        wait_for_enter()
        return

    jpg_files = get_jpg_files(search_dir)
    if not jpg_files:
        print(f"No JPG files found in {search_dir}.")
        wait_for_enter()
        return

    filenames = ensure_filenames_loaded()
    if filenames is None:
        wait_for_enter()
        return

    # Search mode + extra words (with optional profile)
    if config is not None:
        mode = config.mode
        regex_patterns = config.regex_patterns
        extra_words = config.extra_words
        stems_min_len = config.stems_min_len
    else:
        mode = _ask_search_mode()
        if mode == "back":
            return
        if mode == "stems":
            regex_patterns, extra_words = [], []
            stems_min_len = _ask_stems_min_length()
        elif mode == "profile":
            loaded = _load_profile_interactive()
            if loaded:
                mode, regex_patterns, extra_words = loaded
            else:
                mode, regex_patterns, extra_words = "exact", [], []
            stems_min_len = 0
        else:
            regex_patterns = _ask_regex_patterns() if mode == "regex" else []
            extra_words = _ask_extra_search_words()
            stems_min_len = 0

        # Offer to save profile (skipped for stems mode and when using shared config)
        if mode != "stems":
            save_name = input("  Save as profile? (name or Enter to skip): ").strip()
            if save_name:
                save_search_profile(save_name, mode, extra_words,
                                    [rx.pattern for rx in regex_patterns])
                print(f"  {_C.GREEN}Profile '{save_name}' saved{_C.RESET}")
    fuzzy = mode == "fuzzy"
    if extra_words:
        filenames = filenames + extra_words

    try:
        json_data = load_json_data()
    except (OSError, json.JSONDecodeError, ValueError) as e:
        print(f"  {_C.RED}Error: Could not load {FILENAMES_JSON}: {e}{_C.RESET}")
        wait_for_enter()
        return

    if regex_patterns:
        search_fn = partial(_search_file_with_regex, regex_patterns=regex_patterns)
    elif fuzzy:
        search_fn = partial(search_filenames_in_jpg, fuzzy=True)
    else:
        search_fn = search_filenames_in_jpg

    if mode == "stems":
        ps = build_stem_pattern_set(filenames, min_length=stems_min_len)
    elif fuzzy or regex_patterns:
        ps = None
    else:
        ps = build_pattern_set(filenames)
    with ProgressBar(len(jpg_files), "JPG files") as pb:
        results, interrupted = process_batch(jpg_files, filenames, search_fn, json_data,
                                             extra_words=extra_words or None,
                                             progress_callback=pb.advance,
                                             pattern_set=ps,
                                             stems_mode=(mode == "stems"),
                                             search_mode=mode,
                                             stems_min_length=stems_min_len)

    _display_batch_results(results, interrupted)

    total = sum(r.matches for r in results)
    with_matches = sum(1 for r in results if r.matches > 0)
    print_batch_summary("JPG", len(results), with_matches, total, mode=mode)
    with Spinner("Generating CSV summary"):
        csv_path = generate_csv_summary(results)
    print(f"  CSV: {csv_path}")
    wait_for_enter()


def handle_search_docs(
    search_dir: Path = Path("."),
    config: "_SearchConfig | None" = None,
) -> None:
    """Search DOCX, XLSX, PPTX, and ODT files for filenames from filenames.json.

    Args:
        search_dir: Directory to scan for document files.
        config: Pre-built search configuration; when provided the interactive
            mode/extra-words dialog is skipped (used by "All of the above").
    """
    print_header()

    # Optional extension filter
    all_extensions = ['.docx', '.xlsx', '.pptx', '.odt', '.ods', '.odp']
    print(f"  {_C.BOLD}Filter by type{_C.RESET}  [{', '.join(all_extensions)}]")
    print(f"  {_C.DIM}Enter extensions (comma-separated) or Enter for all{_C.RESET}")
    ext_input = input("  Type filter [all]: ").strip().lower()
    if ext_input:
        requested = [e.strip() if e.strip().startswith('.') else f".{e.strip()}"
                     for e in ext_input.split(',') if e.strip()]
        active_ext = [e for e in requested if e in all_extensions]
        if not active_ext:
            print(f"  {_C.RED}No valid extensions entered — using all.{_C.RESET}")
            active_ext = None
        else:
            print(f"  {_C.GREEN}Filtering: {', '.join(active_ext)}{_C.RESET}")
    else:
        active_ext = None

    doc_files = get_doc_files(directory=search_dir, extensions=active_ext)
    if not doc_files:
        print(f"No document files found in {search_dir}.")
        print("Supported formats: .docx, .xlsx, .pptx, .odt, .ods, .odp")
        wait_for_enter()
        return

    filenames = ensure_filenames_loaded()
    if filenames is None:
        wait_for_enter()
        return

    # Search mode + extra words (with optional profile)
    if config is not None:
        mode = config.mode
        regex_patterns = config.regex_patterns
        extra_words = config.extra_words
        stems_min_len = config.stems_min_len
    else:
        mode = _ask_search_mode()
        if mode == "back":
            return
        if mode == "stems":
            regex_patterns, extra_words = [], []
            stems_min_len = _ask_stems_min_length()
        elif mode == "profile":
            loaded = _load_profile_interactive()
            if loaded:
                mode, regex_patterns, extra_words = loaded
            else:
                mode, regex_patterns, extra_words = "exact", [], []
            stems_min_len = 0
        else:
            regex_patterns = _ask_regex_patterns() if mode == "regex" else []
            extra_words = _ask_extra_search_words()
            stems_min_len = 0

        # Offer to save profile (skipped for stems mode and when using shared config)
        if mode != "stems":
            save_name = input("  Save as profile? (name or Enter to skip): ").strip()
            if save_name:
                save_search_profile(save_name, mode, extra_words,
                                    [rx.pattern for rx in regex_patterns])
                print(f"  {_C.GREEN}Profile '{save_name}' saved{_C.RESET}")
    fuzzy = mode == "fuzzy"
    if extra_words:
        filenames = filenames + extra_words

    try:
        json_data = load_json_data()
    except (OSError, json.JSONDecodeError, ValueError) as e:
        print(f"  {_C.RED}Error: Could not load {FILENAMES_JSON}: {e}{_C.RESET}")
        wait_for_enter()
        return

    if regex_patterns:
        search_fn = partial(_search_file_with_regex, regex_patterns=regex_patterns)
    elif fuzzy:
        search_fn = partial(_search_doc_file, fuzzy=True)
    else:
        search_fn = _search_doc_file

    if mode == "stems":
        ps = build_stem_pattern_set(filenames, min_length=stems_min_len)
    elif fuzzy or regex_patterns:
        ps = None
    else:
        ps = build_pattern_set(filenames, strip_single_words=True)
    with ProgressBar(len(doc_files), "document files") as pb:
        results, interrupted = process_batch(doc_files, filenames, search_fn, json_data,
                                             extra_words=extra_words or None,
                                             progress_callback=pb.advance,
                                             pattern_set=ps,
                                             stems_mode=(mode == "stems"),
                                             search_mode=mode,
                                             stems_min_length=stems_min_len)

    _display_batch_results(results, interrupted)

    total = sum(r.matches for r in results)
    with_matches = sum(1 for r in results if r.matches > 0)
    print_batch_summary("Docs", len(results), with_matches, total, mode=mode)
    with Spinner("Generating CSV summary"):
        csv_path = generate_csv_summary(results)
    print(f"  CSV: {csv_path}")
    wait_for_enter()


def handle_search_txt(
    search_dir: Path = Path("."),
    config: "_SearchConfig | None" = None,
) -> None:
    """Search plain-text files for filenames from filenames.json.

    Args:
        search_dir: Directory to scan for TXT files.
        config: Pre-built search configuration; when provided the interactive
            mode/extra-words dialog is skipped (used by "All of the above").
    """
    print_header()

    txt_files = get_txt_files(search_dir)
    if not txt_files:
        print(f"No TXT files found in {search_dir}.")
        wait_for_enter()
        return

    filenames = ensure_filenames_loaded()
    if filenames is None:
        wait_for_enter()
        return

    # Search mode + extra words (with optional profile)
    if config is not None:
        mode = config.mode
        regex_patterns = config.regex_patterns
        extra_words = config.extra_words
        stems_min_len = config.stems_min_len
    else:
        mode = _ask_search_mode()
        if mode == "back":
            return
        if mode == "stems":
            regex_patterns, extra_words = [], []
            stems_min_len = _ask_stems_min_length()
        elif mode == "profile":
            loaded = _load_profile_interactive()
            if loaded:
                mode, regex_patterns, extra_words = loaded
            else:
                mode, regex_patterns, extra_words = "exact", [], []
            stems_min_len = 0
        else:
            regex_patterns = _ask_regex_patterns() if mode == "regex" else []
            extra_words = _ask_extra_search_words()
            stems_min_len = 0

        # Offer to save profile (skipped for stems mode and when using shared config)
        if mode != "stems":
            save_name = input("  Save as profile? (name or Enter to skip): ").strip()
            if save_name:
                save_search_profile(save_name, mode, extra_words,
                                    [rx.pattern for rx in regex_patterns])
                print(f"  {_C.GREEN}Profile '{save_name}' saved{_C.RESET}")
    fuzzy = mode == "fuzzy"
    if extra_words:
        filenames = filenames + extra_words

    try:
        json_data = load_json_data()
    except (OSError, json.JSONDecodeError, ValueError) as e:
        print(f"  {_C.RED}Error: Could not load {FILENAMES_JSON}: {e}{_C.RESET}")
        wait_for_enter()
        return

    if regex_patterns:
        search_fn = partial(_search_file_with_regex, regex_patterns=regex_patterns)
    elif fuzzy:
        search_fn = partial(search_in_txt, fuzzy=True)
    else:
        search_fn = search_in_txt

    if mode == "stems":
        ps = build_stem_pattern_set(filenames, min_length=stems_min_len)
    elif fuzzy or regex_patterns:
        ps = None
    else:
        ps = build_pattern_set(filenames, strip_single_words=True)
    with ProgressBar(len(txt_files), "TXT files") as pb:
        results, interrupted = process_batch(txt_files, filenames, search_fn, json_data,
                                             extra_words=extra_words or None,
                                             progress_callback=pb.advance,
                                             pattern_set=ps,
                                             stems_mode=(mode == "stems"),
                                             search_mode=mode,
                                             stems_min_length=stems_min_len)

    _display_batch_results(results, interrupted)

    total = sum(r.matches for r in results)
    with_matches = sum(1 for r in results if r.matches > 0)
    print_batch_summary("TXT", len(results), with_matches, total, mode=mode)
    with Spinner("Generating CSV summary"):
        csv_path = generate_csv_summary(results)
    print(f"  CSV: {csv_path}")
    wait_for_enter()


def handle_search_hash() -> None:
    """Search filenames.json for entries matching given MD5 or SHA-256 hash values.

    Loads the filenames list once, builds MD5 and SHA-256 lookup indexes, then
    loops: each iteration accepts one or more hash values, displays matches, and
    generates a PDF report.  Loops until the user enters ``b`` or ``q``.
    """
    print_header()
    json_data = _ensure_json_data_loaded()
    if json_data is None:
        wait_for_enter()
        return

    # Warn if JSON contains no hash data (stems-only or no-hashes mode)
    has_hashes = any(e.md5 != "N/A" or e.sha256 != "N/A" for e in json_data)
    if not has_hashes:
        print_header()
        print(f"  {_C.YELLOW}No hash data found in {FILENAMES_JSON}.{_C.RESET}")
        print(f"  Regenerate the filenames list with hashes (Menu 1 → option 1).")
        wait_for_enter()
        return

    # Build lookup indexes once (case-insensitive)
    md5_index: dict[str, list[FileEntry]] = {}
    sha256_index: dict[str, list[FileEntry]] = {}
    for entry in json_data:
        if entry.md5 and entry.md5 != "N/A":
            md5_index.setdefault(entry.md5.lower(), []).append(entry)
        if entry.sha256 and entry.sha256 != "N/A":
            sha256_index.setdefault(entry.sha256.lower(), []).append(entry)

    while True:
        print_header()
        print(f"  {_C.BOLD}Search for MD5/SHA256 Hash{_C.RESET}\n")
        print(f"  Loaded {len(json_data):,} entries from {FILENAMES_JSON}.")
        print(f"  Enter one or more hash values, separated by commas.")
        print(f"  {_C.DIM}MD5 = 32 hex characters   SHA256 = 64 hex characters{_C.RESET}\n")
        print(f"  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

        raw = input("  Hash value(s): ").strip()
        if raw.lower() == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        if raw.lower() == "b" or not raw:
            return

        # Parse, classify and deduplicate input hashes
        raw_tokens = [h.strip() for h in raw.split(",") if h.strip()]
        queried: list[tuple[str, str]] = []
        invalid: list[str] = []
        seen_vals: set[str] = set()

        for token in raw_tokens:
            hash_type = _classify_hash(token)
            key = token.strip().lower()
            if hash_type is None:
                invalid.append(token)
            elif key not in seen_vals:
                seen_vals.add(key)
                queried.append((key, hash_type))

        if invalid:
            print(f"\n  {_C.YELLOW}Skipped {len(invalid)} invalid value(s) "
                  f"(not a valid MD5/SHA256 hex string):{_C.RESET}")
            for inv in invalid[:5]:
                print(f"    {_C.DIM}{inv[:80]}{_C.RESET}")
            if len(invalid) > 5:
                print(f"    {_C.DIM}... and {len(invalid) - 5} more{_C.RESET}")

        if not queried:
            print(f"\n  {_C.RED}No valid hash values entered.{_C.RESET}")
            wait_for_enter()
            continue

        # Display results
        found_count = 0
        not_found_count = 0
        print()

        for hv, ht in queried:
            matches = md5_index.get(hv, []) if ht == "MD5" else sha256_index.get(hv, [])
            if matches:
                found_count += 1
                print(f"  {_C.GREEN}FOUND{_C.RESET}  [{ht}] {hv}")
                for entry in matches:
                    print(f"    Filename: {_C.BOLD}{entry.filename}{_C.RESET}")
                    if entry.filepath:
                        print(f"    Path:     {entry.filepath}")
            else:
                not_found_count += 1
                print(f"  {_C.RED}NOT FOUND{_C.RESET}  [{ht}] {hv}")

        print()
        print(f"  Searched: {len(queried)}   "
              f"Found: {_C.GREEN}{found_count}{_C.RESET}   "
              f"Not found: {_C.RED}{not_found_count}{_C.RESET}")

        # Generate PDF report
        print()
        try:
            with Spinner("Generating hash lookup report"):
                report = generate_hash_search_report(queried, json_data)
            print(f"  Report: {report.name}")
        except (OSError, RuntimeError, ValueError) as e:
            print(f"  {_C.RED}Error generating report: {e}{_C.RESET}")
            logger.error(f"Hash lookup report failed: {e}")

        wait_for_enter()


def handle_match_hash_list() -> None:
    """Compare all hashes in filenames.json against external hash list files.

    Loads all hash files from the ``hashes/`` directory (configurable via
    ``HASHLISTS_DIR``), extracts MD5/SHA-256 hashes via regex, then performs a
    frozenset intersection against hashes in filenames.json.  Generates a PDF
    report with suffix ``HASHLIST_REPORT_SUFFIX`` for all matches found.
    """
    print_header()

    # ── Verify filenames.json ──────────────────────────────────────────────────
    json_data = _ensure_json_data_loaded()
    if json_data is None:
        wait_for_enter()
        return

    has_hashes = any(e.md5 != "N/A" or e.sha256 != "N/A" for e in json_data)
    if not has_hashes:
        print(f"  {_C.YELLOW}No hash data found in {FILENAMES_JSON}.{_C.RESET}")
        print(f"  Regenerate the filenames list with hashes (Menu 1 → option 1).")
        wait_for_enter()
        return

    # ── Verify hashes directory ────────────────────────────────────────────────
    hashes_dir = Path(HASHLISTS_DIR)
    if not hashes_dir.is_dir():
        print(f"  {_C.RED}Directory \"{HASHLISTS_DIR}/\" not found.{_C.RESET}")
        print(f"  Create it and place your hash list files inside, then try again.")
        print(f"  {_C.DIM}Supported formats: any text file containing MD5 (32 hex) or SHA256 (64 hex) values.{_C.RESET}")
        print(f"  {_C.DIM}Lines starting with # are treated as comments and skipped.{_C.RESET}")
        wait_for_enter()
        return

    # ── Load external hashes ───────────────────────────────────────────────────
    print(f"  Loading hash lists from {HASHLISTS_DIR}/ ...")
    try:
        with Spinner("Reading hash list files"):
            external_hashes, source_files, hash_to_sources = _load_hashes_from_dir(hashes_dir)
    except OSError as e:
        print(f"  {_C.RED}Error reading {HASHLISTS_DIR}/: {e}{_C.RESET}")
        wait_for_enter()
        return

    if not external_hashes:
        print(f"  {_C.YELLOW}No valid MD5 or SHA256 hashes found in {HASHLISTS_DIR}/.{_C.RESET}")
        print(f"  Check that the files contain hex hash strings of 32 or 64 characters.")
        wait_for_enter()
        return

    print_header()
    print(f"  {_C.BOLD}Match Hash List{_C.RESET}\n")
    print(f"  Loaded {len(json_data):,} entries from {FILENAMES_JSON}.")
    print(f"  Loaded {len(external_hashes):,} unique hashes from "
          f"{len(source_files)} file(s) in {HASHLISTS_DIR}/:")
    for sf in source_files:
        print(f"    {_C.DIM}{sf.name}{_C.RESET}")
    print()

    # ── Build JSON hash index ──────────────────────────────────────────────────
    hash_index: dict[str, list[FileEntry]] = {}
    for entry in json_data:
        if entry.md5 and entry.md5 != "N/A":
            hash_index.setdefault(entry.md5.lower(), []).append(entry)
        if entry.sha256 and entry.sha256 != "N/A":
            hash_index.setdefault(entry.sha256.lower(), []).append(entry)

    json_hashes: frozenset[str] = frozenset(hash_index.keys())

    # ── Intersection ───────────────────────────────────────────────────────────
    matched_hashes = external_hashes & json_hashes
    matched_entries: list[tuple[str, list[FileEntry]]] = sorted(
        ((h, hash_index[h]) for h in matched_hashes),
        key=lambda t: t[0],
    )

    # ── Display results ────────────────────────────────────────────────────────
    if matched_entries:
        print(f"  {_C.GREEN}{_C.BOLD}{len(matched_entries)} match(es) found:{_C.RESET}\n")
        for hash_val, entries in matched_entries[:10]:
            hash_type = "SHA256" if len(hash_val) == 64 else "MD5"
            print(f"  {_C.GREEN}MATCH{_C.RESET}  [{hash_type}] {hash_val}")
            for entry in entries:
                print(f"    Filename: {_C.BOLD}{entry.filename}{_C.RESET}")
                if entry.filepath:
                    print(f"    Path:     {entry.filepath}")
        if len(matched_entries) > 10:
            print(f"\n  {_C.DIM}... and {len(matched_entries) - 10} more — see report.{_C.RESET}")
    else:
        print(f"  {_C.YELLOW}No matches found.{_C.RESET}")
        print(f"  None of the {len(external_hashes):,} external hashes appear in {FILENAMES_JSON}.")

    print()
    print(f"  External: {len(external_hashes):,}   "
          f"In JSON: {len(json_hashes):,}   "
          f"Matches: {_C.GREEN if matched_entries else _C.RED}{len(matched_entries)}{_C.RESET}")

    # ── Report + CSV — only when matches exist ─────────────────────────────────
    if not matched_entries:
        wait_for_enter()
        return

    print()
    csv_path: Path | None = None
    try:
        now_stem = f"{datetime.now().strftime('%Y%m%d%H%M%S')}{HASHLIST_REPORT_SUFFIX}"
        with Spinner("Generating report and CSV"):
            report = generate_hashlist_comparison_report(
                matched_entries=matched_entries,
                source_files=source_files,
                total_external=len(external_hashes),
                json_entry_count=len(json_data),
                hash_to_sources=hash_to_sources,
                output_stem=now_stem,
            )
            csv_path = generate_hashlist_comparison_csv(
                matched_entries=matched_entries,
                hash_to_sources=hash_to_sources,
                output_path=Path(f"{now_stem}.csv"),
            )
        print(f"  Report: {report.name}")
        print(f"  CSV:    {csv_path}")
    except (OSError, RuntimeError, ValueError) as e:
        print(f"  {_C.RED}Error generating report/CSV: {e}{_C.RESET}")
        logger.error(f"Hash list comparison report/CSV failed: {e}")

    # ── Offer immediate export ─────────────────────────────────────────────────
    if csv_path is not None:
        print()
        ans = input("  Export matched files now? [Y/n]: ").strip().lower()
        if ans != "n":
            source_root = _ask_source_root()
            if source_root is None:
                print(f"  {_C.DIM}Export cancelled.{_C.RESET}")
            else:
                print_header()
                print(f"  {_C.BOLD}Exporting matched files …{_C.RESET}\n")
                print(f"  Source root: {source_root}\n")
                with Spinner("Copying files"):
                    copied, skipped, export_dir = _perform_export(csv_path, source_root)
                if export_dir:
                    print(f"  {_C.GREEN}Export complete.{_C.RESET}")
                    print(f"  Directory: {export_dir.resolve()}")
                    print(f"  Copied: {copied}   Skipped: {skipped}")
                else:
                    print(f"  {_C.RED}Export failed — could not create export directory.{_C.RESET}")

    wait_for_enter()


def _apply_setting(key: str, value: "Any") -> None:
    """Write a validated configuration value to the corresponding module-level constant.

    This is the single authorised write-path for all runtime-configurable
    constants.  Using explicit ``global`` declarations keeps assignments
    statically analysable and avoids the dynamic ``globals()[key] = value``
    pattern.

    Args:
        key: Constant name, must match one of the entries in
            ``handle_settings._SETTINGS``.
        value: Already-converted and range-checked value to assign.
    """
    global FUZZY_THRESHOLD, OCR_DPI, OCR_PSM, OCR_MIN_CONFIDENCE, OCR_LANGUAGES
    global TXT_MARKER, TXT_MARKER_POSITION, MAX_SEARCH_PATTERNS, STEMS_MIN_LENGTH
    global OUTPUT_DIR, OUTPUT_SUFFIX, REPORT_SUFFIX, HASHES_SUFFIX
    global HASHLIST_REPORT_SUFFIX, FILELIST_REPORT_SUFFIX, HITLIST_SUFFIX
    global HASHLISTS_DIR, EXPORT_DIR, HIGHLIGHT_COLOR, MAX_WORKERS
    global MEDIA_EXTENSIONS, DOCUMENT_EXTENSIONS, REPORT_ON_MATCH_ONLY
    match key:
        case "FUZZY_THRESHOLD":          FUZZY_THRESHOLD = value
        case "OCR_DPI":                  OCR_DPI = value
        case "OCR_PSM":                  OCR_PSM = value
        case "OCR_MIN_CONFIDENCE":       OCR_MIN_CONFIDENCE = value
        case "OCR_LANGUAGES":            OCR_LANGUAGES = value
        case "TXT_MARKER":               TXT_MARKER = value
        case "TXT_MARKER_POSITION":      TXT_MARKER_POSITION = value
        case "MAX_SEARCH_PATTERNS":      MAX_SEARCH_PATTERNS = value
        case "STEMS_MIN_LENGTH":         STEMS_MIN_LENGTH = value
        case "OUTPUT_DIR":               OUTPUT_DIR = value
        case "OUTPUT_SUFFIX":            OUTPUT_SUFFIX = value
        case "REPORT_SUFFIX":            REPORT_SUFFIX = value
        case "HASHES_SUFFIX":            HASHES_SUFFIX = value
        case "HASHLIST_REPORT_SUFFIX":   HASHLIST_REPORT_SUFFIX = value
        case "FILELIST_REPORT_SUFFIX":   FILELIST_REPORT_SUFFIX = value
        case "HITLIST_SUFFIX":           HITLIST_SUFFIX = value
        case "HASHLISTS_DIR":            HASHLISTS_DIR = value
        case "EXPORT_DIR":               EXPORT_DIR = value
        case "HIGHLIGHT_COLOR":          HIGHLIGHT_COLOR = value
        case "MAX_WORKERS":              MAX_WORKERS = value
        case "MEDIA_EXTENSIONS":         MEDIA_EXTENSIONS = value
        case "DOCUMENT_EXTENSIONS":      DOCUMENT_EXTENSIONS = value
        case "REPORT_ON_MATCH_ONLY":     REPORT_ON_MATCH_ONLY = value


def handle_settings() -> None:
    """Handle menu option 6: Interactive settings editor."""
    CONFIG_FILE = "txtfinder_config.json"
    # Map setting name → (global var name, type, description, min, max)
    _SETTINGS: list[tuple[str, str, type, str, float | None, float | None]] = [
        ("FUZZY_THRESHOLD",   "FUZZY_THRESHOLD",   float, "Fuzzy similarity threshold (0.0–1.0)", 0.0, 1.0),
        ("OCR_DPI",           "OCR_DPI",           int,   "OCR render DPI (100–600)",              100, 600),
        ("OCR_PSM",           "OCR_PSM",           int,   "Tesseract PSM (0-13; 3=auto, 6=block, 11=sparse)", 0, 13),
        ("OCR_MIN_CONFIDENCE","OCR_MIN_CONFIDENCE",int,   "OCR min confidence (0–100)",             0,   100),
        ("OCR_LANGUAGES",     "OCR_LANGUAGES",     str,   "Tesseract language codes (e.g. eng+deu)",None,None),
        ("TXT_MARKER",        "TXT_MARKER",        str,   "TXT match marker string",                None,None),
        ("TXT_MARKER_POSITION","TXT_MARKER_POSITION",str, 'Marker position ("before" or "after")',  None,None),
        ("MAX_SEARCH_PATTERNS","MAX_SEARCH_PATTERNS",int, "Max unique search patterns (10–5000)",   10,  5000),
        ("STEMS_MIN_LENGTH",  "STEMS_MIN_LENGTH",  int,   "Stems mode: min stem length (0 = no limit)", 0, 500),
        ("OUTPUT_DIR",        "OUTPUT_DIR",        str,   "Output directory (empty = source dir)",  None,None),
        # Output file suffixes
        ("OUTPUT_SUFFIX",     "OUTPUT_SUFFIX",     str,          "Output file suffix for search results (e.g. _checked)",    None, None),
        ("REPORT_SUFFIX",     "REPORT_SUFFIX",     str,          "Output file suffix for search reports (e.g. _report)",     None, None),
        ("HASHES_SUFFIX",         "HASHES_SUFFIX",         str, "Output file suffix for hash lookup reports (e.g. _hashes)",           None, None),
        ("HASHLIST_REPORT_SUFFIX","HASHLIST_REPORT_SUFFIX",str, "Output file suffix for hash list comparison reports (e.g. _hashlists)", None, None),
        ("FILELIST_REPORT_SUFFIX","FILELIST_REPORT_SUFFIX",str, "Output file suffix for file list comparison reports (e.g. _filelist)",   None, None),
        ("HITLIST_SUFFIX",        "HITLIST_SUFFIX",        str, "Output file suffix for file list hit result .txt files (e.g. _hitlist)",  None, None),
        ("HASHLISTS_DIR",         "HASHLISTS_DIR",         str, "Directory containing external hash list files (e.g. hashes)",           None, None),
        ("EXPORT_DIR",            "EXPORT_DIR",            str, "Base directory for hash match file exports (e.g. export)",              None, None),
        # Highlighting and performance
        ("HIGHLIGHT_COLOR",   "HIGHLIGHT_COLOR",   _color_tuple,       "PDF highlight color as R,G,B floats 0.0–1.0 (e.g. '1.0,0.0,0.0' = red)", None, None),
        ("MAX_WORKERS",       "MAX_WORKERS",       int,                "Parallel worker threads (1–32; env: TXTFINDER_WORKERS)",   1,    32),
        # File type filters (comma-separated extension lists for list generation)
        ("MEDIA_EXTENSIONS",    "MEDIA_EXTENSIONS",    _parse_extensions, "Media file extensions, comma-separated (e.g. .mp4,.mkv,.jpg,.png)", None, None),
        ("DOCUMENT_EXTENSIONS", "DOCUMENT_EXTENSIONS", _parse_extensions, "Document file extensions, comma-separated (e.g. .pdf,.docx,.xlsx,.txt)", None, None),
        # Report behaviour
        ("REPORT_ON_MATCH_ONLY", "REPORT_ON_MATCH_ONLY", _parse_bool, "Only generate reports for files with at least one match (true/false)", None, None),
    ]

    while True:
        print_header()
        print(f"  {_C.BOLD}Settings Editor{_C.RESET}\n")
        _cur_globals = globals()
        for i, (name, var, typ, desc, lo, hi) in enumerate(_SETTINGS, 1):
            cur = _cur_globals[var]
            display = ", ".join(sorted(cur)) if isinstance(cur, frozenset) else repr(cur)
            print(f"  {_C.BOLD}{i}{_C.RESET}  {name} = {_C.CYAN}{display}{_C.RESET}")
            print(f"     {_C.DIM}{desc}{_C.RESET}")
        print()
        print(f"  {_C.BOLD}s{_C.RESET}  Save to {CONFIG_FILE}")
        print(f"  {_C.BOLD}l{_C.RESET}  Load from {CONFIG_FILE}")
        print(f"  {_C.BOLD}c{_C.RESET}  Clear disk cache (.txtfinder_cache/)")
        print(f"  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

        choice = input("  Select setting to change [1-{}/s/l/c/b/q]: ".format(len(_SETTINGS))).strip().lower()

        if choice == "b":
            break
        elif choice == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        elif choice == "c":
            if _DISK_CACHE_DIR.exists():
                files = list(_DISK_CACHE_DIR.glob("*.txt"))
                total_bytes = sum(f.stat().st_size for f in files if f.is_file())
                total_mb = total_bytes / (1024 * 1024)
                print(f"  Disk cache: {len(files)} entries, {total_mb:.1f} MB in {_DISK_CACHE_DIR}/")
                confirm = input("  Clear all cache entries? [y/N]: ").strip().lower()
                if confirm == "y":
                    removed = 0
                    for f in files:
                        try:
                            f.unlink()
                            removed += 1
                        except OSError:
                            pass
                    print(f"  {_C.GREEN}Removed {removed} cache entr{'y' if removed == 1 else 'ies'}.{_C.RESET}")
            else:
                print(f"  {_C.DIM}Disk cache is empty ({_DISK_CACHE_DIR}/ does not exist).{_C.RESET}")
            wait_for_enter()
        elif choice == "s":
            try:
                _cur = globals()
                data = {
                    var: sorted(_cur[var]) if isinstance(_cur[var], frozenset) else _cur[var]
                    for _, var, *_ in _SETTINGS
                }
                Path(CONFIG_FILE).write_text(json.dumps(data, indent=2), encoding="utf-8")
                print(f"  {_C.GREEN}Saved to {CONFIG_FILE}{_C.RESET}")
            except OSError as e:
                print(f"  {_C.RED}Save failed: {e}{_C.RESET}")
            wait_for_enter()
        elif choice == "l":
            try:
                raw = json.loads(Path(CONFIG_FILE).read_text(encoding="utf-8"))
                for _, var, typ, _, lo, hi in _SETTINGS:
                    if var in raw:
                        val = typ(raw[var])
                        if lo is not None and val < lo:
                            val = lo  # type: ignore[operator]
                        if hi is not None and val > hi:
                            val = hi  # type: ignore[operator]
                        _apply_setting(var, val)
                print(f"  {_C.GREEN}Loaded from {CONFIG_FILE}{_C.RESET}")
            except FileNotFoundError:
                print(f"  {_C.RED}{CONFIG_FILE} not found.{_C.RESET}")
            except (json.JSONDecodeError, ValueError, OSError) as e:
                print(f"  {_C.RED}Load failed: {e}{_C.RESET}")
            wait_for_enter()
        else:
            try:
                idx = int(choice) - 1
                if not (0 <= idx < len(_SETTINGS)):
                    raise ValueError
            except ValueError:
                print(f"  {_C.RED}Invalid choice.{_C.RESET}")
                wait_for_enter()
                continue

            name, var, typ, desc, lo, hi = _SETTINGS[idx]
            cur = globals()[var]
            cur_display = ", ".join(sorted(cur)) if isinstance(cur, frozenset) else repr(cur)
            print(f"\n  {_C.BOLD}{name}{_C.RESET} = {cur_display}")
            print(f"  {_C.DIM}{desc}{_C.RESET}")
            raw_in = input("  New value (Enter to keep): ").strip()
            if not raw_in:
                continue
            try:
                new_val = typ(raw_in)
                if lo is not None and new_val < lo:  # type: ignore[operator]
                    print(f"  {_C.RED}Value must be >= {lo}{_C.RESET}")
                    wait_for_enter()
                    continue
                if hi is not None and new_val > hi:  # type: ignore[operator]
                    print(f"  {_C.RED}Value must be <= {hi}{_C.RESET}")
                    wait_for_enter()
                    continue
                _apply_setting(var, new_val)
                print(f"  {_C.GREEN}{name} set to {new_val!r}{_C.RESET}")
            except ValueError:
                print(f"  {_C.RED}Invalid value for type {typ.__name__}.{_C.RESET}")
                wait_for_enter()


def _ask_search_dir() -> Path | None:
    """Ask the user which directory to search for files.

    Offers three options matching the list-generation wizard:

    1. ``input/`` sub-directory (if it exists)
    2. A custom path entered by the user
    3. The current working directory

    Validates the selected path with :func:`_validate_scan_directory`.

    Returns:
        Resolved ``Path`` to the chosen search directory, or ``None`` if the
        user cancels (``b`` or ``q``).
    """
    print_header()
    input_dir = Path(INPUT_DIR)
    input_exists = input_dir.is_dir()
    opt1_note = "" if input_exists else f"  {_C.DIM}(not found){_C.RESET}"

    print(f"  {_C.BOLD}Search directory{_C.RESET}\n")
    print(f"  {_C.BOLD}1{_C.RESET}  From \"{INPUT_DIR}\" directory (./{INPUT_DIR}/){opt1_note}")
    print(f"  {_C.BOLD}2{_C.RESET}  From custom path")
    print(f"  {_C.BOLD}3{_C.RESET}  From current directory (.)")
    print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

    choice = input("  Select source [1/2/3/b/q]: ").strip().lower()
    if choice == "q":
        clear_screen()
        print("Thank you for using txtFinder!")
        raise SystemExit(0)
    if choice == "b" or not choice:
        return None

    if choice == "1":
        if not input_exists:
            print(f"\n  {_C.RED}Directory \"./{INPUT_DIR}/\" does not exist.{_C.RESET}")
            wait_for_enter()
            return None
        ok, err = _validate_scan_directory(input_dir)
        if not ok:
            print(f"\n  {_C.RED}{err}{_C.RESET}")
            wait_for_enter()
            return None
        return input_dir.resolve()

    elif choice == "2":
        raw_path = input("\n  Enter directory path: ").strip()
        if not raw_path:
            print(f"\n  {_C.RED}No path entered.{_C.RESET}")
            wait_for_enter()
            return None
        p = Path(raw_path).resolve()
        ok, err = _validate_scan_directory(p)
        if not ok:
            print(f"\n  {_C.RED}{err}{_C.RESET}")
            wait_for_enter()
            return None
        return p

    elif choice == "3":
        return Path(".").resolve()
    else:
        print(f"\n  {_C.RED}Invalid choice.{_C.RESET}")
        wait_for_enter()
        return None


def handle_search_files() -> None:
    """Handle menu option 4: Search filenames in files with a file-type submenu.

    First asks the user which directory to search (input/, custom path, or
    CWD), then presents a submenu for PDF, JPG, Docs (DOCX/XLSX/PPTX/ODT),
    TXT, or all supported file types at once.  The selected directory is
    passed to each handler so the correct files are found.  Loops until the
    user enters ``b`` or ``q``.
    """
    search_dir = _ask_search_dir()
    if search_dir is None:
        return

    while True:
        print_header()
        print(f"  {_C.BOLD}Search in Files{_C.RESET}")
        print(f"  {_C.DIM}Directory: {search_dir}{_C.RESET}\n")
        print(f"  {_C.BOLD}1{_C.RESET}  PDF files")
        print(f"  {_C.BOLD}2{_C.RESET}  JPG / image files")
        print(f"  {_C.BOLD}3{_C.RESET}  Documents (DOCX, XLSX, PPTX, ODT)")
        print(f"  {_C.BOLD}4{_C.RESET}  TXT files")
        print(f"  {_C.BOLD}a{_C.RESET}  All of the above")
        print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

        choice = input("  Select file type [1/2/3/4/a/b/q]: ").strip().lower()
        if choice == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        if choice == "b" or not choice:
            return

        if choice == "1":
            handle_search_pdf(search_dir)
        elif choice == "2":
            handle_search_jpg(search_dir)
        elif choice == "3":
            handle_search_docs(search_dir)
        elif choice == "4":
            handle_search_txt(search_dir)
        elif choice == "a":
            print_header()
            print(f"  {_C.BOLD}All file types — configure search once{_C.RESET}")
            print(f"  {_C.DIM}The same settings will be applied to PDF, JPG, Docs and TXT.{_C.RESET}\n")
            shared_cfg = _ask_search_config()
            if shared_cfg is None:
                continue
            handle_search_pdf(search_dir, config=shared_cfg)
            handle_search_jpg(search_dir, config=shared_cfg)
            handle_search_docs(search_dir, config=shared_cfg)
            handle_search_txt(search_dir, config=shared_cfg)
        else:
            print(f"\n  {_C.RED}Invalid choice.{_C.RESET}")
            wait_for_enter()


def handle_search_filenames_in_json() -> None:
    """Search filenames.json for entries matching a filename substring.

    Performs a case-insensitive substring match on ``entry.filename`` (bare
    name including extension).  Displays filename, filepath, size, MD5, and
    SHA-256 for each match.  Loops until the user enters ``b`` or ``q``.
    """
    print_header()
    json_data = _ensure_json_data_loaded()
    if json_data is None:
        wait_for_enter()
        return

    while True:
        print_header()
        print(f"  {_C.BOLD}Search for Filenames{_C.RESET}\n")
        print(f"  Loaded {len(json_data):,} entries from {FILENAMES_JSON}.")
        print(f"  Enter a filename or partial name to search for (case-insensitive).")
        print(f"  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

        raw = input("  Filename search term: ").strip()
        if raw.lower() == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        if raw.lower() == "b" or not raw:
            return

        term = raw.lower()
        matches = [e for e in json_data if term in e.filename.lower()]

        print()
        if not matches:
            print(f"  {_C.YELLOW}No matches found for \"{raw}\".{_C.RESET}")
        else:
            print(f"  {_C.GREEN}{len(matches):,} match(es) for \"{raw}\":{_C.RESET}\n")
            for entry in matches:
                size_str = f"{entry.size:,} bytes" if entry.size else "unknown size"
                print(f"  {_C.BOLD}{entry.filename}{_C.RESET}")
                if entry.filepath:
                    print(f"    Path:   {entry.filepath}")
                print(f"    Size:   {size_str}")
                if entry.md5 != "N/A":
                    print(f"    MD5:    {entry.md5}")
                if entry.sha256 != "N/A":
                    print(f"    SHA256: {entry.sha256}")
                print()

        wait_for_enter()


def handle_search_keywords_in_json() -> None:
    """Search filenames.json for entries matching comma-separated keywords.

    Matches each keyword case-insensitively against both ``entry.filename``
    and ``entry.filepath``.  Results indicate which field each keyword matched.
    An entry is shown once per keyword that matches it; duplicate entries
    across keywords are shown separately.  Loops until the user enters
    ``b`` or ``q``.
    """
    print_header()
    json_data = _ensure_json_data_loaded()
    if json_data is None:
        wait_for_enter()
        return

    while True:
        print_header()
        print(f"  {_C.BOLD}Search for Keywords{_C.RESET}\n")
        print(f"  Loaded {len(json_data):,} entries from {FILENAMES_JSON}.")
        print(f"  Enter keywords separated by commas (searched in filename and path).")
        print(f"  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

        raw = input("  Keyword(s): ").strip()
        if raw.lower() == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        if raw.lower() == "b" or not raw:
            return

        keywords = [kw.strip() for kw in raw.split(",") if kw.strip()]
        if not keywords:
            print(f"\n  {_C.RED}No valid keywords entered.{_C.RESET}")
            wait_for_enter()
            continue

        print()
        total_hits = 0

        for kw in keywords:
            kw_lower = kw.lower()
            hits: list[tuple[FileEntry, str]] = []
            for entry in json_data:
                matched_fields: list[str] = []
                if kw_lower in entry.filename.lower():
                    matched_fields.append("filename")
                if entry.filepath and kw_lower in entry.filepath.lower():
                    matched_fields.append("path")
                if matched_fields:
                    hits.append((entry, ", ".join(matched_fields)))

            if not hits:
                print(f"  {_C.YELLOW}No matches for \"{kw}\".{_C.RESET}\n")
                continue

            total_hits += len(hits)
            print(f"  {_C.GREEN}{len(hits):,} match(es) for \"{kw}\":{_C.RESET}\n")
            for entry, field_label in hits:
                size_str = f"{entry.size:,} bytes" if entry.size else "unknown size"
                print(f"  {_C.BOLD}{entry.filename}{_C.RESET}  {_C.DIM}[matched: {field_label}]{_C.RESET}")
                if entry.filepath:
                    print(f"    Path:   {entry.filepath}")
                print(f"    Size:   {size_str}")
                if entry.md5 != "N/A":
                    print(f"    MD5:    {entry.md5}")
                if entry.sha256 != "N/A":
                    print(f"    SHA256: {entry.sha256}")
                print()

        if total_hits == 0:
            print(f"  {_C.YELLOW}No matches found for any keyword.{_C.RESET}")

        wait_for_enter()


def handle_export_hash_matches() -> None:
    """Export files matched in a previous hash list comparison.

    Scans the current working directory (and OUTPUT_DIR when configured) for
    CSV files whose names end with ``HASHLIST_REPORT_SUFFIX``, presents a
    numbered selection list, and calls ``_perform_export`` on the chosen file.
    The export copies every matched file into a timestamped subdirectory and
    moves the associated PDF report and CSV to the export root.
    """
    print_header()
    print(f"  {_C.BOLD}Export Matched Files{_C.RESET}\n")

    # Collect candidate CSV files from CWD (and OUTPUT_DIR if set)
    search_dirs: list[Path] = [Path(".")]
    if OUTPUT_DIR:
        od = Path(OUTPUT_DIR)
        if od.is_dir() and od != Path("."):
            search_dirs.append(od)

    candidates: list[Path] = []
    for d in search_dirs:
        for p in sorted(d.iterdir()):
            if p.is_file() and p.suffix.lower() == ".csv" and p.stem.endswith(HASHLIST_REPORT_SUFFIX):
                candidates.append(p)

    if not candidates:
        print(f"  {_C.YELLOW}No hash match CSV files found.{_C.RESET}")
        print(f"  Run 'Match Hash List' first to generate a CSV.")
        print()
        wait_for_enter()
        return

    while True:
        print(f"  Available hash match CSV files:\n")
        for idx, p in enumerate(candidates, 1):
            print(f"  {_C.BOLD}{idx}{_C.RESET}  {p}")
        print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

        choice = input(f"  Select file [1/{len(candidates)}/b/q]: ").strip().lower()
        if choice == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        if choice == "b" or not choice:
            return

        try:
            idx = int(choice)
            if not (1 <= idx <= len(candidates)):
                raise ValueError
            break
        except ValueError:
            print(f"\n  {_C.RED}Invalid input. Enter a number, 'b' to go back, or 'q' to quit.{_C.RESET}\n")

    selected = candidates[idx - 1]
    source_root = _ask_source_root()
    if source_root is None:
        print(f"  {_C.DIM}Export cancelled.{_C.RESET}")
        wait_for_enter()
        return

    print_header()
    print(f"  {_C.BOLD}Exporting matched files …{_C.RESET}\n")
    print(f"  Source root: {source_root}\n")
    with Spinner("Copying files"):
        copied, skipped, export_dir = _perform_export(selected, source_root)

    if export_dir:
        print(f"  {_C.GREEN}Export complete.{_C.RESET}")
        print(f"  Directory: {export_dir.resolve()}")
        print(f"  Copied: {copied}   Skipped: {skipped}")
    else:
        print(f"  {_C.RED}Export failed — could not create export directory.{_C.RESET}")

    wait_for_enter()


# ============================================================================
# FILE LIST COMPARISON
# ============================================================================

def generate_filelist_comparison_report(
    matched: list[tuple[str, list["FileEntry"]]],
    total_queried: int,
    list_path: Path,
    json_data: list["FileEntry"],
    output_dir: Path = Path("."),
) -> Path:
    """Generate a PDF report for a file list comparison against filenames.json.

    Only matched entries are included in the report body.  The summary footer
    always shows the full queried / found / not-found counts.

    Args:
        matched: Sorted list of (queried_name_lower, matching_FileEntry_list) tuples
            for every name found in both the file list and filenames.json.
        total_queried: Number of unique (deduplicated) names from the file list
            that were actually compared against filenames.json.
        list_path: Path to the source file list (.txt) used for the comparison.
        json_data: Entries loaded from filenames.json.
        output_dir: Directory to write the report (overridden by OUTPUT_DIR if set).

    Returns:
        Path to the generated report PDF.
    """
    now = datetime.now()
    report_name = f"{now.strftime('%Y%m%d%H%M%S')}{FILELIST_REPORT_SUFFIX}.pdf"
    if OUTPUT_DIR:
        report_dir = Path(OUTPUT_DIR)
        report_dir.mkdir(parents=True, exist_ok=True)
        report_path = report_dir / report_name
    else:
        report_path = output_dir / report_name

    found_total = len(matched)
    not_found_total = total_queried - found_total

    doc = fitz.open()
    page = doc.new_page()
    y = 72

    # ── Title ─────────────────────────────────────────────────────────────────
    page.insert_text((72, y), "File List Comparison Report",
                      fontsize=14, fontname="helv", color=(0, 0, 0))
    y += 24
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0, 0, 0), width=0.5)
    y += 14

    # ── NOTE (disclaimer) ─────────────────────────────────────────────────────
    note_html = (
        "<b>NOTE:</b> This report has been generated automatically by comparing "
        "a provided file list against filenames.json. Matching is case-insensitive "
        "and limited to the exact filename including its extension. It cannot be "
        "guaranteed that all files have been identified. The results are intended "
        "as a preliminary assessment and require further manual verification."
    )
    note_rect = fitz.Rect(72, y, 523, y + 52)
    page.insert_htmlbox(note_rect, note_html,
                        css="body { font-size: 7pt; color: #666666; }")
    y += 56

    y += 18

    # ── Context paragraph ─────────────────────────────────────────────────────
    context_html = (
        f"This report was generated on <b>{now.strftime('%d.%m.%Y')}</b> at "
        f"<b>{now.strftime('%H:%M:%S')}</b>. "
        f"File list: <b>{list_path.name}</b> ({total_queried:,} names queried). "
        f"Compared against <b>{FILENAMES_JSON}</b> ({len(json_data):,} entries). "
        f"For additional details see <b>txtfinder.log</b>."
    )
    ctx_rect = fitz.Rect(72, y, 523, y + 40)
    page.insert_htmlbox(ctx_rect, context_html,
                        css="body { font-size: 7pt; color: #444444; }")
    y += 44
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0.8, 0.8, 0.8), width=0.3)
    y += 14

    # ── Per-name match blocks ─────────────────────────────────────────────────
    for queried_name, entries in matched:
        _qlines = len(_wrap_text_lines(queried_name, "hebo", 8, 451))
        if y + _qlines * 13 + 31 > 750:
            page = doc.new_page()
            y = 72

        y = _insert_text_wrapped(page, 72, y, queried_name,
                                  8, "hebo", (0.2, 0.2, 0.6),
                                  line_height=13.0)

        for entry in entries:
            if entry.size:
                size_str = (f"{entry.size / 1_048_576:.2f} MB"
                            if entry.size >= 1_048_576
                            else f"{entry.size / 1024:.1f} KB")
            else:
                size_str = "unknown"
            fields = [
                ("Filename", entry.filename),
                ("Filepath", entry.filepath or "—"),
                ("Size",     size_str),
                ("MD5",      entry.md5 or "—"),
                ("SHA256",   entry.sha256 or "—"),
            ]
            for label, value in fields:
                _vlines = _wrap_text_lines(value, "helv", 7, 523 - 155)
                if y + len(_vlines) * 10 > 750:
                    page = doc.new_page()
                    y = 72
                page.insert_text((88, y), f"{label}:",
                                  fontsize=7, fontname="hebo",
                                  color=(0.35, 0.35, 0.35))
                y = _insert_text_wrapped(page, 155, y, value,
                                         7, "helv", (0, 0, 0),
                                         max_x=523.0, line_height=10.0)
            y += 4

        page.draw_line(fitz.Point(72, y), fitz.Point(420, y),
                        color=(0.9, 0.9, 0.9), width=0.2)
        y += 10

    # ── Summary footer ────────────────────────────────────────────────────────
    if y + 24 > 750:
        page = doc.new_page()
        y = 72
    y += 6
    page.draw_line(fitz.Point(72, y), fitz.Point(523, y),
                    color=(0, 0, 0), width=0.5)
    y += 14
    summary = (f"Queried: {total_queried:,}   "
               f"Found: {found_total:,}   "
               f"Not found: {not_found_total:,}")
    page.insert_text((72, y), summary,
                      fontsize=9, fontname="hebo", color=(0, 0, 0))

    doc.set_metadata({
        "author": "ot2i7ba",
        "subject": "github.com/ot2i7ba",
    })
    doc.save(str(report_path))
    doc.close()
    logger.info(f"File list comparison report: {report_path}")
    return report_path


def handle_search_filelist_in_json() -> None:
    """Compare a plain-text file list against filenames.json.

    Reads a user-supplied text file (one filename per line, with extension),
    builds a case-insensitive lookup against filenames.json, and generates a
    PDF report for all matches found.
    """
    print_header()

    # ── Verify filenames.json ──────────────────────────────────────────────────
    json_data = _ensure_json_data_loaded()
    if json_data is None:
        wait_for_enter()
        return

    # ── Step 1: Directory selection ────────────────────────────────────────────
    list_dir = _ask_search_dir()
    if list_dir is None:
        return

    # ── Step 2: Find .txt files in selected directory ──────────────────────────
    txt_files = sorted(
        f for f in list_dir.glob("*.txt")
        if f.is_file() and not _is_output_file(f.stem.lower())
    )
    if not txt_files:
        print_header()
        print(f"  {_C.YELLOW}No .txt files found in {list_dir}.{_C.RESET}")
        print(f"  {_C.DIM}Place a file list (.txt, one filename per line) there and try again.{_C.RESET}")
        wait_for_enter()
        return

    # ── Step 3: File selection loop ────────────────────────────────────────────
    while True:
        print_header()
        print(f"  {_C.BOLD}Search for Filelist{_C.RESET}\n")
        print(f"  Select a file list from {list_dir}:\n")
        for idx, txt_file in enumerate(txt_files, 1):
            print(f"  {_C.BOLD}{idx}{_C.RESET}  {txt_file.name}")
        print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

        choice = input("  Enter choice: ").strip().lower()
        if choice == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        if choice == "b":
            return

        try:
            idx = int(choice) - 1
            if 0 <= idx < len(txt_files):
                list_path = txt_files[idx]
                break
            print(f"  {_C.RED}Invalid selection. Enter a number between 1 and {len(txt_files)}.{_C.RESET}")
        except ValueError:
            print(f"  {_C.RED}Invalid input. Enter a number, 'b' to go back, or 'q' to quit.{_C.RESET}")
        wait_for_enter()

    # ── Read file list ─────────────────────────────────────────────────────────
    try:
        queried_names: list[str] = []
        with open(list_path, "r", encoding="utf-8", errors="replace") as fh:
            for line in fh:
                stripped = line.strip()
                if stripped and not stripped.startswith("#"):
                    queried_names.append(stripped)
    except OSError as e:
        print(f"  {_C.RED}Could not read {list_path.name}: {e}{_C.RESET}")
        wait_for_enter()
        return

    if not queried_names:
        print(f"  {_C.YELLOW}No filenames found in {list_path.name}.{_C.RESET}")
        wait_for_enter()
        return

    total_queried = len(queried_names)

    # ── Build case-insensitive filename index from filenames.json ──────────────
    # O(N) build, O(1) per lookup — optimal for large JSON sets
    filename_index: dict[str, list[FileEntry]] = {}
    for entry in json_data:
        key = entry.filename.lower()
        filename_index.setdefault(key, []).append(entry)

    # ── Match: dedup (case-insensitive) preserving original case for output ────
    seen: set[str] = set()
    unique_names: list[str] = []  # original case from the file list
    for name in queried_names:
        key = name.lower()
        if key not in seen:
            seen.add(key)
            unique_names.append(name)  # preserve original case

    matched: list[tuple[str, list[FileEntry]]] = sorted(
        ((name, filename_index[name.lower()]) for name in unique_names
         if name.lower() in filename_index),
        key=lambda t: t[0].lower(),
    )

    # ── Display summary ────────────────────────────────────────────────────────
    print_header()
    print(f"  {_C.BOLD}Search for Filelist{_C.RESET}\n")
    print(f"  File list : {list_path.name}  ({total_queried:,} names, "
          f"{len(unique_names):,} unique)")
    print(f"  JSON      : {FILENAMES_JSON}  ({len(json_data):,} entries)")
    print()

    if not matched:
        print(f"  {_C.YELLOW}No matches found.{_C.RESET}")
        wait_for_enter()
        return

    print(f"  {_C.GREEN}{_C.BOLD}{len(matched):,} match(es) found:{_C.RESET}\n")
    for queried_name, entries in matched[:10]:
        hit_count = len(entries)
        suffix = f"  {_C.DIM}({hit_count}x){_C.RESET}" if hit_count > 1 else ""
        print(f"    {_C.CYAN}{queried_name}{_C.RESET}{suffix}")
    if len(matched) > 10:
        print(f"    {_C.DIM}… and {len(matched) - 10} more (see report){_C.RESET}")
    print()

    # ── Generate report ────────────────────────────────────────────────────────
    try:
        with Spinner("Generating filelist comparison report"):
            report_path = generate_filelist_comparison_report(
                matched, len(unique_names), list_path, json_data
            )
        print(f"  {_C.GREEN}Report saved: {report_path}{_C.RESET}")
    except (OSError, RuntimeError) as e:
        print(f"  {_C.RED}Report generation failed: {e}{_C.RESET}")
        logger.error(f"generate_filelist_comparison_report failed: {e}")

    # ── Optional hitlist ───────────────────────────────────────────────────────
    print()
    create_hitlist = input("  Create hitlist (.txt with matched filenames)? [y/n]: ").strip().lower()
    if create_hitlist == "y":
        hitlist_path = _output_path(list_path, HITLIST_SUFFIX, ".txt")
        try:
            with open(hitlist_path, "w", encoding="utf-8") as fh:
                for name, _ in matched:
                    fh.write(name + "\n")
            print(f"  {_C.GREEN}Hitlist saved: {hitlist_path}{_C.RESET}")
            logger.info(f"Hitlist saved: {hitlist_path}")
        except OSError as e:
            print(f"  {_C.RED}Could not write hitlist: {e}{_C.RESET}")
            logger.error(f"Hitlist write failed: {e}")

    wait_for_enter()


def handle_search_in_json() -> None:
    """Handle menu option 3: Search within filenames.json with a search-type submenu.

    Presents a submenu with six search/export modes against the loaded
    filenames list: substring search by filename, file list comparison against
    a .txt file, keyword search across filename and filepath, hash lookup with
    PDF report generation, forensic hash list comparison against external files
    in the ``hashes/`` directory, and manual export of matched files from a
    previously generated CSV.
    Loops so that after returning from a sub-handler the submenu is shown
    again rather than jumping to the main menu.  Only ``b`` or ``q`` exits.
    """
    while True:
        print_header()
        print(f"  {_C.BOLD}Search in filenames.json{_C.RESET}\n")
        print(f"  {_C.BOLD}1{_C.RESET}  Search for Filenames   (substring match on filename)")
        print(f"  {_C.BOLD}2{_C.RESET}  Search for Filelist    (compare .txt file list against filenames.json)")
        print(f"  {_C.BOLD}3{_C.RESET}  Search for Keywords    (match in filename or path)")
        print(f"  {_C.BOLD}4{_C.RESET}  Search for Hashes      (MD5 / SHA256 lookup + report)")
        print(f"  {_C.BOLD}5{_C.RESET}  Match Hash List        (compare {HASHLISTS_DIR}/ against filenames.json)")
        print(f"  {_C.BOLD}6{_C.RESET}  Export Matched Files   (copy files from a previous hash match CSV)")
        print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

        choice = input("  Select search type [1/2/3/4/5/6/b/q]: ").strip().lower()
        if choice == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        if choice == "b" or not choice:
            return

        if choice == "1":
            handle_search_filenames_in_json()
        elif choice == "2":
            handle_search_filelist_in_json()
        elif choice == "3":
            handle_search_keywords_in_json()
        elif choice == "4":
            handle_search_hash()
        elif choice == "5":
            handle_match_hash_list()
        elif choice == "6":
            handle_export_hash_matches()
        else:
            print(f"\n  {_C.RED}Invalid choice.{_C.RESET}")
            wait_for_enter()


def display_menu() -> None:
    """Display the main menu."""
    print_header()
    items = [
        ("1", "Generate filenames list"),
        ("2", "Generate hashlist file"),
        ("3", "Search in filenames.json"),
        ("4", "Search in Files"),
        ("5", "Show search history"),
        ("6", "Settings"),
    ]
    for num, label in items:
        print(f"  {_C.BOLD}{num}{_C.RESET}  {label}")
    print()
    print(f"  {_C.DIM}q  Quit{_C.RESET}")
    print()


def handle_generate_list() -> None:
    """Handle menu option 1: Generate filenames list with source selection."""
    print_header()

    input_dir = Path(INPUT_DIR)
    input_exists = input_dir.is_dir()

    # Sub-menu
    print(f"  {_C.BOLD}Generate filenames list{_C.RESET}\n")
    opt1_note = "" if input_exists else f"  {_C.DIM}(not found){_C.RESET}"
    print(f"  {_C.BOLD}1{_C.RESET}  From \"{INPUT_DIR}\" directory (./{INPUT_DIR}/){opt1_note}")
    print(f"  {_C.BOLD}2{_C.RESET}  From custom path")
    print(f"  {_C.BOLD}3{_C.RESET}  From current directory (.)")
    print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

    choice = input("  Select source [1/2/3/b/q]: ").strip().lower()
    if choice == "q":
        clear_screen()
        print("Thank you for using txtFinder!")
        raise SystemExit(0)
    if choice == "b":
        return

    scan_dir: Path | None = None
    exclude_dirs = False

    if choice == "1":
        if not input_exists:
            print(f"\n  {_C.RED}Directory \"./{INPUT_DIR}/\" does not exist.{_C.RESET}")
            print(f"  Create it and place your files inside, then try again.")
            wait_for_enter()
            return
        ok, err = _validate_scan_directory(input_dir)
        if not ok:
            print(f"\n  {_C.RED}{err}{_C.RESET}")
            wait_for_enter()
            return
        scan_dir = input_dir

    elif choice == "2":
        raw_path = input("\n  Enter directory path: ").strip()
        if not raw_path:
            print(f"\n  {_C.RED}No path entered.{_C.RESET}")
            wait_for_enter()
            return
        scan_dir = Path(raw_path).resolve()
        ok, err = _validate_scan_directory(scan_dir)
        if not ok:
            print(f"\n  {_C.RED}{err}{_C.RESET}")
            wait_for_enter()
            return

    elif choice == "3":
        scan_dir = Path(".")
        exclude_dirs = True

    else:
        print(f"\n  {_C.RED}Invalid choice.{_C.RESET}")
        wait_for_enter()
        return

    # Ask file type filter
    print_header()
    print(f"  {_C.BOLD}File type filter{_C.RESET}")
    print(f"  {_C.BOLD}1{_C.RESET}  All files              {_C.DIM}(no filter){_C.RESET}")
    print(f"  {_C.BOLD}2{_C.RESET}  Media files only       {_C.DIM}(video + image — configurable in Settings){_C.RESET}")
    print(f"  {_C.BOLD}3{_C.RESET}  Document files only    {_C.DIM}(PDF, Office, text — configurable in Settings){_C.RESET}")
    print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}")
    filter_choice = input("  Filter [1/2/3/b/q]: ").strip().lower()
    if filter_choice == "q":
        clear_screen()
        print("Thank you for using txtFinder!")
        raise SystemExit(0)
    if filter_choice == "b":
        return
    if filter_choice == "2":
        ext_filter: frozenset | None = MEDIA_EXTENSIONS
        filter_label = f"Media only ({len(MEDIA_EXTENSIONS)} extensions)"
    elif filter_choice == "3":
        ext_filter = DOCUMENT_EXTENSIONS
        filter_label = f"Documents only ({len(DOCUMENT_EXTENSIONS)} extensions)"
    else:
        ext_filter = None
        filter_label = "All files"

    # Preview: quick file count (respects active filter)
    print(f"\n  Counting files in {scan_dir} ...")
    total_files = _count_files_in_dir(scan_dir, extensions=ext_filter)
    if total_files == 0:
        print(f"  {_C.RED}No matching files found.{_C.RESET}")
        wait_for_enter()
        return

    print(f"  Found {_C.BOLD}{total_files:,}{_C.RESET} files ({filter_label}).")

    # Ask list mode
    print_header()
    print(f"  {_C.BOLD}List mode{_C.RESET}")
    print(f"  {_C.BOLD}1{_C.RESET}  Full filename + SHA256 + MD5  {_C.DIM}(both hashes){_C.RESET}")
    print(f"  {_C.BOLD}2{_C.RESET}  Full filename + SHA256 only   {_C.DIM}(faster than both){_C.RESET}")
    print(f"  {_C.BOLD}3{_C.RESET}  Full filename + MD5 only      {_C.DIM}(faster than both){_C.RESET}")
    print(f"  {_C.BOLD}4{_C.RESET}  Full filename without hashes  {_C.DIM}(no hashing, fast){_C.RESET}")
    print(f"  {_C.BOLD}5{_C.RESET}  Filename stem only            {_C.DIM}(no hashing, fastest){_C.RESET}")
    print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}")
    list_mode = input("  Mode [1/2/3/4/5/b/q]: ").strip().lower()
    if list_mode == "q":
        clear_screen()
        print("Thank you for using txtFinder!")
        raise SystemExit(0)
    if list_mode == "b":
        return
    stems_only = list_mode == "5"
    no_hashes  = list_mode == "4"
    hash_mode  = "sha256" if list_mode == "2" else "md5" if list_mode == "3" else "both"

    # Ask deduplication option
    print_header()
    print(f"  {_C.BOLD}Duplicates{_C.RESET}")
    print(f"  {_C.BOLD}1{_C.RESET}  Keep all  {_C.DIM}(include duplicates){_C.RESET}")
    print(f"  {_C.BOLD}2{_C.RESET}  Remove duplicates  {_C.DIM}(each filename only once, case-insensitive){_C.RESET}")
    print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}")
    dedup_choice = input("  Duplicates [1/2/b/q]: ").strip().lower()
    if dedup_choice == "q":
        clear_screen()
        print("Thank you for using txtFinder!")
        raise SystemExit(0)
    if dedup_choice == "b":
        return
    deduplicate = dedup_choice == "2"

    print_header()
    confirm = input("  Continue? [Y/n]: ").strip().lower()
    if confirm == "n":
        return

    # Run scan with spinner + live counter
    _no_hash_mode = stems_only or no_hashes
    if _no_hash_mode:
        _spinner_msg = "Scanning files..."
        _verb = "Scanning"
    elif hash_mode == "sha256":
        _spinner_msg = "Hashing files (SHA256)..."
        _verb = "Hashing (SHA256)"
    elif hash_mode == "md5":
        _spinner_msg = "Hashing files (MD5)..."
        _verb = "Hashing (MD5)"
    else:
        _spinner_msg = "Hashing files (SHA256+MD5)..."
        _verb = "Hashing"
    try:
        with Spinner(_spinner_msg) as sp:
            def _progress(count: int) -> None:
                """Update spinner with the current file-scan count.

                Args:
                    count: Number of files processed so far.
                """
                sp.update_message(f"{_verb} file {count:,} / ~{total_files:,}")

            entries_written, duplicates_removed = generate_filenames_list(
                scan_dir=scan_dir,
                output_dir=Path("."),
                progress_callback=_progress,
                exclude_dirs=exclude_dirs,
                stems_only=stems_only,
                no_hashes=no_hashes,
                deduplicate=deduplicate,
                extensions=ext_filter,
                hash_mode=hash_mode,
            )
        if duplicates_removed:
            print(f"  {FILENAMES_JSON} created with {entries_written:,} entries "
                  f"({_C.YELLOW}{duplicates_removed:,} duplicate(s) removed{_C.RESET}).")
        else:
            print(f"  {FILENAMES_JSON} created with {entries_written:,} files.")
    except (OSError, RuntimeError, ValueError) as e:
        print(f"  {_C.RED}Error: Failed to generate filenames list. Check txtfinder.log for details.{_C.RESET}")
        logger.error(f"Error in handle_generate_list: {e}")

    wait_for_enter()


def handle_generate_hashlist() -> None:
    """Handle menu option 2: Generate hash list file (md5 or sha256 submenu)."""
    if not Path(FILENAMES_JSON).exists():
        print_header()
        print(f"  {_C.YELLOW}{FILENAMES_JSON} not found.{_C.RESET}")
        answer = input("  Generate it now from the current directory? [y/n]: ").strip().lower()
        if answer != "y":
            wait_for_enter()
            return
        try:
            with Spinner("Generating filenames list"):
                entries, _ = generate_filenames_list(scan_dir=Path("."))
            if not entries:
                print(f"  {_C.YELLOW}No files found — {FILENAMES_JSON} not created.{_C.RESET}")
                wait_for_enter()
                return
            print(f"  {_C.GREEN}{FILENAMES_JSON} generated ({len(entries):,} entries).{_C.RESET}")
        except (OSError, RuntimeError, ValueError) as e:
            print(f"  {_C.RED}Generation failed: {e}{_C.RESET}")
            wait_for_enter()
            return

    while True:
        print_header()
        print(f"  {_C.BOLD}Generate hashlist file{_C.RESET}\n")
        print(f"  {_C.BOLD}1{_C.RESET}  Generate hashlist md5")
        print(f"  {_C.BOLD}2{_C.RESET}  Generate hashlist sha256")
        print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

        choice = input("  Select [1/2/b/q]: ").strip().lower()
        if choice == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        if choice == "b":
            return
        if choice in ("1", "2"):
            break
        print(f"  {_C.RED}Invalid input.{_C.RESET}")
        wait_for_enter()

    hash_type = "md5" if choice == "1" else "sha256"
    try:
        with Spinner(f"Generating {hash_type.upper()} hash list"):
            count = generate_hashlist(hash_type)
        print(f"  filenames_{hash_type}.txt created with {count} hashes.")
    except (OSError, RuntimeError, ValueError) as e:
        print(f"  {_C.RED}Error: Failed to generate hash list. Check txtfinder.log for details.{_C.RESET}")
        logger.error(f"Error in handle_generate_hashlist({hash_type}): {e}")
    wait_for_enter()


def _search_single_pdf(selected_pdf: Path, filenames: list[str], fuzzy: bool = False,
                       extra_words: list[str] | None = None,
                       json_data: list[FileEntry] | None = None,
                       stems_mode: bool = False,
                       stems_min_len: int = 0,
                       search_mode: str = "exact") -> None:
    """Search a single PDF interactively and display detailed results.

    Args:
        selected_pdf: Path to the PDF file to search.
        filenames: List of filenames loaded from filenames.json.
        fuzzy: When ``True`` apply fuzzy matching instead of exact search.
        extra_words: Optional additional search terms to include in the report.
        json_data: Pre-loaded FileEntry list; loaded on demand if ``None``.
        stems_mode: When ``True`` use stem-only pattern set for the search.
        stems_min_len: Minimum stem length used when ``stems_mode`` is active.
        search_mode: Search mode label recorded in the report context paragraph.
    """
    print_header()
    mode_label = f" {_C.YELLOW}[FUZZY]{_C.RESET}" if fuzzy else ""
    print(f"Loaded {len(filenames)} filenames from JSON.{mode_label}\n")

    try:
        spinner_msg = f"Searching in {selected_pdf.name}"
        if fuzzy:
            spinner_msg += " (fuzzy)"
        with Spinner(spinner_msg) as sp:
            def _page_progress(page: int, total: int) -> None:
                """Update spinner with the current page number during PDF search.

                Args:
                    page: Current page number being processed.
                    total: Total number of pages in the document.
                """
                sp.update_message(f"Searching in {selected_pdf.name} (page {page}/{total})")

            # In stems mode use the pre-built stem pattern set (no variants)
            pattern_set = build_stem_pattern_set(filenames, min_length=stems_min_len) if stems_mode else None
            hit_count, output_path, has_native_text, ocr_used = search_and_highlight_pdf(
                selected_pdf, filenames, fuzzy=fuzzy,
                progress_callback=_page_progress,
                pattern_set=pattern_set)
            # Check for accumulated MuPDF warnings (suppressed from stdout)
            mupdf_warns = fitz.TOOLS.mupdf_warnings()
            if mupdf_warns:
                warn_lines = [l for l in mupdf_warns.strip().splitlines() if l.strip()]
                sp.warn(f"{len(warn_lines)} MuPDF warning(s)")
                logger.warning(f"MuPDF warnings for {selected_pdf.name}: {mupdf_warns}")

        w = 50
        print(f"\n{_C.CYAN}{'=' * w}{_C.RESET}")
        print(f"  {_C.BOLD}Results for: {selected_pdf.name}{_C.RESET}")
        print(f"{_C.CYAN}{'=' * w}{_C.RESET}")
        print(f"  Native Text: {'Yes' if has_native_text else 'No'}")

        if not has_native_text and not TESSERACT_AVAILABLE:
            print(f"  OCR: Not available (Tesseract not installed)")
        elif ocr_used:
            print(f"  OCR Used: Yes")
        else:
            print(f"  OCR Used: No")

        if fuzzy:
            print(f"  Fuzzy Matching: ON (threshold {FUZZY_THRESHOLD:.0%})")

        print(f"  Total Hits: {hit_count}")
        if output_path:
            print(f"  Output: {output_path.name}")
        else:
            print("  No matches - no new PDF created")
        print(f"{_C.CYAN}{'=' * w}{_C.RESET}")

        # Generate report using cached text (avoids double OCR extraction)
        if json_data is None:
            json_data = load_json_data()
        cached_text = _pop_cached_text(selected_pdf)
        with Spinner("Generating report"):
            report = _generate_report_for_file(selected_pdf, json_data, cached_text,
                                               extra_words=extra_words, stems_mode=stems_mode,
                                               search_mode=search_mode,
                                               stems_min_length=stems_min_len,
                                               ocr_used=ocr_used)
        if report:
            print(f"  Report: {report.name}")
            logger.info(f"Report: {report.name}")
        print()

    except IOError as e:
        print(f"\nError: PDF file cannot be read - {e}")
        print("Suggestions:")
        print("  1. Verify the file is a valid PDF (try opening in a PDF viewer)")
        print("  2. Check file permissions")
        print("  3. Re-download the file if it's damaged")
        logger.error(f"IOError in handle_search_pdf: {e}")
    except ValueError as e:
        print(f"\nError: PDF validation failed - {e}")
        print("Suggestion: The PDF may be empty or corrupt. Try another file.")
        logger.error(f"ValueError in handle_search_pdf: {e}")
    except Exception as e:
        print(f"\nError: Failed to process PDF. Check txtfinder.log for details.")
        logger.error(f"Error in handle_search_pdf: {e}")


def _search_all_pdfs(pdf_files: list[Path], filenames: list[str], fuzzy: bool = False,
                     extra_words: list[str] | None = None,
                     json_data: list[FileEntry] | None = None,
                     mode: str = "exact",
                     stems_min_len: int = 0) -> None:
    """Batch-search all PDFs using parallel processing.

    Args:
        pdf_files: List of PDF paths to process.
        filenames: List of filenames loaded from filenames.json.
        fuzzy: When ``True`` apply fuzzy matching instead of exact search.
        extra_words: Optional additional search terms forwarded to report generation.
        json_data: Pre-loaded FileEntry list; loaded on demand if ``None``.
        mode: Search mode string (``"exact"``, ``"fuzzy"``, or ``"stems"``).
        stems_min_len: Minimum stem length used when ``mode`` is ``"stems"``.
    """
    print_header()
    if json_data is None:
        json_data = load_json_data()

    search_fn = partial(_search_pdf_file, fuzzy=fuzzy) if fuzzy else _search_pdf_file
    label = "PDF files"
    if fuzzy:
        label += " (fuzzy)"

    if mode == "stems":
        ps = build_stem_pattern_set(filenames, min_length=stems_min_len)
    elif fuzzy:
        ps = None
    else:
        ps = build_pattern_set(filenames)
    with ProgressBar(len(pdf_files), label) as pb:
        results, interrupted = process_batch(pdf_files, filenames, search_fn, json_data,
                                             extra_words=extra_words,
                                             progress_callback=pb.advance,
                                             pattern_set=ps,
                                             stems_mode=(mode == "stems"),
                                             search_mode=mode,
                                             stems_min_length=stems_min_len)
    # Check for accumulated MuPDF warnings
    mupdf_warns = fitz.TOOLS.mupdf_warnings()
    if mupdf_warns:
        warn_lines = [l for l in mupdf_warns.strip().splitlines() if l.strip()]
        logger.warning(f"MuPDF warnings during batch PDF search ({len(warn_lines)}): {mupdf_warns}")

    _display_batch_results(results, interrupted)

    total = sum(r.matches for r in results)
    with_matches = sum(1 for r in results if r.matches > 0)
    print_batch_summary("PDF", len(results), with_matches, total, mode=mode)
    with Spinner("Generating CSV summary"):
        csv_path = generate_csv_summary(results)
    print(f"  CSV: {csv_path}")


def handle_search_pdf(
    search_dir: Path = Path("."),
    config: "_SearchConfig | None" = None,
) -> None:
    """Search PDF files for filenames from filenames.json.

    Args:
        search_dir: Directory to scan for PDF files.
        config: Pre-built search configuration; when provided the interactive
            mode/extra-words dialog is skipped and batch mode is forced
            (used by "All of the above").
    """
    print_header()

    pdf_files = get_pdf_files(search_dir)

    if not pdf_files:
        print(f"No PDF files found in {search_dir}.")
        wait_for_enter()
        return

    # When a shared config is provided (called from "All of the above") always
    # run in batch mode so the PDF-selection dialog can be skipped entirely.
    if config is not None:
        batch_mode = True
        selected_pdf = pdf_files[0]  # unused in batch mode; avoids unbound-variable
    else:
        # Display PDF selection menu
        while True:
            print_header()
            print(f"  Select a PDF to search:\n")
            print(f"  {_C.BOLD}a{_C.RESET}  Search all PDFs ({len(pdf_files)} files)")
            print()
            for idx, pdf_file in enumerate(pdf_files, 1):
                print(f"  {_C.BOLD}{idx}{_C.RESET}  {pdf_file.name}")
            print(f"\n  {_C.DIM}b  Back   q  Quit{_C.RESET}\n")

            choice = input("  Enter choice: ").strip().lower()

            if choice == "q":
                clear_screen()
                print("Thank you for using txtFinder!")
                raise SystemExit(0)
            if choice == "b":
                return

            if choice == "a":
                batch_mode = True
                break

            try:
                pdf_index = int(choice) - 1
                if 0 <= pdf_index < len(pdf_files):
                    selected_pdf = pdf_files[pdf_index]
                    batch_mode = False
                    break
                else:
                    print(f"  {_C.RED}Invalid selection. Try again.{_C.RESET}")
                    wait_for_enter()
            except ValueError:
                print(f"  {_C.RED}Invalid input. Enter a number, 'a' for all, 'b' to go back, or 'q' to quit.{_C.RESET}")
                wait_for_enter()

    # Check if filenames.json exists
    json_path = Path(FILENAMES_JSON)
    if not json_path.exists():
        print(f"\n  {FILENAMES_JSON} not found.")
        create_choice = input("  Create first via Menu 1? [y/n]: ").strip().lower()
        if create_choice == "y":
            handle_generate_list()
            if not json_path.exists():
                wait_for_enter()
                return
        else:
            wait_for_enter()
            return

    # Load JSON once — derive both json_data (for reports) and filenames from it
    try:
        json_data = load_json_data(json_path)
        filenames = [e.filename for e in json_data]
    except json.JSONDecodeError:
        print(f"\nError: {FILENAMES_JSON} is malformed (invalid JSON syntax).")
        print("Suggestion: Regenerate the file using Menu Option 1.")
        wait_for_enter()
        return
    except ValueError as e:
        print(f"\nError: {FILENAMES_JSON} has invalid structure: {e}")
        print("Suggestion: Regenerate the file using Menu Option 1.")
        wait_for_enter()
        return
    except Exception as e:
        print(f"\nError: Failed to load {FILENAMES_JSON}. Check txtfinder.log for details.")
        wait_for_enter()
        return

    if not filenames:
        print("No filenames found in JSON.")
        wait_for_enter()
        return

    # Search mode + extra words (with optional profile)
    if config is not None:
        mode = config.mode
        regex_patterns = config.regex_patterns
        extra_words = config.extra_words
        stems_min_len = config.stems_min_len
    else:
        mode = _ask_search_mode()
        if mode == "back":
            return
        if mode == "stems":
            regex_patterns, extra_words = [], []
            stems_min_len = _ask_stems_min_length()
        elif mode == "profile":
            loaded = _load_profile_interactive()
            if loaded:
                mode, regex_patterns, extra_words = loaded
            else:
                mode, regex_patterns, extra_words = "exact", [], []
            stems_min_len = 0
        else:
            regex_patterns = _ask_regex_patterns() if mode == "regex" else []
            extra_words = _ask_extra_search_words()
            stems_min_len = 0

        # Offer to save profile (skipped for stems mode and when using shared config)
        if mode != "stems":
            save_name = input("  Save as profile? (name or Enter to skip): ").strip()
            if save_name:
                save_search_profile(save_name, mode, extra_words,
                                    [rx.pattern for rx in regex_patterns])
                print(f"  {_C.GREEN}Profile '{save_name}' saved{_C.RESET}")

    fuzzy = mode == "fuzzy"
    if extra_words:
        filenames = filenames + extra_words

    if regex_patterns:
        # Regex mode: use generic regex search for PDFs
        search_fn = partial(_search_file_with_regex, regex_patterns=regex_patterns)
        target_files = pdf_files if batch_mode else [selected_pdf]
        with ProgressBar(len(target_files), "PDF files (regex)") as pb:
            results, interrupted = process_batch(target_files, filenames, search_fn, json_data,
                                                 extra_words=extra_words or None,
                                                 progress_callback=pb.advance,
                                                 search_mode="regex")
        _display_batch_results(results, interrupted)
        total = sum(r.matches for r in results)
        with_matches = sum(1 for r in results if r.matches > 0)
        print_batch_summary("PDF", len(results), with_matches, total, mode=mode)
        with Spinner("Generating CSV summary"):
            csv_path = generate_csv_summary(results)
        print(f"  CSV: {csv_path}")
    elif batch_mode:
        _search_all_pdfs(pdf_files, filenames, fuzzy=fuzzy,
                         extra_words=extra_words or None, json_data=json_data,
                         mode=mode, stems_min_len=stems_min_len)
    else:
        _search_single_pdf(selected_pdf, filenames, fuzzy=fuzzy,
                           extra_words=extra_words or None, json_data=json_data,
                           stems_mode=(mode == "stems"), stems_min_len=stems_min_len,
                           search_mode=mode)

    wait_for_enter()


def interactive_menu() -> None:
    """Run the interactive menu loop."""
    while True:
        display_menu()
        choice = input("  Enter choice: ").strip().lower()

        if choice == "1":
            handle_generate_list()
        elif choice == "2":
            handle_generate_hashlist()
        elif choice == "3":
            handle_search_in_json()
        elif choice == "4":
            handle_search_files()
        elif choice == "5":
            handle_search_history()
        elif choice == "6":
            handle_settings()
        elif choice == "q":
            clear_screen()
            print("Thank you for using txtFinder!")
            raise SystemExit(0)
        else:
            print("Invalid choice. Please enter 1-6 or q.")
            wait_for_enter()


def _load_config_on_startup() -> None:
    """Load txtfinder_config.json at startup if it exists, silently ignoring errors."""
    cfg_path = Path("txtfinder_config.json")
    if not cfg_path.exists():
        return
    _VALID_KEYS = {
        "FUZZY_THRESHOLD", "OCR_DPI", "OCR_PSM", "OCR_MIN_CONFIDENCE", "OCR_LANGUAGES",
        "TXT_MARKER", "TXT_MARKER_POSITION", "MAX_SEARCH_PATTERNS", "OUTPUT_DIR",
        "STEMS_MIN_LENGTH", "OUTPUT_SUFFIX", "REPORT_SUFFIX", "HASHES_SUFFIX",
        "HASHLIST_REPORT_SUFFIX", "HASHLISTS_DIR", "EXPORT_DIR",
        "HIGHLIGHT_COLOR", "MAX_WORKERS", "MEDIA_EXTENSIONS", "DOCUMENT_EXTENSIONS",
        "REPORT_ON_MATCH_ONLY", "FILELIST_REPORT_SUFFIX", "HITLIST_SUFFIX",
    }
    _TYPE_MAP = {
        "FUZZY_THRESHOLD": float, "OCR_DPI": int, "OCR_PSM": int, "OCR_MIN_CONFIDENCE": int,
        "MAX_SEARCH_PATTERNS": int, "STEMS_MIN_LENGTH": int, "MAX_WORKERS": int,
        "OCR_LANGUAGES": str, "TXT_MARKER": str, "TXT_MARKER_POSITION": str, "OUTPUT_DIR": str,
        "OUTPUT_SUFFIX": str, "REPORT_SUFFIX": str, "HASHES_SUFFIX": str,
        "HASHLIST_REPORT_SUFFIX": str, "FILELIST_REPORT_SUFFIX": str, "HITLIST_SUFFIX": str,
        "HASHLISTS_DIR": str, "EXPORT_DIR": str,
        "HIGHLIGHT_COLOR": _color_tuple,
        "MEDIA_EXTENSIONS": _parse_extensions, "DOCUMENT_EXTENSIONS": _parse_extensions,
        "REPORT_ON_MATCH_ONLY": _parse_bool,
    }
    try:
        raw = json.loads(cfg_path.read_text(encoding="utf-8"))
        for key, val in raw.items():
            if key in _VALID_KEYS:
                _apply_setting(key, _TYPE_MAP[key](val))
        logger.debug(f"Loaded config from {cfg_path}")
    except (OSError, json.JSONDecodeError, ValueError, TypeError) as e:
        logger.warning(f"Could not load {cfg_path}: {e}")


@app.command()
def main() -> None:
    """
    txtFinder - Interactive CLI tool for file management and PDF searching.

    Launch the interactive menu to generate file lists and search PDFs.
    """
    _load_config_on_startup()
    _cleanup_disk_cache()
    clear_screen()
    try:
        interactive_menu()
    except KeyboardInterrupt:
        clear_screen()
        print("\nInterrupted by user. Exiting...")
        sys.exit(0)
    except Exception as e:
        logger.error(f"Unexpected error in main: {e}")
        print(f"An unexpected error occurred. Check txtfinder.log for details.")
        sys.exit(1)


if __name__ == "__main__":
    app()
